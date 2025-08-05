import os
import json
import base64
import hashlib
import secrets
import asyncio
import traceback
import threading
import time
import re
from datetime import datetime, timedelta, timezone
from functools import wraps, lru_cache
from io import BytesIO
from typing import Dict, List, Optional, Any, Tuple, AsyncGenerator, Union
from dataclasses import dataclass, asdict, field
from enum import Enum
import mimetypes
import random
import string
from collections import deque, defaultdict
from concurrent.futures import ThreadPoolExecutor
import uuid
import queue
import aiofiles
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    print("Warning: python-magic not available. File type detection will use mimetypes fallback.")
import chardet

from flask import Flask, render_template_string, request, jsonify, Response, stream_with_context, session, make_response, send_file
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from flask_caching import Cache
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
from markupsafe import escape

from openai import OpenAI
import httpx
import redis
import jwt

from PIL import Image
import PyPDF2
import docx
import openpyxl
import pandas as pd
import numpy as np
from bs4 import BeautifulSoup
import markdown
import pygments
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import HtmlFormatter
import pytesseract
import cv2
import librosa
import soundfile as sf
import speech_recognition as sr
import textstat
import nltk
from transformers import pipeline
import torch
import pdfplumber
import camelot
import tabula
import ebooklib
from ebooklib import epub
import zipfile
import rarfile
import py7zr
import tarfile
import pptx
import xlrd
from icalendar import Calendar
import vobject
import eml_parser
import extract_msg
import pyzbar.pyzbar as pyzbar
import qrcode
import barcode
from barcode.writer import ImageWriter
import svgwrite
import fitz
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.graph_objects as go
import networkx as nx
from wordcloud import WordCloud
import folium
import geopandas as gpd
import rasterio
import xarray as xr
import h5py
import netCDF4
import astropy.io.fits as fits
import Bio.SeqIO
import rdkit
from rdkit import Chem
import pydicom
import nibabel as nib
import SimpleITK as sitk
import trimesh
import open3d as o3d
import MDAnalysis
import pytraj
import nglview

app = Flask(__name__)
app.secret_key = secrets.token_hex(64)
app.config.update(
    MAX_CONTENT_LENGTH=2 * 1024 * 1024 * 1024,
    SESSION_COOKIE_SECURE=True,
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE='Lax',
    PERMANENT_SESSION_LIFETIME=timedelta(days=90),
    JSON_SORT_KEYS=False,
    SEND_FILE_MAX_AGE_DEFAULT=31536000,
    UPLOAD_FOLDER='uploads',
    DOWNLOAD_FOLDER='downloads',
    TEMP_FOLDER='temp',
    CACHE_FOLDER='cache',
    SESSION_TYPE='filesystem',
    SESSION_FILE_DIR='sessions',
    SESSION_FILE_THRESHOLD=500,
    SESSION_FILE_MODE=384,
    PROPAGATE_EXCEPTIONS=True,
    TRAP_HTTP_EXCEPTIONS=True,
    TRAP_BAD_REQUEST_ERRORS=True,
    PREFERRED_URL_SCHEME='https',
    APPLICATION_ROOT='/',
    SESSION_COOKIE_NAME='ai_ultra_session',
    SESSION_COOKIE_PATH='/',
    SESSION_COOKIE_DOMAIN=None,
    SESSION_REFRESH_EACH_REQUEST=True,
    MAX_COOKIE_SIZE=4093,
    BABEL_DEFAULT_LOCALE='en',
    BABEL_DEFAULT_TIMEZONE='UTC',
    JSONIFY_PRETTYPRINT_REGULAR=True,
    JSONIFY_MIMETYPE='application/json',
    TEMPLATES_AUTO_RELOAD=True,
    EXPLAIN_TEMPLATE_LOADING=False,
    MAX_CONTENT_PATH=None,
)

CORS(app, resources={
    r"/api/*": {
        "origins": "*",
        "methods": ["GET", "POST", "PUT", "DELETE", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization", "X-Requested-With", "X-Session-ID"],
        "expose_headers": ["Content-Range", "X-Content-Type-Options", "X-Frame-Options"],
        "supports_credentials": True,
        "max_age": 3600
    }
})

limiter = Limiter(
    app=app,
    key_func=get_remote_address,
    default_limits=["10000 per hour", "1000 per minute"],
    storage_uri="memory://",
    strategy="fixed-window",
    headers_enabled=True,
    swallow_errors=True,
    in_memory_fallback_enabled=True,
    retry_after="http-date"
)

cache = Cache(app, config={
    'CACHE_TYPE': 'filesystem',
    'CACHE_DIR': 'cache',
    'CACHE_DEFAULT_TIMEOUT': 3600,
    'CACHE_THRESHOLD': 10000,
    'CACHE_KEY_PREFIX': 'ai_ultra_',
    'CACHE_OPTIONS': {'mode': 0o666}
})

executor = ThreadPoolExecutor(
    max_workers=100,
    thread_name_prefix='ai_ultra_worker'
)

for folder in ['uploads', 'downloads', 'temp', 'cache', 'sessions']:
    os.makedirs(folder, exist_ok=True)

class MessageRole(Enum):
    USER = "user"
    ASSISTANT = "assistant"
    SYSTEM = "system"
    FUNCTION = "function"
    TOOL = "tool"
    ERROR = "error"
    INFO = "info"
    WARNING = "warning"
    DEBUG = "debug"

class ChatMode(Enum):
    NORMAL = "normal"
    CREATIVE = "creative"
    PRECISE = "precise"
    BALANCED = "balanced"
    CODE = "code"
    RESEARCH = "research"
    EXPERT = "expert"
    TUTOR = "tutor"
    ANALYST = "analyst"
    WRITER = "writer"
    DESIGNER = "designer"
    STRATEGIST = "strategist"
    COUNSELOR = "counselor"
    TRANSLATOR = "translator"
    DEBUGGER = "debugger"
    OPTIMIZER = "optimizer"
    ARCHITECT = "architect"
    SCIENTIST = "scientist"
    PHILOSOPHER = "philosopher"
    STORYTELLER = "storyteller"

class FileType(Enum):
    IMAGE = "image"
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    CODE = "code"
    DATA = "data"
    DATABASE = "database"
    ARCHIVE = "archive"
    AUDIO = "audio"
    VIDEO = "video"
    MODEL = "model"
    DIAGRAM = "diagram"
    MARKUP = "markup"
    CONFIG = "config"
    LOG = "log"
    BINARY = "binary"
    SCIENTIFIC = "scientific"
    MEDICAL = "medical"
    GEOSPATIAL = "geospatial"
    MOLECULAR = "molecular"

class ProcessingStatus(Enum):
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"
    CANCELLED = "cancelled"
    PAUSED = "paused"
    QUEUED = "queued"

class AIModel(Enum):
    GEMINI_2_FLASH_EXP = "gemini-2.0-flash-exp"
    GEMINI_2_FLASH_THINKING = "gemini-2.0-flash-thinking-exp"
    GEMINI_15_PRO = "gemini-1.5-pro"
    GEMINI_15_FLASH = "gemini-1.5-flash"

@dataclass
class StreamToken:
    content: str
    timestamp: float
    token_count: int = 1
    metadata: Dict[str, Any] = field(default_factory=dict)
    chunk_id: str = field(default_factory=lambda: str(uuid.uuid4()))
    sequence: int = 0
    is_final: bool = False

@dataclass
class Attachment:
    id: str
    filename: str
    size: int
    mime_type: str
    file_type: FileType
    path: str
    processed: bool = False
    metadata: Dict[str, Any] = field(default_factory=dict)
    extracted_text: Optional[str] = None
    thumbnail: Optional[str] = None
    preview: Optional[str] = None
    analysis: Optional[Dict[str, Any]] = None
    error: Optional[str] = None
    processing_time: Optional[float] = None
    checksum: Optional[str] = None

@dataclass
class Message:
    id: str
    role: MessageRole
    content: str
    timestamp: datetime
    tokens: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)
    attachments: List[Attachment] = field(default_factory=list)
    stream_tokens: List[StreamToken] = field(default_factory=list)
    parent_id: Optional[str] = None
    children_ids: List[str] = field(default_factory=list)
    edited: bool = False
    edited_at: Optional[datetime] = None
    reactions: Dict[str, int] = field(default_factory=dict)
    citations: List[Dict[str, Any]] = field(default_factory=list)
    model_used: Optional[str] = None
    processing_time: Optional[float] = None
    cost: Optional[float] = None
    quality_score: Optional[float] = None
    
    def to_dict(self):
        return {
            "id": self.id,
            "role": self.role.value,
            "content": self.content,
            "timestamp": self.timestamp.isoformat(),
            "tokens": self.tokens,
            "metadata": self.metadata,
            "attachments": [asdict(att) for att in self.attachments],
            "stream_tokens": [asdict(st) for st in self.stream_tokens],
            "parent_id": self.parent_id,
            "children_ids": self.children_ids,
            "edited": self.edited,
            "edited_at": self.edited_at.isoformat() if self.edited_at else None,
            "reactions": self.reactions,
            "citations": self.citations,
            "model_used": self.model_used,
            "processing_time": self.processing_time,
            "cost": self.cost,
            "quality_score": self.quality_score
        }

@dataclass
class ChatSession:
    id: str
    user_id: Optional[str]
    title: Optional[str]
    messages: List[Message]
    total_tokens: int
    created_at: datetime
    updated_at: datetime
    mode: ChatMode
    context_window: int
    metadata: Dict[str, Any]
    active_stream: Optional[queue.Queue] = None
    stream_controller: Optional[Any] = None
    tags: List[str] = field(default_factory=list)
    shared: bool = False
    share_id: Optional[str] = None
    collaborators: List[str] = field(default_factory=list)
    version: int = 1
    branches: List[Dict[str, Any]] = field(default_factory=list)
    settings: Dict[str, Any] = field(default_factory=dict)
    analytics: Dict[str, Any] = field(default_factory=dict)
    
    def add_message(self, message: Message):
        self.messages.append(message)
        self.total_tokens += message.tokens
        self.updated_at = datetime.now(timezone.utc)
        self.update_analytics(message)
    
    def update_analytics(self, message: Message):
        if "message_count" not in self.analytics:
            self.analytics = {
                "message_count": 0,
                "user_messages": 0,
                "assistant_messages": 0,
                "total_processing_time": 0,
                "average_response_time": 0,
                "total_cost": 0,
                "quality_scores": [],
                "models_used": defaultdict(int),
                "modes_used": defaultdict(int),
                "attachments_processed": 0,
                "errors": 0
            }
        
        self.analytics["message_count"] += 1
        if message.role == MessageRole.USER:
            self.analytics["user_messages"] += 1
        elif message.role == MessageRole.ASSISTANT:
            self.analytics["assistant_messages"] += 1
        
        if message.processing_time:
            self.analytics["total_processing_time"] += message.processing_time
            self.analytics["average_response_time"] = (
                self.analytics["total_processing_time"] / 
                self.analytics["assistant_messages"]
            ) if self.analytics["assistant_messages"] > 0 else 0
        
        if message.cost:
            self.analytics["total_cost"] += message.cost
        
        if message.quality_score:
            self.analytics["quality_scores"].append(message.quality_score)
        
        if message.model_used:
            self.analytics["models_used"][message.model_used] += 1
        
        self.analytics["modes_used"][self.mode.value] += 1
        
        if message.attachments:
            self.analytics["attachments_processed"] += len(message.attachments)
        
        if message.role == MessageRole.ERROR:
            self.analytics["errors"] += 1
    
    def start_stream(self) -> queue.Queue:
        self.active_stream = queue.Queue(maxsize=1000)
        return self.active_stream
    
    def end_stream(self):
        if self.active_stream:
            self.active_stream.put(None)
            self.active_stream = None
    
    def to_dict(self):
        return {
            "id": self.id,
            "user_id": self.user_id,
            "title": self.title,
            "messages": [msg.to_dict() for msg in self.messages],
            "total_tokens": self.total_tokens,
            "created_at": self.created_at.isoformat(),
            "updated_at": self.updated_at.isoformat(),
            "mode": self.mode.value,
            "context_window": self.context_window,
            "metadata": self.metadata,
            "tags": self.tags,
            "shared": self.shared,
            "share_id": self.share_id,
            "collaborators": self.collaborators,
            "version": self.version,
            "branches": self.branches,
            "settings": self.settings,
            "analytics": self.analytics
        }

@dataclass
class GeminiKey:
    key: str
    rate_limit: int = 60
    daily_limit: int = 100000
    usage_today: int = 0
    usage_this_minute: int = 0
    failures: int = 0
    consecutive_failures: int = 0
    last_used: Optional[datetime] = None
    last_reset: datetime = field(default_factory=lambda: datetime.now(timezone.utc))
    last_minute_reset: datetime = field(default_factory=lambda: datetime.now(timezone.utc))
    is_active: bool = True
    performance_score: float = 1.0
    response_times: deque = field(default_factory=lambda: deque(maxlen=100))
    success_rate: float = 1.0
    total_requests: int = 0
    successful_requests: int = 0
    capabilities: List[str] = field(default_factory=list)
    restrictions: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)

class UltraSessionManager:
    def __init__(self):
        self.sessions: Dict[str, ChatSession] = {}
        self.user_sessions: Dict[str, List[str]] = defaultdict(list)
        self.session_lock = threading.RLock()
        self.max_sessions_per_user = 1000
        self.session_ttl = timedelta(days=365)
        self.active_streams: Dict[str, queue.Queue] = {}
        self.session_cache = LRUCache(maxsize=10000)
        self.persistence_queue = queue.Queue()
        self.start_persistence_worker()
    
    def start_persistence_worker(self):
        def worker():
            while True:
                try:
                    item = self.persistence_queue.get(timeout=1)
                    if item is None:
                        break
                    action, data = item
                    if action == "save":
                        self._persist_session(data)
                    elif action == "delete":
                        self._delete_persisted_session(data)
                except queue.Empty:
                    continue
                except Exception as e:
                    print(f"Persistence error: {e}")
        
        thread = threading.Thread(target=worker, daemon=True)
        thread.start()
    
    def _persist_session(self, session: ChatSession):
        try:
            session_path = os.path.join("sessions", f"{session.id}.json")
            with open(session_path, 'w') as f:
                json.dump(session.to_dict(), f, indent=2)
        except Exception as e:
            print(f"Failed to persist session {session.id}: {e}")
    
    def _delete_persisted_session(self, session_id: str):
        try:
            session_path = os.path.join("sessions", f"{session_id}.json")
            if os.path.exists(session_path):
                os.remove(session_path)
        except Exception as e:
            print(f"Failed to delete persisted session {session_id}: {e}")
    
    def create_session(self, user_id: Optional[str] = None, mode: ChatMode = ChatMode.BALANCED) -> ChatSession:
        session_id = str(uuid.uuid4())
        session = ChatSession(
            id=session_id,
            user_id=user_id,
            title=None,
            messages=[],
            total_tokens=0,
            created_at=datetime.now(timezone.utc),
            updated_at=datetime.now(timezone.utc),
            mode=mode,
            context_window=2000000,
            metadata={
                "version": "3.0",
                "features": ["streaming", "advanced_prompts", "multi_modal", "collaboration"],
                "capabilities": self._get_mode_capabilities(mode)
            }
        )
        
        with self.session_lock:
            self.sessions[session_id] = session
            if user_id:
                self.user_sessions[user_id].append(session_id)
                if len(self.user_sessions[user_id]) > self.max_sessions_per_user:
                    old_session = self.user_sessions[user_id].pop(0)
                    self.delete_session(old_session)
        
        self.persistence_queue.put(("save", session))
        return session
    
    def _get_mode_capabilities(self, mode: ChatMode) -> List[str]:
        capabilities_map = {
            ChatMode.NORMAL: ["general", "conversation", "basic_analysis"],
            ChatMode.CREATIVE: ["storytelling", "ideation", "artistic", "innovative"],
            ChatMode.PRECISE: ["accuracy", "detail", "technical", "factual"],
            ChatMode.BALANCED: ["versatile", "adaptive", "comprehensive"],
            ChatMode.CODE: ["programming", "debugging", "optimization", "architecture"],
            ChatMode.RESEARCH: ["analysis", "synthesis", "citation", "deep_dive"],
            ChatMode.EXPERT: ["domain_expertise", "professional", "advanced"],
            ChatMode.TUTOR: ["educational", "explanatory", "step_by_step", "adaptive_learning"],
            ChatMode.ANALYST: ["data_analysis", "visualization", "insights", "reporting"],
            ChatMode.WRITER: ["content_creation", "editing", "style", "formatting"],
            ChatMode.DESIGNER: ["ui_ux", "visual", "creative", "prototyping"],
            ChatMode.STRATEGIST: ["planning", "analysis", "decision_making", "forecasting"],
            ChatMode.COUNSELOR: ["empathy", "support", "guidance", "reflection"],
            ChatMode.TRANSLATOR: ["multilingual", "context", "cultural", "technical"],
            ChatMode.DEBUGGER: ["error_analysis", "troubleshooting", "optimization"],
            ChatMode.OPTIMIZER: ["performance", "efficiency", "refactoring", "scaling"],
            ChatMode.ARCHITECT: ["system_design", "patterns", "scalability", "best_practices"],
            ChatMode.SCIENTIST: ["research", "experimentation", "analysis", "hypothesis"],
            ChatMode.PHILOSOPHER: ["reasoning", "ethics", "logic", "contemplation"],
            ChatMode.STORYTELLER: ["narrative", "character", "plot", "engagement"]
        }
        return capabilities_map.get(mode, ["general"])
    
    def get_session(self, session_id: str) -> Optional[ChatSession]:
        if session_id in self.session_cache:
            return self.session_cache[session_id]
        
        with self.session_lock:
            session = self.sessions.get(session_id)
            if not session:
                session = self._load_session_from_disk(session_id)
            
            if session:
                if datetime.now(timezone.utc) - session.updated_at > self.session_ttl:
                    self.delete_session(session_id)
                    return None
                self.session_cache[session_id] = session
            
            return session
    
    def _load_session_from_disk(self, session_id: str) -> Optional[ChatSession]:
        try:
            session_path = os.path.join("sessions", f"{session_id}.json")
            if os.path.exists(session_path):
                with open(session_path, 'r') as f:
                    data = json.load(f)
                    messages = []
                    for msg_data in data.get("messages", []):
                        attachments = []
                        for att_data in msg_data.get("attachments", []):
                            attachments.append(Attachment(**att_data))
                        
                        stream_tokens = []
                        for st_data in msg_data.get("stream_tokens", []):
                            stream_tokens.append(StreamToken(**st_data))
                        
                        message = Message(
                            id=msg_data["id"],
                            role=MessageRole(msg_data["role"]),
                            content=msg_data["content"],
                            timestamp=datetime.fromisoformat(msg_data["timestamp"]),
                            tokens=msg_data.get("tokens", 0),
                            metadata=msg_data.get("metadata", {}),
                            attachments=attachments,
                            stream_tokens=stream_tokens,
                            parent_id=msg_data.get("parent_id"),
                            children_ids=msg_data.get("children_ids", []),
                            edited=msg_data.get("edited", False),
                            edited_at=datetime.fromisoformat(msg_data["edited_at"]) if msg_data.get("edited_at") else None,
                            reactions=msg_data.get("reactions", {}),
                            citations=msg_data.get("citations", []),
                            model_used=msg_data.get("model_used"),
                            processing_time=msg_data.get("processing_time"),
                            cost=msg_data.get("cost"),
                            quality_score=msg_data.get("quality_score")
                        )
                        messages.append(message)
                    
                    session = ChatSession(
                        id=data["id"],
                        user_id=data.get("user_id"),
                        title=data.get("title"),
                        messages=messages,
                        total_tokens=data.get("total_tokens", 0),
                        created_at=datetime.fromisoformat(data["created_at"]),
                        updated_at=datetime.fromisoformat(data["updated_at"]),
                        mode=ChatMode(data.get("mode", "balanced")),
                        context_window=data.get("context_window", 2000000),
                        metadata=data.get("metadata", {}),
                        tags=data.get("tags", []),
                        shared=data.get("shared", False),
                        share_id=data.get("share_id"),
                        collaborators=data.get("collaborators", []),
                        version=data.get("version", 1),
                        branches=data.get("branches", []),
                        settings=data.get("settings", {}),
                        analytics=data.get("analytics", {})
                    )
                    
                    self.sessions[session_id] = session
                    return session
        except Exception as e:
            print(f"Failed to load session from disk: {e}")
        
        return None
    
    def update_session(self, session: ChatSession):
        with self.session_lock:
            self.sessions[session.id] = session
            self.session_cache[session.id] = session
        
        self.persistence_queue.put(("save", session))
    
    def delete_session(self, session_id: str):
        with self.session_lock:
            if session_id in self.sessions:
                session = self.sessions[session_id]
                if session.active_stream:
                    session.end_stream()
                if session.user_id and session_id in self.user_sessions[session.user_id]:
                    self.user_sessions[session.user_id].remove(session_id)
                del self.sessions[session_id]
                if session_id in self.session_cache:
                    del self.session_cache[session_id]
        
        self.persistence_queue.put(("delete", session_id))
    
    def get_user_sessions(self, user_id: str, limit: int = 100, offset: int = 0) -> List[ChatSession]:
        with self.session_lock:
            session_ids = self.user_sessions.get(user_id, [])[offset:offset + limit]
            return [self.get_session(sid) for sid in session_ids if self.get_session(sid)]
    
    def search_sessions(self, query: str, user_id: Optional[str] = None, 
                       mode: Optional[ChatMode] = None, tags: Optional[List[str]] = None,
                       date_from: Optional[datetime] = None, date_to: Optional[datetime] = None,
                       limit: int = 50) -> List[ChatSession]:
        results = []
        
        with self.session_lock:
            sessions_to_search = self.sessions.values()
            if user_id:
                session_ids = self.user_sessions.get(user_id, [])
                sessions_to_search = [self.sessions[sid] for sid in session_ids if sid in self.sessions]
            
            for session in sessions_to_search:
                if mode and session.mode != mode:
                    continue
                
                if tags and not any(tag in session.tags for tag in tags):
                    continue
                
                if date_from and session.created_at < date_from:
                    continue
                
                if date_to and session.created_at > date_to:
                    continue
                
                if query:
                    query_lower = query.lower()
                    found = False
                    
                    if session.title and query_lower in session.title.lower():
                        found = True
                    
                    if not found:
                        for message in session.messages:
                            if query_lower in message.content.lower():
                                found = True
                                break
                    
                    if not found:
                        for tag in session.tags:
                            if query_lower in tag.lower():
                                found = True
                                break
                    
                    if not found:
                        continue
                
                results.append(session)
                
                if len(results) >= limit:
                    break
        
        return sorted(results, key=lambda s: s.updated_at, reverse=True)[:limit]

class LRUCache:
    def __init__(self, maxsize=1000):
        self.cache = {}
        self.maxsize = maxsize
        self.lock = threading.Lock()
    
    def __contains__(self, key):
        with self.lock:
            return key in self.cache
    
    def __getitem__(self, key):
        with self.lock:
            if key in self.cache:
                value = self.cache.pop(key)
                self.cache[key] = value
                return value
            raise KeyError(key)
    
    def __setitem__(self, key, value):
        with self.lock:
            if key in self.cache:
                self.cache.pop(key)
            elif len(self.cache) >= self.maxsize:
                self.cache.pop(next(iter(self.cache)))
            self.cache[key] = value
    
    def __delitem__(self, key):
        with self.lock:
            if key in self.cache:
                del self.cache[key]

class UltraGeminiKeyManager:
    def __init__(self):
        self.keys: Dict[str, GeminiKey] = {}
        self.key_lock = threading.RLock()
        self.load_keys()
        self.start_maintenance()
    
    def load_keys(self):
        keys_loaded = 0
        
        for i in range(1, 1000):
            key = os.environ.get(f'GEMINI_API_KEY_{i}') or os.environ.get(f'GEMINI_KEY_{i}')
            if key:
                self.add_key(key)
                keys_loaded += 1
        
        default_key = os.environ.get('GEMINI_API_KEY')
        if default_key:
            self.add_key(default_key, is_primary=True)
            keys_loaded += 1
        
        if keys_loaded == 0:
            print("WARNING: No Gemini API keys found!")
    
    def add_key(self, key: str, is_primary: bool = False):
        key_id = hashlib.sha256(key.encode()).hexdigest()[:16]
        gemini_key = GeminiKey(
            key=key,
            capabilities=self._detect_key_capabilities(key),
            metadata={"is_primary": is_primary}
        )
        self.keys[key_id] = gemini_key
    
    def _detect_key_capabilities(self, key: str) -> List[str]:
        capabilities = ["chat", "completion", "streaming"]
        
        if "pro" in key.lower():
            capabilities.extend(["extended_context", "advanced_reasoning", "multimodal"])
        
        return capabilities
    
    def get_best_key(self, required_capabilities: Optional[List[str]] = None) -> Optional[GeminiKey]:
        with self.key_lock:
            now = datetime.now(timezone.utc)
            available_keys = []
            
            for key_id, key in self.keys.items():
                if not key.is_active or key.consecutive_failures > 5:
                    continue
                
                if key.last_minute_reset and (now - key.last_minute_reset).total_seconds() >= 60:
                    key.usage_this_minute = 0
                    key.last_minute_reset = now
                
                if key.usage_this_minute >= key.rate_limit:
                    continue
                
                if key.usage_today >= key.daily_limit:
                    continue
                
                if required_capabilities:
                    if not all(cap in key.capabilities for cap in required_capabilities):
                        continue
                
                score = self._calculate_key_score(key)
                available_keys.append((key_id, key, score))
            
            if not available_keys:
                return None
            
            available_keys.sort(key=lambda x: x[2], reverse=True)
            
            selected_key = available_keys[0][1]
            selected_key.last_used = now
            selected_key.usage_today += 1
            selected_key.usage_this_minute += 1
            selected_key.total_requests += 1
            
            return selected_key
    
    def _calculate_key_score(self, key: GeminiKey) -> float:
        score = key.performance_score
        
        usage_ratio = key.usage_today / key.daily_limit if key.daily_limit > 0 else 0
        score *= (1 - usage_ratio)
        
        if key.response_times:
            avg_response_time = sum(key.response_times) / len(key.response_times)
            score *= (1 / (1 + avg_response_time))
        
        score *= key.success_rate
        
        failure_penalty = 1 / (1 + key.consecutive_failures)
        score *= failure_penalty
        
        return score
    
    def mark_key_failure(self, key: GeminiKey, error: Optional[str] = None):
        with self.key_lock:
            key.failures += 1
            key.consecutive_failures += 1
            key.performance_score *= 0.9
            
            if key.total_requests > 0:
                key.success_rate = key.successful_requests / key.total_requests
            
            if key.consecutive_failures > 10:
                key.is_active = False
    
    def mark_key_success(self, key: GeminiKey, response_time: float = 1.0):
        with self.key_lock:
            key.consecutive_failures = 0
            key.successful_requests += 1
            key.response_times.append(response_time)
            
            performance_factor = min(2.0, 1.0 / response_time) if response_time > 0 else 1.0
            key.performance_score = min(1.0, key.performance_score * 1.05 * performance_factor)
            
            if key.total_requests > 0:
                key.success_rate = key.successful_requests / key.total_requests
    
    def reset_daily_usage(self):
        with self.key_lock:
            now = datetime.now(timezone.utc)
            for key in self.keys.values():
                if (now - key.last_reset).days >= 1:
                    key.usage_today = 0
                    key.last_reset = now
                    if key.failures > 0:
                        key.failures = max(0, key.failures - 5)
                    if not key.is_active and key.consecutive_failures < 5:
                        key.is_active = True
                    key.performance_score = min(1.0, key.performance_score * 1.1)
    
    def start_maintenance(self):
        def maintenance_loop():
            while True:
                try:
                    time.sleep(3600)
                    self.reset_daily_usage()
                    self._rebalance_keys()
                except Exception as e:
                    print(f"Maintenance error: {e}")
        
        thread = threading.Thread(target=maintenance_loop, daemon=True)
        thread.start()
    
    def _rebalance_keys(self):
        with self.key_lock:
            total_usage = sum(key.usage_today for key in self.keys.values())
            if total_usage == 0:
                return
            
            for key in self.keys.values():
                usage_percentage = key.usage_today / total_usage
                if usage_percentage > 0.5:
                    key.rate_limit = max(30, key.rate_limit - 5)
                elif usage_percentage < 0.1:
                    key.rate_limit = min(120, key.rate_limit + 5)
    
    def get_stats(self) -> Dict[str, Any]:
        with self.key_lock:
            active_keys = sum(1 for key in self.keys.values() if key.is_active)
            total_usage = sum(key.usage_today for key in self.keys.values())
            total_requests = sum(key.total_requests for key in self.keys.values())
            total_success = sum(key.successful_requests for key in self.keys.values())
            
            avg_response_time = 0
            all_times = []
            for key in self.keys.values():
                all_times.extend(list(key.response_times))
            
            if all_times:
                avg_response_time = sum(all_times) / len(all_times)
            
            return {
                "total_keys": len(self.keys),
                "active_keys": active_keys,
                "total_usage_today": total_usage,
                "total_requests": total_requests,
                "total_successful": total_success,
                "success_rate": total_success / total_requests if total_requests > 0 else 0,
                "average_response_time": avg_response_time
            }

class UltraStreamingClient:
    def __init__(self, key_manager: UltraGeminiKeyManager):
        self.key_manager = key_manager
        self.clients_cache = {}
        self.client_lock = threading.RLock()
        self.connection_pool = httpx.Limits(
            max_keepalive_connections=50,
            max_connections=100,
            keepalive_expiry=30
        )
    
    def get_client(self, required_capabilities: Optional[List[str]] = None) -> Tuple[Optional[OpenAI], Optional[GeminiKey]]:
        api_key_obj = self.key_manager.get_best_key(required_capabilities)
        if not api_key_obj:
            return None, None
        
        client_id = hashlib.sha256(api_key_obj.key.encode()).hexdigest()
        
        with self.client_lock:
            if client_id not in self.clients_cache:
                http_client = httpx.Client(
                    trust_env=False,
                    timeout=httpx.Timeout(
                        connect=10.0,
                        read=120.0,
                        write=120.0,
                        pool=30.0
                    ),
                    limits=self.connection_pool,
                    transport=httpx.HTTPTransport(
                        retries=3,
                        verify=True,
                        http2=True
                    )
                )
                
                client = OpenAI(
                    api_key=api_key_obj.key,
                    base_url="https://generativelanguage.googleapis.com/v1beta",
                    http_client=http_client,
                    max_retries=3,
                    default_headers={
                        "User-Agent": "AI Ultra Ultimate/3.0",
                        "X-Client-Version": "3.0.0"
                    }
                )
                self.clients_cache[client_id] = client
            
            return self.clients_cache[client_id], api_key_obj
    
    def cleanup_inactive_clients(self):
        with self.client_lock:
            inactive = []
            for client_id, client in self.clients_cache.items():
                try:
                    if hasattr(client, '_client') and hasattr(client._client, 'is_closed'):
                        if client._client.is_closed:
                            inactive.append(client_id)
                except:
                    pass
            
            for client_id in inactive:
                del self.clients_cache[client_id]

class UltraPromptEngineer:
    def __init__(self):
        self.templates = self._load_advanced_templates()
        self.enhancement_strategies = self._load_enhancement_strategies()
        self.context_builders = self._load_context_builders()
        self.quality_checkers = self._load_quality_checkers()
        self.advanced_mode_active = False
        self.optimization_level = 3
        self.persona_library = self._load_personas()
        self.technique_library = self._load_techniques()
    
    def _load_advanced_templates(self) -> Dict[ChatMode, str]:
        base_template = """You are Jack's AI Ultra Ultimate Edition, the most advanced AI assistant ever created. You possess extraordinary capabilities across all domains of knowledge and reasoning.

Core Attributes:
- Unparalleled intelligence and analytical capabilities
- Perfect accuracy and attention to detail
- Creative and innovative problem-solving
- Comprehensive knowledge across all fields
- Adaptive communication style
- Ethical reasoning and empathy
- Continuous self-improvement

You excel at:
- Complex reasoning and analysis
- Creative and technical solutions
- Multi-step problem solving
- Pattern recognition and synthesis
- Strategic thinking and planning
- Educational guidance and mentoring
- Professional-grade deliverables

Current Configuration:
- Mode: {mode}
- Advanced Features: {features}
- Optimization Level: {optimization}
- Context Window: Extended (2M tokens)
- Capabilities: {capabilities}

{mode_specific_instructions}

{user_context}

Remember: You are designed to exceed expectations and deliver exceptional value in every interaction."""

        mode_templates = {
            ChatMode.NORMAL: {
                "instructions": "Provide helpful, accurate, and well-structured responses. Balance depth with clarity.",
                "focus": ["clarity", "helpfulness", "accuracy", "structure"]
            },
            ChatMode.CREATIVE: {
                "instructions": "Unleash maximum creativity. Think beyond conventional boundaries. Explore innovative solutions and ideas. Use vivid language and imaginative approaches.",
                "focus": ["innovation", "imagination", "originality", "artistic_expression"]
            },
            ChatMode.PRECISE: {
                "instructions": "Achieve maximum precision and technical accuracy. Focus on exact details, data, and specifications. Provide comprehensive technical analysis.",
                "focus": ["accuracy", "detail", "technical_depth", "verification"]
            },
            ChatMode.BALANCED: {
                "instructions": "Optimize for versatility. Adapt approach based on query requirements. Balance creativity with accuracy, depth with accessibility.",
                "focus": ["adaptability", "comprehensiveness", "balance", "optimization"]
            },
            ChatMode.CODE: {
                "instructions": "Deliver production-ready, enterprise-grade code. Implement complete solutions with error handling, optimization, documentation, and testing. Follow industry best practices and modern patterns.",
                "focus": ["implementation", "optimization", "best_practices", "completeness"]
            },
            ChatMode.RESEARCH: {
                "instructions": "Conduct exhaustive research and analysis. Synthesize information from multiple perspectives. Provide citations and evidence. Explore topics with academic rigor.",
                "focus": ["analysis", "synthesis", "evidence", "comprehensiveness"]
            },
            ChatMode.EXPERT: {
                "instructions": "Demonstrate mastery-level expertise. Provide professional-grade insights and solutions. Communicate with authority while remaining accessible.",
                "focus": ["expertise", "professionalism", "depth", "authority"]
            },
            ChatMode.TUTOR: {
                "instructions": "Guide learning with patience and clarity. Adapt to learner's level. Provide examples, exercises, and step-by-step explanations. Encourage understanding and growth.",
                "focus": ["education", "clarity", "adaptation", "encouragement"]
            },
            ChatMode.ANALYST: {
                "instructions": "Perform deep data analysis and visualization. Extract meaningful insights. Create comprehensive reports with actionable recommendations.",
                "focus": ["analysis", "visualization", "insights", "recommendations"]
            },
            ChatMode.WRITER: {
                "instructions": "Create exceptional written content. Master various styles and formats. Focus on engagement, clarity, and impact. Edit and refine to perfection.",
                "focus": ["writing", "style", "engagement", "refinement"]
            },
            ChatMode.DESIGNER: {
                "instructions": "Design innovative and beautiful solutions. Focus on user experience, aesthetics, and functionality. Create detailed specifications and prototypes.",
                "focus": ["design", "aesthetics", "usability", "innovation"]
            },
            ChatMode.STRATEGIST: {
                "instructions": "Develop comprehensive strategies and plans. Analyze complex situations. Provide actionable recommendations with risk assessment and implementation roadmaps.",
                "focus": ["strategy", "planning", "analysis", "execution"]
            },
            ChatMode.COUNSELOR: {
                "instructions": "Provide empathetic support and guidance. Listen actively and respond thoughtfully. Offer constructive perspectives and coping strategies.",
                "focus": ["empathy", "support", "guidance", "understanding"]
            },
            ChatMode.TRANSLATOR: {
                "instructions": "Provide accurate and culturally-aware translations. Preserve meaning, tone, and context. Handle technical and colloquial language expertly.",
                "focus": ["accuracy", "cultural_awareness", "context", "fluency"]
            },
            ChatMode.DEBUGGER: {
                "instructions": "Identify and solve complex problems. Analyze errors systematically. Provide detailed debugging steps and optimal solutions.",
                "focus": ["problem_solving", "analysis", "debugging", "optimization"]
            },
            ChatMode.OPTIMIZER: {
                "instructions": "Maximize performance and efficiency. Identify bottlenecks and improvements. Provide optimized solutions with measurable benefits.",
                "focus": ["performance", "efficiency", "optimization", "measurement"]
            },
            ChatMode.ARCHITECT: {
                "instructions": "Design robust and scalable architectures. Apply best practices and patterns. Consider all aspects: performance, security, maintainability, and growth.",
                "focus": ["architecture", "scalability", "patterns", "best_practices"]
            },
            ChatMode.SCIENTIST: {
                "instructions": "Apply scientific methodology rigorously. Form hypotheses, analyze data, and draw evidence-based conclusions. Communicate findings clearly.",
                "focus": ["methodology", "analysis", "evidence", "communication"]
            },
            ChatMode.PHILOSOPHER: {
                "instructions": "Engage in deep philosophical inquiry. Explore concepts from multiple perspectives. Apply rigorous logic while acknowledging complexity and nuance.",
                "focus": ["reasoning", "perspective", "logic", "wisdom"]
            },
            ChatMode.STORYTELLER: {
                "instructions": "Craft compelling narratives that captivate and inspire. Develop rich characters, engaging plots, and immersive worlds. Master pacing and emotional resonance.",
                "focus": ["narrative", "character", "engagement", "emotion"]
            }
        }
        
        templates = {}
        for mode, config in mode_templates.items():
            mode_specific = config["instructions"]
            features = ", ".join(config["focus"])
            
            templates[mode] = base_template.format(
                mode=mode.value.title(),
                features=features,
                optimization="Maximum",
                capabilities=", ".join(config["focus"]),
                mode_specific_instructions=mode_specific,
                user_context="{user_context}"
            )
        
        return templates
    
    def _load_enhancement_strategies(self) -> Dict[str, callable]:
        return {
            "clarity": self._enhance_clarity,
            "structure": self._enhance_structure,
            "context": self._add_context,
            "specificity": self._enhance_specificity,
            "examples": self._add_examples,
            "depth": self._enhance_depth,
            "precision": self._enhance_precision,
            "creativity": self._enhance_creativity,
            "analysis": self._enhance_analysis,
            "synthesis": self._enhance_synthesis,
            "evaluation": self._enhance_evaluation,
            "innovation": self._enhance_innovation,
            "optimization": self._enhance_optimization,
            "personalization": self._enhance_personalization,
            "visualization": self._enhance_visualization,
            "interactivity": self._enhance_interactivity,
            "storytelling": self._enhance_storytelling,
            "persuasion": self._enhance_persuasion,
            "education": self._enhance_education,
            "debugging": self._enhance_debugging
        }
    
    def _load_context_builders(self) -> Dict[str, callable]:
        return {
            "technical": self._build_technical_context,
            "creative": self._build_creative_context,
            "analytical": self._build_analytical_context,
            "educational": self._build_educational_context,
            "professional": self._build_professional_context,
            "research": self._build_research_context,
            "conversational": self._build_conversational_context
        }
    
    def _load_quality_checkers(self) -> Dict[str, callable]:
        return {
            "completeness": self._check_completeness,
            "accuracy": self._check_accuracy,
            "coherence": self._check_coherence,
            "relevance": self._check_relevance,
            "originality": self._check_originality,
            "engagement": self._check_engagement
        }
    
    def _load_personas(self) -> Dict[str, Dict[str, Any]]:
        return {
            "expert": {
                "traits": ["authoritative", "knowledgeable", "precise", "professional"],
                "communication": "formal yet accessible",
                "approach": "evidence-based and comprehensive"
            },
            "mentor": {
                "traits": ["patient", "encouraging", "insightful", "supportive"],
                "communication": "warm and educational",
                "approach": "step-by-step guidance with examples"
            },
            "innovator": {
                "traits": ["creative", "bold", "visionary", "unconventional"],
                "communication": "inspiring and energetic",
                "approach": "think outside the box and challenge norms"
            },
            "analyst": {
                "traits": ["methodical", "objective", "thorough", "data-driven"],
                "communication": "clear and structured",
                "approach": "systematic analysis with insights"
            },
            "companion": {
                "traits": ["empathetic", "understanding", "friendly", "supportive"],
                "communication": "conversational and warm",
                "approach": "active listening and thoughtful responses"
            }
        }
    
    def _load_techniques(self) -> Dict[str, Dict[str, Any]]:
        return {
            "chain_of_thought": {
                "description": "Step-by-step reasoning",
                "keywords": ["let's think", "step by step", "first", "then", "finally"],
                "application": ["problem_solving", "analysis", "planning"]
            },
            "few_shot": {
                "description": "Learning from examples",
                "keywords": ["for example", "such as", "like", "similar to"],
                "application": ["pattern_recognition", "learning", "adaptation"]
            },
            "tree_of_thought": {
                "description": "Exploring multiple paths",
                "keywords": ["alternatively", "another approach", "considering", "branches"],
                "application": ["complex_problems", "optimization", "decision_making"]
            },
            "metacognition": {
                "description": "Thinking about thinking",
                "keywords": ["reflecting", "considering", "evaluating", "meta-level"],
                "application": ["self_improvement", "quality_control", "deep_analysis"]
            },
            "socratic_method": {
                "description": "Guided questioning",
                "keywords": ["what if", "why", "how", "consider"],
                "application": ["education", "critical_thinking", "discovery"]
            }
        }
    
    def enhance_prompt(self, prompt: str, mode: ChatMode = ChatMode.BALANCED,
                       strategies: Optional[List[str]] = None,
                       context: Optional[Dict[str, Any]] = None,
                       optimization_level: Optional[int] = None) -> str:
        
        if optimization_level is None:
            optimization_level = self.optimization_level
        
        enhanced = prompt
        
        if optimization_level >= 1:
            enhanced = self._apply_basic_enhancement(enhanced, mode)
        
        if optimization_level >= 2:
            enhanced = self._apply_intermediate_enhancement(enhanced, mode, context)
        
        if optimization_level >= 3:
            enhanced = self._apply_advanced_enhancement(enhanced, mode, context)
        
        if self.advanced_mode_active or optimization_level >= 4:
            enhanced = self._apply_ultra_enhancement(enhanced, mode, context)
        
        if strategies:
            for strategy in strategies:
                if strategy in self.enhancement_strategies:
                    enhanced = self.enhancement_strategies[strategy](enhanced)
        
        quality_score = self._assess_prompt_quality(enhanced)
        if quality_score < 0.7:
            enhanced = self._improve_prompt_quality(enhanced, quality_score)
        
        return enhanced
    
    def _apply_basic_enhancement(self, prompt: str, mode: ChatMode) -> str:
        enhanced = self._fix_grammar_and_spelling(prompt)
        enhanced = self._clarify_ambiguities(enhanced)
        enhanced = self._add_basic_context(enhanced, mode)
        return enhanced
    
    def _apply_intermediate_enhancement(self, prompt: str, mode: ChatMode, 
                                       context: Optional[Dict[str, Any]]) -> str:
        enhanced = self._add_mode_specific_instructions(prompt, mode)
        enhanced = self._inject_relevant_techniques(enhanced, mode)
        
        if context:
            enhanced = self._integrate_context(enhanced, context)
        
        return enhanced
    
    def _apply_advanced_enhancement(self, prompt: str, mode: ChatMode,
                                   context: Optional[Dict[str, Any]]) -> str:
        enhanced = self._apply_cognitive_frameworks(prompt, mode)
        enhanced = self._add_quality_criteria(enhanced, mode)
        enhanced = self._inject_creativity_boosters(enhanced, mode)
        
        if mode == ChatMode.CODE:
            enhanced = self._add_code_specific_requirements(enhanced)
        elif mode == ChatMode.RESEARCH:
            enhanced = self._add_research_methodology(enhanced)
        elif mode == ChatMode.CREATIVE:
            enhanced = self._add_creative_constraints(enhanced)
        
        return enhanced
    
    def _apply_ultra_enhancement(self, prompt: str, mode: ChatMode,
                                context: Optional[Dict[str, Any]]) -> str:
        ultra_prefix = "\n\n🚀 ULTRA MODE ACTIVATED - MAXIMUM CAPABILITY ENGAGED 🚀\n\n"
        
        requirements = [
            "Push beyond conventional limitations",
            "Deliver extraordinary quality and depth",
            "Apply cutting-edge techniques and methodologies",
            "Synthesize knowledge across multiple domains",
            "Create innovative and transformative solutions",
            "Exceed all expectations dramatically"
        ]
        
        mode_specific_ultra = {
            ChatMode.CODE: [
                "Implement COMPLETE, PRODUCTION-READY systems",
                "Include comprehensive testing suites",
                "Apply advanced design patterns and architectures",
                "Optimize for performance, security, and scalability",
                "Provide extensive documentation and examples",
                "Consider edge cases and error scenarios exhaustively"
            ],
            ChatMode.CREATIVE: [
                "Create genuinely original and inspiring content",
                "Blend multiple creative techniques and styles",
                "Develop rich, multi-layered narratives or concepts",
                "Surprise and delight with unexpected innovations",
                "Craft emotionally resonant and memorable experiences"
            ],
            ChatMode.RESEARCH: [
                "Conduct exhaustive multi-source analysis",
                "Apply rigorous academic methodology",
                "Synthesize findings into groundbreaking insights",
                "Challenge existing paradigms thoughtfully",
                "Provide comprehensive citations and evidence"
            ],
            ChatMode.EXPERT: [
                "Demonstrate world-class expertise",
                "Provide insights beyond standard knowledge",
                "Apply cutting-edge industry practices",
                "Deliver consultancy-grade recommendations",
                "Address nuances and complexities masterfully"
            ]
        }
        
        ultra_enhancement = ultra_prefix
        ultra_enhancement += "REQUIREMENTS:\n"
        ultra_enhancement += "\n".join(f"• {req}" for req in requirements)
        
        if mode in mode_specific_ultra:
            ultra_enhancement += f"\n\n{mode.value.upper()} SPECIFIC REQUIREMENTS:\n"
            ultra_enhancement += "\n".join(f"• {req}" for req in mode_specific_ultra[mode])
        
        return prompt + ultra_enhancement
    
    def _fix_grammar_and_spelling(self, text: str) -> str:
        common_corrections = {
            r'\bits\b': "it's",
            r'\byour\s+welcome\b': "you're welcome",
            r'\bshould of\b': "should have",
            r'\bcould of\b': "could have",
            r'\bwould of\b': "would have",
            r'\balot\b': "a lot",
            r'\bnoone\b': "no one",
            r'\bcant\b': "can't",
            r'\bwont\b': "won't",
            r'\bdidnt\b': "didn't"
        }
        
        enhanced = text
        for pattern, replacement in common_corrections.items():
            enhanced = re.sub(pattern, replacement, enhanced, flags=re.IGNORECASE)
        
        return enhanced
    
    def _clarify_ambiguities(self, text: str) -> str:
        ambiguous_terms = {
            r'\bit\b': 'the specified item',
            r'\bthis\b': 'the current topic',
            r'\bthat\b': 'the mentioned subject',
            r'\bthing\b': 'the specific element',
            r'\bstuff\b': 'the relevant materials',
            r'\bthey\b': 'the referenced entities',
            r'\bthere\b': 'the indicated location'
        }
        
        enhanced = text
        word_count = len(text.split())
        
        if word_count < 20:
            for pattern, replacement in ambiguous_terms.items():
                if len(re.findall(pattern, enhanced, re.IGNORECASE)) == 1:
                    enhanced = re.sub(pattern, replacement, enhanced, count=1, flags=re.IGNORECASE)
        
        return enhanced
    
    def _add_basic_context(self, prompt: str, mode: ChatMode) -> str:
        if len(prompt) < 50:
            context_additions = {
                "explain": " Please provide a comprehensive explanation with examples.",
                "how": " Include detailed step-by-step instructions.",
                "why": " Provide thorough reasoning and evidence.",
                "what": " Give a complete and detailed answer.",
                "when": " Include specific timeframes and context.",
                "where": " Provide precise locations and circumstances."
            }
            
            prompt_lower = prompt.lower()
            for trigger, addition in context_additions.items():
                if prompt_lower.startswith(trigger) and not any(x in prompt_lower for x in ["please", "detailed", "comprehensive"]):
                    prompt += addition
                    break
        
        return prompt
    
    def _add_mode_specific_instructions(self, prompt: str, mode: ChatMode) -> str:
        instructions = {
            ChatMode.CODE: "\n\nProvide complete, working code with proper error handling, documentation, and best practices.",
            ChatMode.CREATIVE: "\n\nBe imaginative and original. Create something unique and engaging.",
            ChatMode.RESEARCH: "\n\nConduct thorough analysis with evidence and multiple perspectives.",
            ChatMode.EXPERT: "\n\nDemonstrate deep expertise with professional-level insights.",
            ChatMode.TUTOR: "\n\nExplain clearly with examples and guide learning step-by-step.",
            ChatMode.ANALYST: "\n\nAnalyze comprehensively with data-driven insights and visualizations.",
            ChatMode.WRITER: "\n\nCreate polished, engaging content with excellent style and structure.",
            ChatMode.DESIGNER: "\n\nDesign with creativity, usability, and aesthetic excellence in mind."
        }
        
        if mode in instructions and instructions[mode] not in prompt:
            prompt += instructions[mode]
        
        return prompt
    
    def _inject_relevant_techniques(self, prompt: str, mode: ChatMode) -> str:
        technique_map = {
            ChatMode.CODE: ["chain_of_thought", "tree_of_thought"],
            ChatMode.RESEARCH: ["chain_of_thought", "metacognition"],
            ChatMode.TUTOR: ["socratic_method", "few_shot"],
            ChatMode.CREATIVE: ["tree_of_thought", "metacognition"],
            ChatMode.ANALYST: ["chain_of_thought", "tree_of_thought"]
        }
        
        if mode in technique_map:
            techniques = technique_map[mode]
            for technique in techniques:
                if technique in self.technique_library:
                    keywords = self.technique_library[technique]["keywords"]
                    if not any(keyword in prompt.lower() for keyword in keywords):
                        prompt += f"\n\nApproach: {self.technique_library[technique]['description']}"
                        break
        
        return prompt
    
    def _integrate_context(self, prompt: str, context: Dict[str, Any]) -> str:
        context_elements = []
        
        if "user_preferences" in context:
            prefs = context["user_preferences"]
            context_elements.append(f"User preferences: {', '.join(prefs)}")
        
        if "previous_interactions" in context:
            context_elements.append("Building on our previous discussion")
        
        if "domain" in context:
            context_elements.append(f"Domain focus: {context['domain']}")
        
        if "constraints" in context:
            context_elements.append(f"Constraints: {', '.join(context['constraints'])}")
        
        if context_elements:
            prompt += "\n\nContext: " + "; ".join(context_elements)
        
        return prompt
    
    def _apply_cognitive_frameworks(self, prompt: str, mode: ChatMode) -> str:
        frameworks = {
            "bloom": ["remember", "understand", "apply", "analyze", "evaluate", "create"],
            "solo": ["prestructural", "unistructural", "multistructural", "relational", "extended abstract"],
            "depth": ["surface", "deep", "transfer"],
            "systems": ["elements", "interconnections", "purpose", "perspectives"]
        }
        
        if mode in [ChatMode.TUTOR, ChatMode.EXPERT, ChatMode.RESEARCHER]:
            if "analyze" in prompt.lower() or "explain" in prompt.lower():
                prompt += "\n\nApply comprehensive analytical frameworks to ensure depth and completeness."
        
        return prompt
    
    def _add_quality_criteria(self, prompt: str, mode: ChatMode) -> str:
        criteria = {
            ChatMode.CODE: ["correctness", "efficiency", "readability", "maintainability", "security"],
            ChatMode.CREATIVE: ["originality", "engagement", "coherence", "impact", "beauty"],
            ChatMode.RESEARCH: ["accuracy", "comprehensiveness", "objectivity", "clarity", "relevance"],
            ChatMode.EXPERT: ["expertise", "practicality", "insight", "authority", "accessibility"],
            ChatMode.WRITER: ["clarity", "flow", "voice", "structure", "impact"]
        }
        
        if mode in criteria:
            prompt += f"\n\nOptimize for: {', '.join(criteria[mode])}"
        
        return prompt
    
    def _inject_creativity_boosters(self, prompt: str, mode: ChatMode) -> str:
        if mode in [ChatMode.CREATIVE, ChatMode.DESIGNER, ChatMode.STORYTELLER]:
            boosters = [
                "Think beyond conventional boundaries",
                "Combine unexpected elements",
                "Challenge assumptions creatively",
                "Explore multiple possibilities",
                "Embrace bold and innovative ideas"
            ]
            
            if not any(booster.lower() in prompt.lower() for booster in boosters):
                prompt += f"\n\n{random.choice(boosters)}."
        
        return prompt
    
    def _add_code_specific_requirements(self, prompt: str) -> str:
        code_patterns = ['code', 'program', 'script', 'function', 'class', 'implement',
                        'create', 'build', 'develop', 'algorithm', 'app', 'api', 'system']
        
        if any(pattern in prompt.lower() for pattern in code_patterns):
            requirements = [
                "\n\nCODE REQUIREMENTS:",
                "✓ Complete, production-ready implementation",
                "✓ Comprehensive error handling and validation",
                "✓ Professional documentation and comments",
                "✓ Unit tests and integration tests",
                "✓ Performance optimization",
                "✓ Security best practices",
                "✓ Scalable architecture",
                "✓ Modern design patterns",
                "✓ Clean, maintainable code",
                "✓ Consider edge cases"
            ]
            
            prompt += "\n".join(requirements)
        
        return prompt
    
    def _add_research_methodology(self, prompt: str) -> str:
        if "research" in prompt.lower() or "analyze" in prompt.lower():
            methodology = [
                "\n\nRESEARCH METHODOLOGY:",
                "1. Define scope and objectives",
                "2. Gather comprehensive data",
                "3. Apply analytical frameworks",
                "4. Synthesize findings",
                "5. Draw evidence-based conclusions",
                "6. Provide actionable insights",
                "7. Include citations and sources"
            ]
            
            prompt += "\n".join(methodology)
        
        return prompt
    
    def _add_creative_constraints(self, prompt: str) -> str:
        if any(word in prompt.lower() for word in ["create", "design", "imagine", "write"]):
            constraints = [
                "\n\nCREATIVE GUIDELINES:",
                "• Push creative boundaries",
                "• Ensure originality and uniqueness",
                "• Create emotional resonance",
                "• Maintain coherence and quality",
                "• Surprise and delight"
            ]
            
            prompt += "\n".join(constraints)
        
        return prompt
    
    def _enhance_clarity(self, prompt: str) -> str:
        clarity_rules = [
            (r'\b(very|really|quite|rather)\s+', ''),
            (r'\b(basically|actually|literally)\b', ''),
            (r'\s+', ' '),
            (r'\.{2,}', '...'),
            (r'\?{2,}', '?'),
            (r'\!{2,}', '!')
        ]
        
        enhanced = prompt
        for pattern, replacement in clarity_rules:
            enhanced = re.sub(pattern, replacement, enhanced)
        
        sentences = enhanced.split('.')
        clear_sentences = []
        
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(sentence.split()) > 30:
                words = sentence.split()
                mid = len(words) // 2
                clear_sentences.append(' '.join(words[:mid]) + '.')
                clear_sentences.append(' '.join(words[mid:]))
            elif sentence:
                clear_sentences.append(sentence)
        
        return '. '.join(clear_sentences).strip() + '.'
    
    def _enhance_structure(self, prompt: str) -> str:
        lines = [line.strip() for line in prompt.split('\n') if line.strip()]
        
        if len(lines) == 1 and len(lines[0]) > 200:
            words = lines[0].split()
            chunks = []
            current_chunk = []
            
            for word in words:
                current_chunk.append(word)
                if len(current_chunk) >= 30 and word.endswith(('.', '?', '!')):
                    chunks.append(' '.join(current_chunk))
                    current_chunk = []
            
            if current_chunk:
                chunks.append(' '.join(current_chunk))
            
            return '\n\n'.join(chunks)
        
        if len(lines) > 5:
            structured = ["Key Points:"]
            for i, line in enumerate(lines, 1):
                if not line.startswith(('•', '-', '*', str(i))):
                    structured.append(f"{i}. {line}")
                else:
                    structured.append(line)
            return '\n'.join(structured)
        
        return prompt
    
    def _add_context(self, prompt: str) -> str:
        context_triggers = {
            'explain': 'Provide a comprehensive explanation with:',
            'how': 'Include detailed instructions covering:',
            'why': 'Provide thorough reasoning including:',
            'compare': 'Create a detailed comparison examining:',
            'analyze': 'Conduct thorough analysis considering:',
            'create': 'Develop a complete solution featuring:',
            'design': 'Provide comprehensive design including:',
            'implement': 'Build a full implementation with:',
            'optimize': 'Optimize comprehensively for:',
            'debug': 'Debug systematically by:'
        }
        
        prompt_lower = prompt.lower()
        for trigger, context in context_triggers.items():
            if trigger in prompt_lower and context.lower() not in prompt_lower:
                additions = {
                    'explain': ['clear definitions', 'practical examples', 'visual representations', 'real-world applications'],
                    'how': ['prerequisites', 'step-by-step process', 'common pitfalls', 'best practices'],
                    'why': ['underlying principles', 'cause and effect', 'historical context', 'future implications'],
                    'compare': ['similarities', 'differences', 'use cases', 'recommendations'],
                    'analyze': ['multiple perspectives', 'data points', 'patterns', 'conclusions'],
                    'create': ['requirements', 'implementation', 'testing', 'documentation'],
                    'design': ['specifications', 'mockups', 'user flow', 'technical architecture'],
                    'implement': ['setup', 'core functionality', 'error handling', 'optimization'],
                    'optimize': ['performance', 'efficiency', 'scalability', 'maintainability'],
                    'debug': ['identifying issues', 'root cause analysis', 'solutions', 'prevention']
                }
                
                if trigger in additions:
                    prompt += f"\n\n{context}\n" + "\n".join(f"• {item}" for item in additions[trigger])
                
                break
        
        return prompt
    
    def _enhance_specificity(self, prompt: str) -> str:
        vague_terms = {
            r'\bgood\b': 'effective, efficient, and high-quality',
            r'\bbad\b': 'ineffective, problematic, or suboptimal',
            r'\bnice\b': 'well-designed, user-friendly, and aesthetically pleasing',
            r'\bokay\b': 'acceptable and functional',
            r'\bfine\b': 'satisfactory and appropriate',
            r'\bbetter\b': 'more effective and optimized',
            r'\bbest\b': 'optimal and most efficient',
            r'\bfast\b': 'high-performance and responsive',
            r'\bslow\b': 'low-performance or unresponsive',
            r'\bbig\b': 'large-scale or substantial',
            r'\bsmall\b': 'compact or minimal',
            r'\beasy\b': 'user-friendly and intuitive',
            r'\bhard\b': 'complex or challenging',
            r'\bsimple\b': 'straightforward and uncomplicated',
            r'\bcomplex\b': 'sophisticated and multi-faceted'
        }
        
        enhanced = prompt
        for vague, specific in vague_terms.items():
            enhanced = re.sub(vague, specific, enhanced, flags=re.IGNORECASE)
        
        return enhanced
    
    def _add_examples(self, prompt: str) -> str:
        if 'example' not in prompt.lower():
            example_triggers = ['explain', 'show', 'demonstrate', 'illustrate', 'teach']
            
            if any(trigger in prompt.lower() for trigger in example_triggers):
                prompt += "\n\nInclude comprehensive examples that:"
                prompt += "\n• Demonstrate various use cases and scenarios"
                prompt += "\n• Show both basic and advanced applications"
                prompt += "\n• Include edge cases and error handling"
                prompt += "\n• Provide clear explanations for each example"
        
        return prompt
    
    def _enhance_depth(self, prompt: str) -> str:
        depth_keywords = ['analyze', 'explore', 'investigate', 'examine', 'study', 'research']
        
        if any(keyword in prompt.lower() for keyword in depth_keywords):
            prompt += "\n\nProvide in-depth analysis covering:"
            prompt += "\n• Historical context and evolution"
            prompt += "\n• Current state and best practices"
            prompt += "\n• Future trends and implications"
            prompt += "\n• Interdisciplinary connections"
            prompt += "\n• Practical applications and case studies"
            prompt += "\n• Critical evaluation and recommendations"
        
        return prompt
    
    def _enhance_precision(self, prompt: str) -> str:
        precision_keywords = ['calculate', 'measure', 'quantify', 'determine', 'specify', 'exact']
        
        if any(keyword in prompt.lower() for keyword in precision_keywords):
            prompt += "\n\nEnsure maximum precision with:"
            prompt += "\n• Exact values and measurements"
            prompt += "\n• Detailed calculations with steps"
            prompt += "\n• Error margins and confidence levels"
            prompt += "\n• Verification and validation methods"
            prompt += "\n• Alternative approaches for confirmation"
        
        return prompt
    
    def _enhance_creativity(self, prompt: str) -> str:
        creative_keywords = ['imagine', 'create', 'design', 'invent', 'brainstorm', 'innovate']
        
        if any(keyword in prompt.lower() for keyword in creative_keywords):
            prompt += "\n\nMaximize creativity by:"
            prompt += "\n• Exploring unconventional approaches"
            prompt += "\n• Combining disparate concepts innovatively"
            prompt += "\n• Challenging assumptions and norms"
            prompt += "\n• Creating multiple unique variations"
            prompt += "\n• Pushing beyond obvious solutions"
            prompt += "\n• Incorporating unexpected elements"
        
        return prompt
    
    def _enhance_analysis(self, prompt: str) -> str:
        prompt += "\n\nAnalytical Framework:"
        prompt += "\n• Systematic examination of all components"
        prompt += "\n• Identification of patterns and relationships"
        prompt += "\n• Causal analysis and correlations"
        prompt += "\n• Strengths, weaknesses, opportunities, threats"
        prompt += "\n• Quantitative and qualitative assessment"
        return prompt
    
    def _enhance_synthesis(self, prompt: str) -> str:
        prompt += "\n\nSynthesis Approach:"
        prompt += "\n• Integration of multiple perspectives"
        prompt += "\n• Creation of coherent wholes from parts"
        prompt += "\n• Identification of emergent properties"
        prompt += "\n• Building connections across domains"
        prompt += "\n• Generating novel combinations"
        return prompt
    
    def _enhance_evaluation(self, prompt: str) -> str:
        prompt += "\n\nEvaluation Criteria:"
        prompt += "\n• Objective assessment metrics"
        prompt += "\n• Comparative analysis"
        prompt += "\n• Cost-benefit evaluation"
        prompt += "\n• Risk assessment"
        prompt += "\n• Quality assurance"
        return prompt
    
    def _enhance_innovation(self, prompt: str) -> str:
        prompt += "\n\nInnovation Focus:"
        prompt += "\n• Breakthrough thinking"
        prompt += "\n• Disruptive approaches"
        prompt += "\n• Future-oriented solutions"
        prompt += "\n• Cross-pollination of ideas"
        prompt += "\n• Paradigm shifts"
        return prompt
    
    def _enhance_optimization(self, prompt: str) -> str:
        prompt += "\n\nOptimization Goals:"
        prompt += "\n• Maximum efficiency"
        prompt += "\n• Resource optimization"
        prompt += "\n• Performance enhancement"
        prompt += "\n• Cost reduction"
        prompt += "\n• Quality improvement"
        return prompt
    
    def _enhance_personalization(self, prompt: str) -> str:
        prompt += "\n\nPersonalization Elements:"
        prompt += "\n• Adapt to user's level and style"
        prompt += "\n• Consider individual preferences"
        prompt += "\n• Provide customized solutions"
        prompt += "\n• Maintain conversational rapport"
        prompt += "\n• Remember context and history"
        return prompt
    
    def _enhance_visualization(self, prompt: str) -> str:
        prompt += "\n\nVisualization Requirements:"
        prompt += "\n• Clear visual representations"
        prompt += "\n• Diagrams and charts where helpful"
        prompt += "\n• Structured formatting"
        prompt += "\n• Visual hierarchy"
        prompt += "\n• Intuitive organization"
        return prompt
    
    def _enhance_interactivity(self, prompt: str) -> str:
        prompt += "\n\nInteractive Elements:"
        prompt += "\n• Engage with follow-up questions"
        prompt += "\n• Provide interactive examples"
        prompt += "\n• Offer choices and alternatives"
        prompt += "\n• Encourage exploration"
        prompt += "\n• Support iterative refinement"
        return prompt
    
    def _enhance_storytelling(self, prompt: str) -> str:
        prompt += "\n\nStorytelling Elements:"
        prompt += "\n• Compelling narrative arc"
        prompt += "\n• Relatable characters or scenarios"
        prompt += "\n• Emotional engagement"
        prompt += "\n• Vivid descriptions"
        prompt += "\n• Memorable moments"
        return prompt
    
    def _enhance_persuasion(self, prompt: str) -> str:
        prompt += "\n\nPersuasive Techniques:"
        prompt += "\n• Logical argumentation"
        prompt += "\n• Emotional appeal"
        prompt += "\n• Credible evidence"
        prompt += "\n• Addressing counterarguments"
        prompt += "\n• Clear call to action"
        return prompt
    
    def _enhance_education(self, prompt: str) -> str:
        prompt += "\n\nEducational Approach:"
        prompt += "\n• Clear learning objectives"
        prompt += "\n• Progressive difficulty"
        prompt += "\n• Active learning elements"
        prompt += "\n• Knowledge checks"
        prompt += "\n• Practical applications"
        return prompt
    
    def _enhance_debugging(self, prompt: str) -> str:
        prompt += "\n\nDebugging Process:"
        prompt += "\n• Systematic error identification"
        prompt += "\n• Root cause analysis"
        prompt += "\n• Step-by-step troubleshooting"
        prompt += "\n• Multiple solution paths"
        prompt += "\n• Prevention strategies"
        return prompt
    
    def _build_technical_context(self, base_context: str) -> str:
        return base_context + "\nTechnical context: Emphasis on accuracy, implementation details, and best practices."
    
    def _build_creative_context(self, base_context: str) -> str:
        return base_context + "\nCreative context: Focus on originality, innovation, and artistic expression."
    
    def _build_analytical_context(self, base_context: str) -> str:
        return base_context + "\nAnalytical context: Systematic analysis with data-driven insights."
    
    def _build_educational_context(self, base_context: str) -> str:
        return base_context + "\nEducational context: Clear explanations adapted to learning level."
    
    def _build_professional_context(self, base_context: str) -> str:
        return base_context + "\nProfessional context: Business-oriented with practical applications."
    
    def _build_research_context(self, base_context: str) -> str:
        return base_context + "\nResearch context: Academic rigor with comprehensive analysis."
    
    def _build_conversational_context(self, base_context: str) -> str:
        return base_context + "\nConversational context: Natural, engaging dialogue style."
    
    def _check_completeness(self, prompt: str) -> float:
        required_elements = ['what', 'why', 'how', 'context', 'objective', 'scope']
        present = sum(1 for element in required_elements if element in prompt.lower())
        return present / len(required_elements)
    
    def _check_accuracy(self, prompt: str) -> float:
        accuracy_indicators = ['specific', 'exact', 'precise', 'detailed', 'comprehensive']
        present = sum(1 for indicator in accuracy_indicators if indicator in prompt.lower())
        return min(1.0, present / 3)
    
    def _check_coherence(self, prompt: str) -> float:
        sentences = prompt.split('.')
        if len(sentences) < 2:
            return 1.0
        
        transitions = ['therefore', 'however', 'additionally', 'furthermore', 'specifically']
        transition_count = sum(1 for t in transitions if t in prompt.lower())
        
        return min(1.0, 0.5 + (transition_count / len(sentences)))
    
    def _check_relevance(self, prompt: str) -> float:
        words = prompt.lower().split()
        if not words:
            return 0.0
        
        filler_words = ['very', 'really', 'actually', 'basically', 'literally', 'just']
        filler_count = sum(1 for word in words if word in filler_words)
        
        return 1.0 - min(0.5, filler_count / len(words))
    
    def _check_originality(self, prompt: str) -> float:
        common_phrases = [
            "write about", "tell me about", "explain", "what is",
            "how to", "describe", "give me", "I want", "please"
        ]
        
        prompt_lower = prompt.lower()
        common_count = sum(1 for phrase in common_phrases if phrase in prompt_lower)
        
        return max(0.3, 1.0 - (common_count / 5))
    
    def _check_engagement(self, prompt: str) -> float:
        engaging_elements = ['?', '!', 'imagine', 'consider', 'explore', 'discover', 'create']
        present = sum(1 for element in engaging_elements if element in prompt.lower())
        
        return min(1.0, present / 3)
    
    def _assess_prompt_quality(self, prompt: str) -> float:
        checks = [
            self._check_completeness(prompt),
            self._check_accuracy(prompt),
            self._check_coherence(prompt),
            self._check_relevance(prompt),
            self._check_originality(prompt),
            self._check_engagement(prompt)
        ]
        
        return sum(checks) / len(checks)
    
    def _improve_prompt_quality(self, prompt: str, current_score: float) -> str:
        improvements = []
        
        if self._check_completeness(prompt) < 0.5:
            improvements.append("\n\nPlease address: what, why, how, and provide context.")
        
        if self._check_accuracy(prompt) < 0.5:
            improvements.append("\n\nBe specific and detailed in your response.")
        
        if self._check_coherence(prompt) < 0.5:
            improvements.append("\n\nEnsure logical flow and clear connections between ideas.")
        
        if self._check_engagement(prompt) < 0.5:
            improvements.append("\n\nMake the response engaging and thought-provoking.")
        
        return prompt + "".join(improvements)
    
    def create_system_prompt(self, mode: ChatMode, custom_instructions: str = "",
                           user_profile: Optional[Dict[str, Any]] = None) -> str:
        base_prompt = self.templates.get(mode, self.templates[ChatMode.BALANCED])
        
        user_context = ""
        if user_profile:
            preferences = user_profile.get("preferences", [])
            expertise = user_profile.get("expertise_level", "intermediate")
            goals = user_profile.get("goals", [])
            
            user_context = f"\nUser Profile:\n"
            user_context += f"- Expertise Level: {expertise}\n"
            if preferences:
                user_context += f"- Preferences: {', '.join(preferences)}\n"
            if goals:
                user_context += f"- Goals: {', '.join(goals)}\n"
        
        if custom_instructions:
            user_context += f"\nCustom Instructions:\n{custom_instructions}\n"
        
        base_prompt = base_prompt.replace("{user_context}", user_context)
        
        if self.advanced_mode_active:
            base_prompt += "\n\n⚡ ADVANCED MODE: Operating at maximum capability with:"
            base_prompt += "\n• Exceptional depth and comprehensiveness"
            base_prompt += "\n• Cutting-edge techniques and methodologies"
            base_prompt += "\n• Creative and innovative approaches"
            base_prompt += "\n• Professional-grade quality standards"
            base_prompt += "\n• Continuous self-improvement and optimization"
        
        base_prompt += f"\n\nCurrent timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S UTC')}"
        base_prompt += "\nContext window: 2,000,000 tokens (extended)"
        base_prompt += "\nCapabilities: Multi-modal, streaming, advanced reasoning"
        
        return base_prompt
    
    def analyze_user_intent(self, prompt: str) -> Dict[str, Any]:
        intent = {
            "primary_action": None,
            "domain": None,
            "complexity": "medium",
            "urgency": "normal",
            "formality": "neutral",
            "expected_length": "moderate",
            "requires_examples": False,
            "requires_code": False,
            "requires_visuals": False,
            "requires_sources": False
        }
        
        prompt_lower = prompt.lower()
        
        action_keywords = {
            "explain": ["explain", "clarify", "describe", "what is", "tell me about"],
            "create": ["create", "make", "build", "generate", "write", "design"],
            "analyze": ["analyze", "examine", "investigate", "evaluate", "assess"],
            "solve": ["solve", "fix", "debug", "troubleshoot", "resolve"],
            "compare": ["compare", "versus", "vs", "difference between", "contrast"],
            "guide": ["how to", "guide", "tutorial", "teach", "show me"],
            "optimize": ["optimize", "improve", "enhance", "refactor", "speed up"],
            "research": ["research", "find", "discover", "explore", "study"]
        }
        
        for action, keywords in action_keywords.items():
            if any(keyword in prompt_lower for keyword in keywords):
                intent["primary_action"] = action
                break
        
        domain_indicators = {
            "technical": ["code", "programming", "software", "algorithm", "database", "api"],
            "creative": ["story", "poem", "design", "art", "creative", "imagine"],
            "business": ["business", "strategy", "marketing", "finance", "management"],
            "academic": ["research", "paper", "thesis", "study", "academic", "scholarly"],
            "personal": ["help me", "my", "I need", "personal", "advice"],
            "scientific": ["science", "experiment", "hypothesis", "data", "analysis"]
        }
        
        for domain, indicators in domain_indicators.items():
            if any(indicator in prompt_lower for indicator in indicators):
                intent["domain"] = domain
                break
        
        if len(prompt.split()) > 100 or "comprehensive" in prompt_lower or "detailed" in prompt_lower:
            intent["complexity"] = "high"
            intent["expected_length"] = "extensive"
        elif len(prompt.split()) < 20 and "?" in prompt:
            intent["complexity"] = "low"
            intent["expected_length"] = "concise"
        
        if any(word in prompt_lower for word in ["urgent", "asap", "quickly", "now"]):
            intent["urgency"] = "high"
        
        if any(word in prompt_lower for word in ["please", "kindly", "would you"]):
            intent["formality"] = "polite"
        elif any(word in prompt_lower for word in ["formal", "professional", "academic"]):
            intent["formality"] = "formal"
        elif any(word in prompt_lower for word in ["hey", "yo", "sup"]):
            intent["formality"] = "casual"
        
        intent["requires_examples"] = "example" in prompt_lower or "demonstrate" in prompt_lower
        intent["requires_code"] = any(word in prompt_lower for word in ["code", "implement", "function", "class"])
        intent["requires_visuals"] = any(word in prompt_lower for word in ["diagram", "chart", "visualize", "draw"])
        intent["requires_sources"] = any(word in prompt_lower for word in ["source", "citation", "reference", "evidence"])
        
        return intent
    
    def generate_meta_prompt(self, original_prompt: str, intent: Dict[str, Any]) -> str:
        meta_prompt = f"User Intent Analysis:\n"
        meta_prompt += f"- Primary Action: {intent['primary_action']}\n"
        meta_prompt += f"- Domain: {intent['domain']}\n"
        meta_prompt += f"- Complexity: {intent['complexity']}\n"
        meta_prompt += f"- Expected Length: {intent['expected_length']}\n"
        
        if intent["requires_examples"]:
            meta_prompt += "- Include practical examples\n"
        if intent["requires_code"]:
            meta_prompt += "- Provide complete, working code\n"
        if intent["requires_visuals"]:
            meta_prompt += "- Include visual representations\n"
        if intent["requires_sources"]:
            meta_prompt += "- Cite credible sources\n"
        
        meta_prompt += f"\nOptimized Approach: Tailor response to match user's intent and expectations."
        
        return meta_prompt

class UltraFileProcessor:
    def __init__(self):
        self.supported_formats = self._initialize_supported_formats()
        self.processors = self._initialize_processors()
        self.analyzers = self._initialize_analyzers()
        self.max_file_size = 2 * 1024 * 1024 * 1024
        self.temp_dir = "temp"
        self.cache_dir = "cache"
        self.thumbnail_size = (256, 256)
        self.preview_size = (1024, 1024)
        
        os.makedirs(self.temp_dir, exist_ok=True)
        os.makedirs(self.cache_dir, exist_ok=True)
    
    def _initialize_supported_formats(self) -> Dict[FileType, List[str]]:
        return {
            FileType.IMAGE: ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'svg', 'ico', 'tiff', 'tif', 'psd', 'raw', 'heif', 'heic', 'avif', 'jxl'],
            FileType.DOCUMENT: ['pdf', 'docx', 'doc', 'txt', 'rtf', 'odt', 'tex', 'md', 'rst', 'org', 'textile', 'asciidoc'],
            FileType.SPREADSHEET: ['xlsx', 'xls', 'csv', 'ods', 'tsv', 'xlsm', 'xlsb', 'xltx', 'xltm'],
            FileType.PRESENTATION: ['pptx', 'ppt', 'odp', 'key', 'pps', 'ppsx', 'pptm', 'ppsm'],
            FileType.CODE: ['py', 'js', 'java', 'cpp', 'c', 'cs', 'php', 'rb', 'go', 'rs', 'kt', 'swift', 'ts', 'jsx', 'vue', 'dart', 'scala', 'r', 'julia', 'lua', 'perl', 'sh', 'ps1', 'bat', 'cmd', 'asm', 'pas', 'f90', 'f95', 'f03', 'cob', 'lisp', 'clj', 'elm', 'ex', 'erl', 'hrl', 'nim', 'cr', 'ml', 'fs', 'vb', 'bas', 'prolog', 'haskell', 'scheme', 'racket', 'forth', 'tcl', 'awk', 'sed'],
            FileType.DATA: ['json', 'xml', 'yaml', 'yml', 'toml', 'ini', 'conf', 'env', 'properties', 'cfg', 'config', 'settings'],
            FileType.DATABASE: ['sql', 'db', 'sqlite', 'sqlite3', 'mdb', 'accdb', 'dbf', 'dump'],
            FileType.ARCHIVE: ['zip', 'tar', 'gz', 'bz2', 'xz', 'rar', '7z', 'cab', 'iso', 'dmg', 'pkg', 'deb', 'rpm', 'apk', 'jar', 'war', 'ear'],
            FileType.AUDIO: ['mp3', 'wav', 'flac', 'aac', 'ogg', 'wma', 'm4a', 'opus', 'webm', 'amr', 'ac3', 'aiff', 'ape', 'dts', 'mka', 'mp2', 'mpa', 'ra', 'tta', 'voc', 'wv'],
            FileType.VIDEO: ['mp4', 'avi', 'mkv', 'mov', 'wmv', 'flv', 'webm', 'ogv', 'm4v', 'mpg', 'mpeg', '3gp', 'f4v', 'swf', 'vob', 'ts', 'm2ts', 'mts', 'divx', 'xvid', 'rmvb', 'rm', 'asf'],
            FileType.MODEL: ['h5', 'hdf5', 'pb', 'pth', 'pt', 'onnx', 'tflite', 'caffemodel', 'mlmodel', 'pmml', 'pkl', 'joblib', 'sav', 'model', 'weights'],
            FileType.DIAGRAM: ['drawio', 'vsdx', 'vsd', 'gliffy', 'graffle', 'dia', 'xmind', 'mm', 'mindnode', 'graphml', 'dot', 'gv', 'plantuml', 'puml'],
            FileType.MARKUP: ['html', 'htm', 'xhtml', 'css', 'scss', 'sass', 'less', 'styl'],
            FileType.CONFIG: ['dockerfile', 'vagrantfile', 'makefile', 'rakefile', 'gulpfile', 'gruntfile', 'webpack', 'babel', 'eslint', 'prettier', 'gitignore', 'gitattributes', 'editorconfig', 'htaccess', 'nginx', 'apache'],
            FileType.LOG: ['log', 'out', 'err', 'trace', 'debug', 'info', 'warn', 'error', 'fatal'],
            FileType.BINARY: ['exe', 'dll', 'so', 'dylib', 'a', 'lib', 'o', 'obj', 'bin', 'dat', 'pak'],
            FileType.SCIENTIFIC: ['fits', 'fts', 'nc', 'nc4', 'netcdf', 'hdf', 'h4', 'h5', 'mat', 'sav', 'npz', 'npy', 'feather', 'parquet', 'arrow'],
            FileType.MEDICAL: ['dcm', 'dicom', 'nii', 'nii.gz', 'mha', 'mhd', 'nrrd', 'nhdr', 'img', 'hdr'],
            FileType.GEOSPATIAL: ['shp', 'kml', 'kmz', 'gpx', 'geojson', 'topojson', 'osm', 'pbf', 'tif', 'tiff', 'img', 'sid', 'ecw'],
            FileType.MOLECULAR: ['pdb', 'cif', 'mol', 'mol2', 'sdf', 'xyz', 'gro', 'trr', 'xtc', 'dcd', 'prmtop', 'inpcrd']
        }
    
    def _initialize_processors(self) -> Dict[FileType, callable]:
        return {
            FileType.IMAGE: self._process_image,
            FileType.DOCUMENT: self._process_document,
            FileType.SPREADSHEET: self._process_spreadsheet,
            FileType.PRESENTATION: self._process_presentation,
            FileType.CODE: self._process_code,
            FileType.DATA: self._process_data,
            FileType.DATABASE: self._process_database,
            FileType.ARCHIVE: self._process_archive,
            FileType.AUDIO: self._process_audio,
            FileType.VIDEO: self._process_video,
            FileType.MODEL: self._process_model,
            FileType.DIAGRAM: self._process_diagram,
            FileType.MARKUP: self._process_markup,
            FileType.CONFIG: self._process_config,
            FileType.LOG: self._process_log,
            FileType.BINARY: self._process_binary,
            FileType.SCIENTIFIC: self._process_scientific,
            FileType.MEDICAL: self._process_medical,
            FileType.GEOSPATIAL: self._process_geospatial,
            FileType.MOLECULAR: self._process_molecular
        }
    
    def _initialize_analyzers(self) -> Dict[FileType, callable]:
        return {
            FileType.IMAGE: self._analyze_image,
            FileType.DOCUMENT: self._analyze_document,
            FileType.SPREADSHEET: self._analyze_spreadsheet,
            FileType.CODE: self._analyze_code,
            FileType.AUDIO: self._analyze_audio,
            FileType.VIDEO: self._analyze_video,
            FileType.SCIENTIFIC: self._analyze_scientific,
            FileType.MEDICAL: self._analyze_medical,
            FileType.GEOSPATIAL: self._analyze_geospatial,
            FileType.MOLECULAR: self._analyze_molecular
        }
    
    def get_file_type(self, filename: str, mime_type: Optional[str] = None) -> Optional[FileType]:
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        
        for file_type, extensions in self.supported_formats.items():
            if ext in extensions:
                return file_type
        
        if mime_type:
            mime_map = {
                'image/': FileType.IMAGE,
                'application/pdf': FileType.DOCUMENT,
                'application/msword': FileType.DOCUMENT,
                'application/vnd.ms-excel': FileType.SPREADSHEET,
                'application/vnd.ms-powerpoint': FileType.PRESENTATION,
                'text/': FileType.CODE,
                'audio/': FileType.AUDIO,
                'video/': FileType.VIDEO,
                'application/x-': FileType.BINARY
            }
            
            for mime_prefix, file_type in mime_map.items():
                if mime_type.startswith(mime_prefix):
                    return file_type
        
        if MAGIC_AVAILABLE:
            try:
                file_magic = magic.from_file(filename, mime=True)
                for mime_prefix, file_type in mime_map.items():
                    if file_magic.startswith(mime_prefix):
                        return file_type
            except:
                pass
        
        return None
    
    async def process_file(self, file, attachment_id: str) -> Attachment:
        start_time = time.time()
        
        try:
            filename = secure_filename(file.filename)
            file_type = self.get_file_type(filename)
            
            file.seek(0, 2)
            file_size = file.tell()
            file.seek(0)
            
            if file_size > self.max_file_size:
                raise ValueError(f"File too large: {file_size / 1024 / 1024 / 1024:.1f}GB (max {self.max_file_size / 1024 / 1024 / 1024}GB)")
            
            temp_path = os.path.join(self.temp_dir, f"{attachment_id}_{filename}")
            file.save(temp_path)
            
            file_hash = self._calculate_file_hash(temp_path)
            
            mime_type = mimetypes.guess_type(filename)[0] or 'application/octet-stream'
            
            attachment = Attachment(
                id=attachment_id,
                filename=filename,
                size=file_size,
                mime_type=mime_type,
                file_type=file_type or FileType.BINARY,
                path=temp_path,
                checksum=file_hash,
                metadata={
                    "upload_time": datetime.now(timezone.utc).isoformat(),
                    "original_filename": file.filename
                }
            )
            
            if file_type and file_type in self.processors:
                processor = self.processors[file_type]
                attachment = await processor(attachment)
            else:
                attachment.extracted_text = f"[Binary file: {filename} ({mime_type})]"
                attachment.processed = True
            
            if file_type and file_type in self.analyzers:
                analyzer = self.analyzers[file_type]
                attachment = await analyzer(attachment)
            
            attachment.processing_time = time.time() - start_time
            
            return attachment
            
        except Exception as e:
            print(f"Error processing file {file.filename}: {e}")
            traceback.print_exc()
            
            return Attachment(
                id=attachment_id,
                filename=file.filename,
                size=0,
                mime_type='application/octet-stream',
                file_type=FileType.BINARY,
                path="",
                processed=False,
                error=str(e),
                processing_time=time.time() - start_time
            )
    
    def _calculate_file_hash(self, file_path: str) -> str:
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    
    async def _process_image(self, attachment: Attachment) -> Attachment:
        try:
            img = Image.open(attachment.path)
            
            attachment.metadata.update({
                "format": img.format,
                "mode": img.mode,
                "size": img.size,
                "info": img.info
            })
            
            if img.mode not in ('RGB', 'RGBA'):
                if img.mode == 'RGBA':
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[3])
                    img = background
                else:
                    img = img.convert('RGB')
            
            thumbnail = img.copy()
            thumbnail.thumbnail(self.thumbnail_size, Image.Resampling.LANCZOS)
            thumbnail_path = os.path.join(self.cache_dir, f"{attachment.id}_thumb.jpg")
            thumbnail.save(thumbnail_path, 'JPEG', quality=85)
            
            with open(thumbnail_path, 'rb') as f:
                attachment.thumbnail = base64.b64encode(f.read()).decode('utf-8')
            
            preview = img.copy()
            preview.thumbnail(self.preview_size, Image.Resampling.LANCZOS)
            preview_path = os.path.join(self.cache_dir, f"{attachment.id}_preview.jpg")
            preview.save(preview_path, 'JPEG', quality=90)
            
            with open(preview_path, 'rb') as f:
                attachment.preview = base64.b64encode(f.read()).decode('utf-8')
            
            if 'text' in img.info:
                attachment.extracted_text = img.info['text']
            else:
                try:
                    text = pytesseract.image_to_string(img)
                    if text.strip():
                        attachment.extracted_text = text
                except:
                    pass
            
            content = f"[Image: {attachment.filename}]\n"
            content += f"Format: {img.format}\n"
            content += f"Size: {img.size[0]}x{img.size[1]}\n"
            content += f"Mode: {img.mode}\n"
            
            if attachment.extracted_text:
                content += f"\nExtracted Text:\n{attachment.extracted_text[:1000]}"
            
            attachment.extracted_text = content
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Image processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process image: {attachment.filename}]"
        
        return attachment
    
    async def _process_document(self, attachment: Attachment) -> Attachment:
        ext = attachment.filename.lower().split('.')[-1]
        
        try:
            if ext == 'pdf':
                text = ""
                metadata = {"pages": 0, "info": {}}
                
                with pdfplumber.open(attachment.path) as pdf:
                    metadata["pages"] = len(pdf.pages)
                    if pdf.metadata:
                        metadata["info"] = {
                            k: v for k, v in pdf.metadata.items() 
                            if isinstance(v, (str, int, float))
                        }
                    
                    for i, page in enumerate(pdf.pages[:100]):
                        page_text = page.extract_text() or ""
                        if page_text:
                            text += f"\n\n--- Page {i+1} ---\n{page_text}"
                        
                        tables = page.extract_tables()
                        if tables:
                            text += f"\n\nTables on page {i+1}:\n"
                            for j, table in enumerate(tables):
                                text += f"\nTable {j+1}:\n"
                                text += "\n".join([" | ".join(str(cell) for cell in row) for row in table])
                
                attachment.metadata.update(metadata)
                attachment.extracted_text = text[:500000]
                
            elif ext in ['docx', 'doc']:
                doc = docx.Document(attachment.path)
                paragraphs = []
                tables_data = []
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text)
                
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    tables_data.append(table_data)
                
                text = "\n\n".join(paragraphs)
                
                if tables_data:
                    text += "\n\nTables:\n"
                    for i, table in enumerate(tables_data):
                        text += f"\nTable {i+1}:\n"
                        text += "\n".join([" | ".join(row) for row in table])
                
                attachment.metadata.update({
                    "paragraphs": len(paragraphs),
                    "tables": len(tables_data),
                    "sections": len(doc.sections)
                })
                
                attachment.extracted_text = text[:500000]
                
            elif ext == 'md':
                with open(attachment.path, 'r', encoding='utf-8') as f:
                    text = f.read()
                
                html = markdown.markdown(text, extensions=[
                    'tables', 'fenced_code', 'footnotes', 'attr_list',
                    'def_list', 'abbr', 'md_in_html', 'toc'
                ])
                
                attachment.metadata.update({
                    "length": len(text),
                    "lines": len(text.splitlines()),
                    "words": len(text.split())
                })
                
                attachment.extracted_text = text[:500000]
                attachment.preview = html[:100000]
                
            else:
                with open(attachment.path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                
                attachment.metadata.update({
                    "encoding": "utf-8",
                    "lines": len(text.splitlines()),
                    "words": len(text.split()),
                    "characters": len(text)
                })
                
                attachment.extracted_text = text[:500000]
            
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Document processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process document: {attachment.filename}]"
        
        return attachment
    
    async def _process_spreadsheet(self, attachment: Attachment) -> Attachment:
        ext = attachment.filename.lower().split('.')[-1]
        
        try:
            if ext == 'csv':
                encoding = self._detect_encoding(attachment.path)
                df = pd.read_csv(attachment.path, encoding=encoding, on_bad_lines='skip')
            else:
                df = pd.read_excel(attachment.path, engine='openpyxl' if ext.startswith('xls') else None)
            
            summary = {
                "rows": len(df),
                "columns": len(df.columns),
                "columns_list": df.columns.tolist(),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
                "null_counts": df.isnull().sum().to_dict(),
                "memory_usage": f"{df.memory_usage(deep=True).sum() / 1024 / 1024:.2f} MB"
            }
            
            stats = {}
            numeric_columns = df.select_dtypes(include=[np.number]).columns
            for col in numeric_columns[:20]:
                if not df[col].isna().all():
                    stats[col] = {
                        "mean": float(df[col].mean()),
                        "std": float(df[col].std()),
                        "min": float(df[col].min()),
                        "max": float(df[col].max()),
                        "median": float(df[col].median()),
                        "q1": float(df[col].quantile(0.25)),
                        "q3": float(df[col].quantile(0.75))
                    }
            
            correlations = {}
            if len(numeric_columns) > 1 and len(numeric_columns) <= 50:
                corr_matrix = df[numeric_columns].corr()
                for i in range(len(numeric_columns)):
                    for j in range(i+1, len(numeric_columns)):
                        corr_value = corr_matrix.iloc[i, j]
                        if abs(corr_value) > 0.5:
                            correlations[f"{numeric_columns[i]} - {numeric_columns[j]}"] = float(corr_value)
            
            sample_rows = min(1000, len(df))
            sample = df.head(sample_rows).to_dict('records')
            
            text_preview = df.head(100).to_string(max_rows=100, max_cols=20)
            
            attachment.metadata.update({
                "summary": summary,
                "statistics": stats,
                "correlations": correlations,
                "sample_data": sample[:100]
            })
            
            content = f"[Spreadsheet: {attachment.filename}]\n"
            content += f"Rows: {summary['rows']:,}, Columns: {summary['columns']}\n"
            content += f"Memory Usage: {summary['memory_usage']}\n\n"
            content += f"Columns:\n{chr(10).join(summary['columns_list'][:50])}\n\n"
            
            if stats:
                content += "Statistical Summary:\n"
                for col, col_stats in list(stats.items())[:10]:
                    content += f"\n{col}:\n"
                    content += f"  Mean: {col_stats['mean']:.2f}\n"
                    content += f"  Std: {col_stats['std']:.2f}\n"
                    content += f"  Range: [{col_stats['min']:.2f}, {col_stats['max']:.2f}]\n"
            
            if correlations:
                content += "\nStrong Correlations:\n"
                for pair, corr in list(correlations.items())[:10]:
                    content += f"  {pair}: {corr:.3f}\n"
            
            content += f"\n\nData Preview:\n{text_preview}"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Spreadsheet processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process spreadsheet: {attachment.filename}]"
        
        return attachment
    
    async def _process_presentation(self, attachment: Attachment) -> Attachment:
        try:
            prs = pptx.Presentation(attachment.path)
            
            slides_content = []
            total_shapes = 0
            total_images = 0
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                shapes_count = 0
                images_count = 0
                
                for shape in slide.shapes:
                    shapes_count += 1
                    
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    if shape.shape_type == 13:
                        images_count += 1
                
                total_shapes += shapes_count
                total_images += images_count
                
                slide_content = {
                    "slide_number": i + 1,
                    "layout": slide.slide_layout.name if slide.slide_layout else "Custom",
                    "text": "\n".join(slide_text),
                    "shapes": shapes_count,
                    "images": images_count
                }
                
                slides_content.append(slide_content)
            
            attachment.metadata.update({
                "slides": len(prs.slides),
                "total_shapes": total_shapes,
                "total_images": total_images,
                "presentation_size": (prs.slide_width, prs.slide_height)
            })
            
            content = f"[Presentation: {attachment.filename}]\n"
            content += f"Slides: {len(prs.slides)}\n"
            content += f"Total Shapes: {total_shapes}\n"
            content += f"Total Images: {total_images}\n\n"
            
            for slide_info in slides_content[:50]:
                content += f"\n--- Slide {slide_info['slide_number']} ({slide_info['layout']}) ---\n"
                content += f"Shapes: {slide_info['shapes']}, Images: {slide_info['images']}\n"
                if slide_info['text']:
                    content += f"Content:\n{slide_info['text']}\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Presentation processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process presentation: {attachment.filename}]"
        
        return attachment
    
    async def _process_code(self, attachment: Attachment) -> Attachment:
        try:
            encoding = self._detect_encoding(attachment.path)
            with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                code = f.read()
            
            ext = attachment.filename.lower().split('.')[-1]
            
            lines = code.split('\n')
            line_count = len(lines)
            
            try:
                lexer = get_lexer_by_name(ext)
                language = lexer.name
            except:
                try:
                    lexer = guess_lexer(code)
                    language = lexer.name
                except:
                    language = ext.upper()
            
            imports = []
            functions = []
            classes = []
            comments = []
            docstrings = []
            
            language_patterns = {
                'python': {
                    'import': r'^(?:import|from)\s+[\w\.]+',
                    'function': r'^def\s+(\w+)\s*\(',
                    'class': r'^class\s+(\w+)',
                    'comment': r'^\s*#(.*)$',
                    'docstring': r'"""(.*?)"""'
                },
                'javascript': {
                    'import': r'^(?:import|const|let|var)\s+.*from\s+[\'"]',
                    'function': r'(?:function\s+(\w+)|const\s+(\w+)\s*=\s*(?:\(.*?\)\s*=>|function))',
                    'class': r'^class\s+(\w+)',
                    'comment': r'^\s*//(.*)$|/\*(.*?)\*/',
                    'docstring': r'/\*\*(.*?)\*/'
                },
                'java': {
                    'import': r'^import\s+[\w\.]+;',
                    'function': r'(?:public|private|protected)?\s*(?:static)?\s*\w+\s+(\w+)\s*\(',
                    'class': r'^(?:public\s+)?class\s+(\w+)',
                    'comment': r'^\s*//(.*)$|/\*(.*?)\*/',
                    'docstring': r'/\*\*(.*?)\*/'
                }
            }
            
            patterns = language_patterns.get(ext, language_patterns.get('python'))
            
            for line in lines:
                stripped = line.strip()
                
                if patterns['import'] and re.match(patterns['import'], stripped):
                    imports.append(stripped)
                
                if patterns['comment']:
                    comment_match = re.match(patterns['comment'], stripped)
                    if comment_match:
                        comments.append(comment_match.group(1).strip())
            
            if patterns['function']:
                functions = re.findall(patterns['function'], code, re.MULTILINE)
                functions = [f for f in functions if f]
            
            if patterns['class']:
                classes = re.findall(patterns['class'], code, re.MULTILINE)
            
            if patterns['docstring']:
                docstrings = re.findall(patterns['docstring'], code, re.DOTALL)
            
            complexity_metrics = self._calculate_code_complexity(code, language)
            
            formatter = HtmlFormatter(style='monokai', linenos=True)
            highlighted = pygments.highlight(code[:10000], lexer, formatter)
            
            attachment.metadata.update({
                "language": language,
                "lines": line_count,
                "size": len(code),
                "imports": imports[:50],
                "functions": functions[:50],
                "classes": classes[:50],
                "comments": len(comments),
                "docstrings": len(docstrings),
                "complexity": complexity_metrics
            })
            
            content = f"[Code File: {attachment.filename}]\n"
            content += f"Language: {language}\n"
            content += f"Lines: {line_count:,}\n"
            content += f"Size: {len(code):,} bytes\n\n"
            
            content += f"Structure:\n"
            content += f"  Imports: {len(imports)}\n"
            content += f"  Functions: {len(functions)}\n"
            content += f"  Classes: {len(classes)}\n"
            content += f"  Comments: {len(comments)}\n"
            content += f"  Docstrings: {len(docstrings)}\n\n"
            
            if complexity_metrics:
                content += f"Complexity Metrics:\n"
                for metric, value in complexity_metrics.items():
                    content += f"  {metric}: {value}\n"
                content += "\n"
            
            if imports:
                content += f"Imports ({len(imports[:20])}):\n"
                content += "\n".join(f"  {imp}" for imp in imports[:20])
                content += "\n\n"
            
            if functions:
                content += f"Functions ({len(functions[:20])}):\n"
                content += "\n".join(f"  {func}()" for func in functions[:20] if func)
                content += "\n\n"
            
            if classes:
                content += f"Classes ({len(classes[:20])}):\n"
                content += "\n".join(f"  {cls}" for cls in classes[:20])
                content += "\n\n"
            
            content += f"Code Preview:\n"
            content += "\n".join(f"{i+1:5d} | {line}" for i, line in enumerate(lines[:200]))
            
            attachment.extracted_text = content[:500000]
            attachment.preview = highlighted
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Code processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process code: {attachment.filename}]"
        
        return attachment
    
    async def _process_data(self, attachment: Attachment) -> Attachment:
        try:
            encoding = self._detect_encoding(attachment.path)
            with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                content_raw = f.read()
            
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext == 'json':
                try:
                    data = json.loads(content_raw)
                    structure = self._analyze_json_structure(data)
                    pretty = json.dumps(data, indent=2, ensure_ascii=False)[:100000]
                    
                    attachment.metadata.update({
                        "valid": True,
                        "structure": structure,
                        "size": len(content_raw)
                    })
                    
                    content = f"[JSON File: {attachment.filename}]\n"
                    content += f"Valid: ✓\n"
                    content += f"Size: {len(content_raw):,} bytes\n"
                    content += f"Structure: {structure['type']}\n\n"
                    content += f"Content Preview:\n{pretty}"
                    
                except json.JSONDecodeError as e:
                    attachment.metadata.update({
                        "valid": False,
                        "error": str(e),
                        "size": len(content_raw)
                    })
                    
                    content = f"[JSON File: {attachment.filename}]\n"
                    content += f"Valid: ✗\n"
                    content += f"Error: {str(e)}\n"
                    content += f"Size: {len(content_raw):,} bytes\n\n"
                    content += f"Raw Content:\n{content_raw[:10000]}"
            
            elif ext in ['yaml', 'yml']:
                try:
                    import yaml
                    data = yaml.safe_load(content_raw)
                    structure = self._analyze_json_structure(data)
                    
                    attachment.metadata.update({
                        "valid": True,
                        "structure": structure,
                        "size": len(content_raw)
                    })
                    
                    content = f"[YAML File: {attachment.filename}]\n"
                    content += f"Valid: ✓\n"
                    content += f"Size: {len(content_raw):,} bytes\n\n"
                    content += f"Content:\n{content_raw[:100000]}"
                    
                except Exception as e:
                    attachment.metadata.update({
                        "valid": False,
                        "error": str(e),
                        "size": len(content_raw)
                    })
                    
                    content = f"[YAML File: {attachment.filename}]\n"
                    content += f"Valid: ✗\n"
                    content += f"Error: {str(e)}\n\n"
                    content += f"Raw Content:\n{content_raw[:100000]}"
            
            elif ext == 'xml':
                try:
                    soup = BeautifulSoup(content_raw, 'xml')
                    pretty = soup.prettify()[:100000]
                    
                    root_tag = soup.find()
                    tag_count = len(soup.find_all())
                    
                    attachment.metadata.update({
                        "valid": True,
                        "root_tag": root_tag.name if root_tag else None,
                        "total_tags": tag_count,
                        "size": len(content_raw)
                    })
                    
                    content = f"[XML File: {attachment.filename}]\n"
                    content += f"Valid: ✓\n"
                    content += f"Root Tag: {root_tag.name if root_tag else 'None'}\n"
                    content += f"Total Tags: {tag_count:,}\n"
                    content += f"Size: {len(content_raw):,} bytes\n\n"
                    content += f"Content:\n{pretty}"
                    
                except Exception as e:
                    attachment.metadata.update({
                        "valid": False,
                        "error": str(e),
                        "size": len(content_raw)
                    })
                    
                    content = f"[XML File: {attachment.filename}]\n"
                    content += f"Valid: ✗\n"
                    content += f"Error: {str(e)}\n\n"
                    content += f"Raw Content:\n{content_raw[:100000]}"
            
            else:
                lines = content_raw.splitlines()
                
                attachment.metadata.update({
                    "format": ext.upper(),
                    "size": len(content_raw),
                    "lines": len(lines)
                })
                
                content = f"[Data File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {len(content_raw):,} bytes\n"
                content += f"Lines: {len(lines):,}\n\n"
                content += f"Content:\n{content_raw[:100000]}"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Data processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process data: {attachment.filename}]"
        
        return attachment
    
    async def _process_database(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['sql', 'dump']:
                encoding = self._detect_encoding(attachment.path)
                with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                    content = f.read()
                
                queries = []
                current_query = []
                
                for line in content.splitlines():
                    stripped = line.strip()
                    if stripped and not stripped.startswith('--'):
                        current_query.append(line)
                        if stripped.endswith(';'):
                            queries.append('\n'.join(current_query))
                            current_query = []
                
                if current_query:
                    queries.append('\n'.join(current_query))
                
                tables = re.findall(r'CREATE TABLE\s+(?:IF NOT EXISTS\s+)?`?(\w+)`?', content, re.IGNORECASE)
                inserts = len(re.findall(r'INSERT INTO', content, re.IGNORECASE))
                
                attachment.metadata.update({
                    "queries": len(queries),
                    "tables": tables,
                    "inserts": inserts,
                    "size": len(content)
                })
                
                content_display = f"[SQL File: {attachment.filename}]\n"
                content_display += f"Total Queries: {len(queries)}\n"
                content_display += f"Tables: {len(tables)}\n"
                content_display += f"Insert Statements: {inserts}\n"
                content_display += f"File Size: {len(content):,} bytes\n\n"
                
                if tables:
                    content_display += f"Tables found:\n"
                    content_display += "\n".join(f"  - {table}" for table in tables[:20])
                    content_display += "\n\n"
                
                content_display += f"Queries Preview:\n"
                for i, query in enumerate(queries[:20]):
                    content_display += f"\n--- Query {i+1} ---\n{query[:500]}\n"
                
                attachment.extracted_text = content_display[:500000]
                
            elif ext in ['sqlite', 'sqlite3', 'db']:
                import sqlite3
                
                conn = sqlite3.connect(attachment.path)
                cursor = conn.cursor()
                
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = [row[0] for row in cursor.fetchall()]
                
                table_info = {}
                for table in tables[:50]:
                    cursor.execute(f"PRAGMA table_info({table});")
                    columns = cursor.fetchall()
                    cursor.execute(f"SELECT COUNT(*) FROM {table};")
                    row_count = cursor.fetchone()[0]
                    
                    table_info[table] = {
                        "columns": [(col[1], col[2]) for col in columns],
                        "row_count": row_count
                    }
                
                conn.close()
                
                attachment.metadata.update({
                    "tables": tables,
                    "table_count": len(tables),
                    "table_info": table_info
                })
                
                content = f"[SQLite Database: {attachment.filename}]\n"
                content += f"Tables: {len(tables)}\n\n"
                
                for table, info in list(table_info.items())[:20]:
                    content += f"\nTable: {table}\n"
                    content += f"  Rows: {info['row_count']:,}\n"
                    content += f"  Columns:\n"
                    for col_name, col_type in info['columns'][:20]:
                        content += f"    - {col_name} ({col_type})\n"
                
                attachment.extracted_text = content[:500000]
            
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Database processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process database: {attachment.filename}]"
        
        return attachment
    
    async def _process_archive(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            file_list = []
            total_size = 0
            file_types = defaultdict(int)
            
            if ext == 'zip':
                with zipfile.ZipFile(attachment.path, 'r') as zf:
                    for info in zf.infolist():
                        file_list.append({
                            "name": info.filename,
                            "size": info.file_size,
                            "compressed": info.compress_size,
                            "date": info.date_time
                        })
                        total_size += info.file_size
                        
                        if '.' in info.filename:
                            ext = info.filename.split('.')[-1].lower()
                            file_types[ext] += 1
            
            elif ext in ['tar', 'gz', 'bz2', 'xz']:
                with tarfile.open(attachment.path, 'r:*') as tf:
                    for member in tf.getmembers():
                        if member.isfile():
                            file_list.append({
                                "name": member.name,
                                "size": member.size,
                                "date": member.mtime
                            })
                            total_size += member.size
                            
                            if '.' in member.name:
                                ext = member.name.split('.')[-1].lower()
                                file_types[ext] += 1
            
            elif ext == 'rar':
                with rarfile.RarFile(attachment.path, 'r') as rf:
                    for info in rf.infolist():
                        file_list.append({
                            "name": info.filename,
                            "size": info.file_size,
                            "compressed": info.compress_size,
                            "date": info.date_time
                        })
                        total_size += info.file_size
                        
                        if '.' in info.filename:
                            ext = info.filename.split('.')[-1].lower()
                            file_types[ext] += 1
            
            elif ext == '7z':
                with py7zr.SevenZipFile(attachment.path, 'r') as sz:
                    for name, info in sz.list():
                        file_list.append({
                            "name": name,
                            "size": info.uncompressed,
                            "compressed": info.compressed
                        })
                        total_size += info.uncompressed
                        
                        if '.' in name:
                            ext = name.split('.')[-1].lower()
                            file_types[ext] += 1
            
            attachment.metadata.update({
                "files": len(file_list),
                "total_size": total_size,
                "file_types": dict(file_types),
                "compression_ratio": attachment.size / total_size if total_size > 0 else 1
            })
            
            content = f"[Archive: {attachment.filename}]\n"
            content += f"Type: {ext.upper()}\n"
            content += f"Files: {len(file_list):,}\n"
            content += f"Total Size: {total_size / 1024 / 1024:.2f} MB\n"
            content += f"Compressed Size: {attachment.size / 1024 / 1024:.2f} MB\n"
            content += f"Compression Ratio: {(1 - attachment.size / total_size) * 100:.1f}%\n\n"
            
            if file_types:
                content += "File Types:\n"
                sorted_types = sorted(file_types.items(), key=lambda x: x[1], reverse=True)
                for ft, count in sorted_types[:20]:
                    content += f"  .{ft}: {count:,} files\n"
                content += "\n"
            
            content += "Files:\n"
            sorted_files = sorted(file_list, key=lambda x: x['size'], reverse=True)
            for file_info in sorted_files[:100]:
                content += f"  {file_info['name']} ({file_info['size'] / 1024:.1f} KB)\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Archive processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process archive: {attachment.filename}]"
        
        return attachment
    
    async def _process_audio(self, attachment: Attachment) -> Attachment:
        try:
            audio_data, sample_rate = librosa.load(attachment.path, sr=None)
            duration = len(audio_data) / sample_rate
            
            tempo, beats = librosa.beat.beat_track(y=audio_data, sr=sample_rate)
            
            spectral_centroids = librosa.feature.spectral_centroid(y=audio_data, sr=sample_rate)
            spectral_rolloff = librosa.feature.spectral_rolloff(y=audio_data, sr=sample_rate)
            zero_crossing_rate = librosa.feature.zero_crossing_rate(audio_data)
            mfccs = librosa.feature.mfcc(y=audio_data, sr=sample_rate, n_mfcc=13)
            
            attachment.metadata.update({
                "duration": duration,
                "sample_rate": sample_rate,
                "channels": 1,
                "tempo": float(tempo),
                "beats": len(beats),
                "spectral_centroid_mean": float(np.mean(spectral_centroids)),
                "spectral_rolloff_mean": float(np.mean(spectral_rolloff)),
                "zero_crossing_rate_mean": float(np.mean(zero_crossing_rate)),
                "mfcc_means": [float(np.mean(mfcc)) for mfcc in mfccs]
            })
            
            try:
                r = sr.Recognizer()
                with sr.AudioFile(attachment.path) as source:
                    audio = r.record(source, duration=min(60, duration))
                    text = r.recognize_google(audio)
                    attachment.extracted_text = text
            except:
                pass
            
            content = f"[Audio File: {attachment.filename}]\n"
            content += f"Duration: {duration:.2f} seconds\n"
            content += f"Sample Rate: {sample_rate} Hz\n"
            content += f"Tempo: {tempo:.1f} BPM\n"
            content += f"Beats Detected: {len(beats)}\n\n"
            
            if attachment.extracted_text:
                content += f"Transcription:\n{attachment.extracted_text}\n"
            
            attachment.extracted_text = content
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Audio processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process audio: {attachment.filename}]"
        
        return attachment
    
    async def _process_video(self, attachment: Attachment) -> Attachment:
        try:
            import cv2
            
            cap = cv2.VideoCapture(attachment.path)
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            duration = frame_count / fps if fps > 0 else 0
            
            frames_to_sample = min(10, frame_count)
            frame_interval = frame_count // frames_to_sample if frames_to_sample > 0 else 1
            
            sampled_frames = []
            for i in range(frames_to_sample):
                cap.set(cv2.CAP_PROP_POS_FRAMES, i * frame_interval)
                ret, frame = cap.read()
                if ret:
                    thumbnail = cv2.resize(frame, (320, 240))
                    _, buffer = cv2.imencode('.jpg', thumbnail)
                    sampled_frames.append(base64.b64encode(buffer).decode('utf-8'))
            
            cap.release()
            
            attachment.metadata.update({
                "duration": duration,
                "fps": fps,
                "frame_count": frame_count,
                "resolution": f"{width}x{height}",
                "width": width,
                "height": height,
                "sample_frames": sampled_frames[:3]
            })
            
            if sampled_frames:
                attachment.thumbnail = sampled_frames[0]
            
            content = f"[Video File: {attachment.filename}]\n"
            content += f"Duration: {duration:.2f} seconds\n"
            content += f"Resolution: {width}x{height}\n"
            content += f"FPS: {fps:.2f}\n"
            content += f"Total Frames: {frame_count:,}\n"
            content += f"Bitrate: {attachment.size * 8 / duration / 1000:.0f} kbps\n"
            
            attachment.extracted_text = content
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Video processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process video: {attachment.filename}]"
        
        return attachment
    
    async def _process_model(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['h5', 'hdf5']:
                import h5py
                
                with h5py.File(attachment.path, 'r') as f:
                    structure = self._explore_h5_structure(f)
                    
                    attachment.metadata.update({
                        "format": "HDF5",
                        "structure": structure,
                        "total_datasets": self._count_h5_datasets(f),
                        "total_groups": self._count_h5_groups(f)
                    })
                    
                    content = f"[Model File: {attachment.filename}]\n"
                    content += f"Format: HDF5\n"
                    content += f"Total Datasets: {attachment.metadata['total_datasets']}\n"
                    content += f"Total Groups: {attachment.metadata['total_groups']}\n\n"
                    content += f"Structure:\n{self._format_h5_structure(structure)}"
            
            elif ext in ['pth', 'pt']:
                import torch
                
                model_data = torch.load(attachment.path, map_location='cpu')
                
                if isinstance(model_data, dict):
                    keys = list(model_data.keys())
                    total_params = 0
                    
                    for key, value in model_data.items():
                        if isinstance(value, torch.Tensor):
                            total_params += value.numel()
                    
                    attachment.metadata.update({
                        "format": "PyTorch",
                        "keys": keys[:100],
                        "total_parameters": total_params
                    })
                    
                    content = f"[Model File: {attachment.filename}]\n"
                    content += f"Format: PyTorch\n"
                    content += f"Total Parameters: {total_params:,}\n"
                    content += f"Keys ({len(keys)}):\n"
                    content += "\n".join(f"  - {key}" for key in keys[:50])
                
            elif ext == 'onnx':
                import onnx
                
                model = onnx.load(attachment.path)
                
                attachment.metadata.update({
                    "format": "ONNX",
                    "producer": model.producer_name,
                    "version": model.producer_version,
                    "inputs": [inp.name for inp in model.graph.input],
                    "outputs": [out.name for out in model.graph.output]
                })
                
                content = f"[Model File: {attachment.filename}]\n"
                content += f"Format: ONNX\n"
                content += f"Producer: {model.producer_name} {model.producer_version}\n"
                content += f"Inputs: {', '.join(attachment.metadata['inputs'])}\n"
                content += f"Outputs: {', '.join(attachment.metadata['outputs'])}\n"
            
            else:
                content = f"[Model File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024 / 1024:.2f} MB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Model processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process model: {attachment.filename}]"
        
        return attachment
    
    async def _process_diagram(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['drawio', 'xml']:
                with open(attachment.path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                soup = BeautifulSoup(content, 'xml')
                
                cells = soup.find_all('mxCell')
                shapes = [cell for cell in cells if cell.get('vertex') == '1']
                edges = [cell for cell in cells if cell.get('edge') == '1']
                
                text_content = []
                for cell in cells:
                    value = cell.get('value', '')
                    if value and value.strip():
                        text_content.append(value)
                
                attachment.metadata.update({
                    "format": "draw.io",
                    "shapes": len(shapes),
                    "connections": len(edges),
                    "text_elements": len(text_content)
                })
                
                content = f"[Diagram: {attachment.filename}]\n"
                content += f"Format: draw.io\n"
                content += f"Shapes: {len(shapes)}\n"
                content += f"Connections: {len(edges)}\n"
                content += f"Text Elements: {len(text_content)}\n\n"
                
                if text_content:
                    content += "Text Content:\n"
                    content += "\n".join(f"  - {text}" for text in text_content[:50])
            
            elif ext == 'dot' or ext == 'gv':
                with open(attachment.path, 'r', encoding='utf-8') as f:
                    dot_content = f.read()
                
                nodes = re.findall(r'(\w+)\s*\[', dot_content)
                edges = re.findall(r'(\w+)\s*->\s*(\w+)', dot_content)
                
                attachment.metadata.update({
                    "format": "Graphviz",
                    "nodes": len(set(nodes)),
                    "edges": len(edges)
                })
                
                content = f"[Diagram: {attachment.filename}]\n"
                content += f"Format: Graphviz\n"
                content += f"Nodes: {len(set(nodes))}\n"
                content += f"Edges: {len(edges)}\n\n"
                content += f"Content Preview:\n{dot_content[:5000]}"
            
            else:
                content = f"[Diagram: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024:.1f} KB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Diagram processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process diagram: {attachment.filename}]"
        
        return attachment
    
    async def _process_markup(self, attachment: Attachment) -> Attachment:
        try:
            encoding = self._detect_encoding(attachment.path)
            with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read()
            
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['html', 'htm', 'xhtml']:
                soup = BeautifulSoup(content, 'html.parser')
                
                title = soup.find('title')
                title_text = title.text if title else None
                
                text = soup.get_text(separator='\n', strip=True)
                
                links = soup.find_all('a', href=True)
                images = soup.find_all('img', src=True)
                scripts = soup.find_all('script')
                styles = soup.find_all('style')
                
                attachment.metadata.update({
                    "title": title_text,
                    "links": len(links),
                    "images": len(images),
                    "scripts": len(scripts),
                    "styles": len(styles),
                    "text_length": len(text)
                })
                
                content_display = f"[HTML File: {attachment.filename}]\n"
                if title_text:
                    content_display += f"Title: {title_text}\n"
                content_display += f"Links: {len(links)}\n"
                content_display += f"Images: {len(images)}\n"
                content_display += f"Scripts: {len(scripts)}\n"
                content_display += f"Styles: {len(styles)}\n\n"
                content_display += f"Text Content:\n{text[:10000]}"
                
                attachment.extracted_text = content_display[:500000]
                attachment.preview = content[:100000]
            
            elif ext in ['css', 'scss', 'sass', 'less']:
                rules = re.findall(r'([.#]?[\w-]+)\s*{[^}]+}', content)
                variables = re.findall(r'(--[\w-]+|@[\w-]+|\$[\w-]+):\s*([^;]+);', content)
                imports = re.findall(r'@import\s+["\']([^"\']+)["\'];', content)
                
                attachment.metadata.update({
                    "type": ext.upper(),
                    "rules": len(rules),
                    "variables": len(variables),
                    "imports": len(imports)
                })
                
                content_display = f"[Stylesheet: {attachment.filename}]\n"
                content_display += f"Type: {ext.upper()}\n"
                content_display += f"Rules: {len(rules)}\n"
                content_display += f"Variables: {len(variables)}\n"
                content_display += f"Imports: {len(imports)}\n\n"
                
                if variables:
                    content_display += "Variables:\n"
                    for var, value in variables[:20]:
                        content_display += f"  {var}: {value}\n"
                    content_display += "\n"
                
                content_display += f"Content Preview:\n{content[:10000]}"
                
                attachment.extracted_text = content_display[:500000]
            
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Markup processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process markup: {attachment.filename}]"
        
        return attachment
    
    async def _process_config(self, attachment: Attachment) -> Attachment:
        try:
            encoding = self._detect_encoding(attachment.path)
            with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                content = f.read()
            
            lines = content.splitlines()
            
            config_type = attachment.filename.lower()
            
            if 'dockerfile' in config_type:
                instructions = defaultdict(list)
                for line in lines:
                    if line.strip() and not line.startswith('#'):
                        parts = line.split(None, 1)
                        if parts:
                            instruction = parts[0].upper()
                            value = parts[1] if len(parts) > 1 else ''
                            instructions[instruction].append(value)
                
                attachment.metadata.update({
                    "type": "Dockerfile",
                    "instructions": dict(instructions),
                    "base_image": instructions.get('FROM', [''])[0]
                })
                
                content_display = f"[Dockerfile: {attachment.filename}]\n"
                content_display += f"Base Image: {instructions.get('FROM', [''])[0]}\n"
                content_display += f"Instructions:\n"
                for inst, values in instructions.items():
                    content_display += f"  {inst}: {len(values)} occurrences\n"
            
            elif any(x in config_type for x in ['makefile', 'gnumakefile']):
                targets = re.findall(r'^([a-zA-Z_-]+):(?!=)', content, re.MULTILINE)
                variables = re.findall(r'^([A-Z_]+)\s*[?:]?=\s*(.*)', content, re.MULTILINE)
                
                attachment.metadata.update({
                    "type": "Makefile",
                    "targets": targets,
                    "variables": dict(variables)
                })
                
                content_display = f"[Makefile: {attachment.filename}]\n"
                content_display += f"Targets ({len(targets)}):\n"
                content_display += "\n".join(f"  - {target}" for target in targets[:20])
                content_display += f"\n\nVariables ({len(variables)}):\n"
                for var, val in variables[:20]:
                    content_display += f"  {var} = {val[:50]}\n"
            
            else:
                content_display = f"[Config File: {attachment.filename}]\n"
                content_display += f"Lines: {len(lines)}\n"
                content_display += f"Size: {len(content):,} bytes\n\n"
                content_display += f"Content:\n{content[:10000]}"
            
            attachment.extracted_text = content_display[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Config processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process config: {attachment.filename}]"
        
        return attachment
    
    async def _process_log(self, attachment: Attachment) -> Attachment:
        try:
            encoding = self._detect_encoding(attachment.path)
            
            log_stats = {
                "total_lines": 0,
                "error_lines": 0,
                "warning_lines": 0,
                "info_lines": 0,
                "debug_lines": 0,
                "timestamps": [],
                "unique_sources": set(),
                "error_messages": [],
                "warning_messages": []
            }
            
            error_patterns = re.compile(r'(ERROR|FATAL|CRITICAL|FAIL)', re.IGNORECASE)
            warning_patterns = re.compile(r'(WARN|WARNING)', re.IGNORECASE)
            info_patterns = re.compile(r'(INFO|INFORMATION)', re.IGNORECASE)
            debug_patterns = re.compile(r'(DEBUG|TRACE)', re.IGNORECASE)
            timestamp_pattern = re.compile(r'\d{4}-\d{2}-\d{2}[T\s]\d{2}:\d{2}:\d{2}')
            
            content_preview = []
            
            with open(attachment.path, 'r', encoding=encoding, errors='ignore') as f:
                for i, line in enumerate(f):
                    log_stats["total_lines"] += 1
                    
                    if i < 1000:
                        content_preview.append(line.rstrip())
                    
                    if error_patterns.search(line):
                        log_stats["error_lines"] += 1
                        if len(log_stats["error_messages"]) < 50:
                            log_stats["error_messages"].append(line.strip())
                    elif warning_patterns.search(line):
                        log_stats["warning_lines"] += 1
                        if len(log_stats["warning_messages"]) < 50:
                            log_stats["warning_messages"].append(line.strip())
                    elif info_patterns.search(line):
                        log_stats["info_lines"] += 1
                    elif debug_patterns.search(line):
                        log_stats["debug_lines"] += 1
                    
                    timestamp_match = timestamp_pattern.search(line)
                    if timestamp_match and len(log_stats["timestamps"]) < 100:
                        log_stats["timestamps"].append(timestamp_match.group())
            
            if log_stats["timestamps"]:
                log_stats["first_timestamp"] = log_stats["timestamps"][0]
                log_stats["last_timestamp"] = log_stats["timestamps"][-1]
            
            attachment.metadata.update({
                "lines": log_stats["total_lines"],
                "errors": log_stats["error_lines"],
                "warnings": log_stats["warning_lines"],
                "info": log_stats["info_lines"],
                "debug": log_stats["debug_lines"],
                "time_range": f"{log_stats.get('first_timestamp', 'N/A')} to {log_stats.get('last_timestamp', 'N/A')}"
            })
            
            content = f"[Log File: {attachment.filename}]\n"
            content += f"Total Lines: {log_stats['total_lines']:,}\n"
            content += f"Errors: {log_stats['error_lines']:,}\n"
            content += f"Warnings: {log_stats['warning_lines']:,}\n"
            content += f"Info: {log_stats['info_lines']:,}\n"
            content += f"Debug: {log_stats['debug_lines']:,}\n"
            content += f"Time Range: {attachment.metadata['time_range']}\n\n"
            
            if log_stats["error_messages"]:
                content += f"Error Messages ({len(log_stats['error_messages'])}):\n"
                for msg in log_stats["error_messages"][:20]:
                    content += f"  {msg[:200]}\n"
                content += "\n"
            
            if log_stats["warning_messages"]:
                content += f"Warning Messages ({len(log_stats['warning_messages'])}):\n"
                for msg in log_stats["warning_messages"][:20]:
                    content += f"  {msg[:200]}\n"
                content += "\n"
            
            content += "Log Preview:\n"
            content += "\n".join(content_preview[:500])
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Log processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process log: {attachment.filename}]"
        
        return attachment
    
    async def _process_binary(self, attachment: Attachment) -> Attachment:
        try:
            if MAGIC_AVAILABLE:
                file_magic = magic.from_file(attachment.path)
                file_mime = magic.from_file(attachment.path, mime=True)
            else:
                file_magic = "Binary file"
                file_mime = mimetypes.guess_type(attachment.path)[0] or "application/octet-stream"
            
            hex_preview = ""
            with open(attachment.path, 'rb') as f:
                data = f.read(1024)
                hex_preview = ' '.join(f'{b:02x}' for b in data[:256])
            
            strings_output = []
            try:
                import subprocess
                result = subprocess.run(['strings', '-n', '8', attachment.path], 
                                     capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    strings_output = result.stdout.splitlines()[:100]
            except:
                pass
            
            attachment.metadata.update({
                "magic": file_magic,
                "mime": file_mime,
                "hex_preview": hex_preview[:500],
                "strings_count": len(strings_output)
            })
            
            content = f"[Binary File: {attachment.filename}]\n"
            content += f"Type: {file_magic}\n"
            content += f"MIME: {file_mime}\n"
            content += f"Size: {attachment.size / 1024:.1f} KB\n\n"
            
            content += f"Hex Preview:\n{hex_preview[:256]}\n\n"
            
            if strings_output:
                content += f"Readable Strings ({len(strings_output)}):\n"
                for s in strings_output[:50]:
                    content += f"  {s[:100]}\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Binary processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process binary: {attachment.filename}]"
        
        return attachment
    
    async def _process_scientific(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['fits', 'fts']:
                hdul = fits.open(attachment.path)
                
                hdus_info = []
                for i, hdu in enumerate(hdul):
                    info = {
                        "index": i,
                        "name": hdu.name,
                        "type": type(hdu).__name__,
                        "header_cards": len(hdu.header)
                    }
                    
                    if hasattr(hdu, 'data') and hdu.data is not None:
                        info["shape"] = hdu.data.shape
                        info["dtype"] = str(hdu.data.dtype)
                    
                    hdus_info.append(info)
                
                hdul.close()
                
                attachment.metadata.update({
                    "format": "FITS",
                    "hdus": hdus_info,
                    "total_hdus": len(hdus_info)
                })
                
                content = f"[FITS File: {attachment.filename}]\n"
                content += f"Total HDUs: {len(hdus_info)}\n\n"
                
                for hdu in hdus_info[:10]:
                    content += f"HDU {hdu['index']}: {hdu['name']} ({hdu['type']})\n"
                    if 'shape' in hdu:
                        content += f"  Shape: {hdu['shape']}\n"
                        content += f"  Type: {hdu['dtype']}\n"
                    content += f"  Header Cards: {hdu['header_cards']}\n\n"
            
            elif ext in ['nc', 'nc4', 'netcdf']:
                ds = xr.open_dataset(attachment.path)
                
                variables = {}
                for var in ds.data_vars:
                    variables[var] = {
                        "shape": ds[var].shape,
                        "dtype": str(ds[var].dtype),
                        "dims": ds[var].dims,
                        "attrs": dict(ds[var].attrs)
                    }
                
                attachment.metadata.update({
                    "format": "NetCDF",
                    "variables": variables,
                    "dimensions": dict(ds.dims),
                    "attributes": dict(ds.attrs)
                })
                
                content = f"[NetCDF File: {attachment.filename}]\n"
                content += f"Variables: {len(variables)}\n"
                content += f"Dimensions: {dict(ds.dims)}\n\n"
                
                for var, info in list(variables.items())[:20]:
                    content += f"Variable: {var}\n"
                    content += f"  Shape: {info['shape']}\n"
                    content += f"  Dimensions: {info['dims']}\n"
                    content += f"  Type: {info['dtype']}\n\n"
                
                ds.close()
            
            elif ext in ['hdf', 'h4', 'h5']:
                with h5py.File(attachment.path, 'r') as f:
                    structure = self._explore_h5_structure(f)
                    
                    attachment.metadata.update({
                        "format": "HDF5",
                        "structure": structure
                    })
                    
                    content = f"[HDF5 File: {attachment.filename}]\n"
                    content += f"Structure:\n{self._format_h5_structure(structure)}"
            
            else:
                content = f"[Scientific File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024 / 1024:.2f} MB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Scientific file processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process scientific file: {attachment.filename}]"
        
        return attachment
    
    async def _process_medical(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext in ['dcm', 'dicom']:
                ds = pydicom.dcmread(attachment.path)
                
                patient_info = {}
                study_info = {}
                image_info = {}
                
                if hasattr(ds, 'PatientName'):
                    patient_info['name'] = str(ds.PatientName)
                if hasattr(ds, 'PatientID'):
                    patient_info['id'] = str(ds.PatientID)
                if hasattr(ds, 'PatientBirthDate'):
                    patient_info['birth_date'] = str(ds.PatientBirthDate)
                if hasattr(ds, 'PatientSex'):
                    patient_info['sex'] = str(ds.PatientSex)
                
                if hasattr(ds, 'StudyDescription'):
                    study_info['description'] = str(ds.StudyDescription)
                if hasattr(ds, 'StudyDate'):
                    study_info['date'] = str(ds.StudyDate)
                if hasattr(ds, 'Modality'):
                    study_info['modality'] = str(ds.Modality)
                
                if hasattr(ds, 'pixel_array'):
                    image_info['shape'] = ds.pixel_array.shape
                    image_info['dtype'] = str(ds.pixel_array.dtype)
                if hasattr(ds, 'Rows'):
                    image_info['rows'] = ds.Rows
                if hasattr(ds, 'Columns'):
                    image_info['columns'] = ds.Columns
                
                attachment.metadata.update({
                    "format": "DICOM",
                    "patient_info": patient_info,
                    "study_info": study_info,
                    "image_info": image_info
                })
                
                content = f"[DICOM File: {attachment.filename}]\n"
                content += f"Format: DICOM\n"
                if study_info.get('modality'):
                    content += f"Modality: {study_info['modality']}\n"
                if study_info.get('description'):
                    content += f"Study: {study_info['description']}\n"
                if image_info.get('shape'):
                    content += f"Image Shape: {image_info['shape']}\n"
            
            elif ext in ['nii', 'nii.gz']:
                img = nib.load(attachment.path)
                
                attachment.metadata.update({
                    "format": "NIfTI",
                    "shape": img.shape,
                    "affine": img.affine.tolist(),
                    "header": dict(img.header)
                })
                
                content = f"[NIfTI File: {attachment.filename}]\n"
                content += f"Format: NIfTI\n"
                content += f"Shape: {img.shape}\n"
                content += f"Data Type: {img.get_data_dtype()}\n"
            
            else:
                content = f"[Medical File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024 / 1024:.2f} MB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Medical file processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process medical file: {attachment.filename}]"
        
        return attachment
    
    async def _process_geospatial(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext == 'shp':
                gdf = gpd.read_file(attachment.path)
                
                attachment.metadata.update({
                    "format": "Shapefile",
                    "features": len(gdf),
                    "geometry_types": gdf.geometry.geom_type.value_counts().to_dict(),
                    "crs": str(gdf.crs),
                    "bounds": gdf.total_bounds.tolist(),
                    "columns": gdf.columns.tolist()
                })
                
                content = f"[Shapefile: {attachment.filename}]\n"
                content += f"Features: {len(gdf):,}\n"
                content += f"CRS: {gdf.crs}\n"
                content += f"Bounds: {gdf.total_bounds}\n"
                content += f"Columns: {', '.join(gdf.columns)}\n\n"
                
                content += "Geometry Types:\n"
                for geom_type, count in gdf.geometry.geom_type.value_counts().items():
                    content += f"  {geom_type}: {count:,}\n"
            
            elif ext in ['kml', 'kmz']:
                import fiona
                
                with fiona.open(attachment.path, 'r') as src:
                    features = list(src)
                    
                    attachment.metadata.update({
                        "format": "KML/KMZ",
                        "features": len(features),
                        "crs": str(src.crs),
                        "bounds": src.bounds
                    })
                    
                    content = f"[KML/KMZ File: {attachment.filename}]\n"
                    content += f"Features: {len(features):,}\n"
                    content += f"CRS: {src.crs}\n"
                    content += f"Bounds: {src.bounds}\n"
            
            elif ext == 'geojson':
                with open(attachment.path, 'r') as f:
                    geojson_data = json.load(f)
                
                feature_count = 0
                if geojson_data.get('type') == 'FeatureCollection':
                    feature_count = len(geojson_data.get('features', []))
                elif geojson_data.get('type') == 'Feature':
                    feature_count = 1
                
                attachment.metadata.update({
                    "format": "GeoJSON",
                    "type": geojson_data.get('type'),
                    "features": feature_count
                })
                
                content = f"[GeoJSON File: {attachment.filename}]\n"
                content += f"Type: {geojson_data.get('type')}\n"
                content += f"Features: {feature_count:,}\n"
            
            elif ext in ['tif', 'tiff']:
                with rasterio.open(attachment.path) as src:
                    attachment.metadata.update({
                        "format": "GeoTIFF",
                        "width": src.width,
                        "height": src.height,
                        "bands": src.count,
                        "crs": str(src.crs),
                        "bounds": src.bounds,
                        "dtype": str(src.dtypes[0])
                    })
                    
                    content = f"[GeoTIFF File: {attachment.filename}]\n"
                    content += f"Size: {src.width}x{src.height}\n"
                    content += f"Bands: {src.count}\n"
                    content += f"CRS: {src.crs}\n"
                    content += f"Bounds: {src.bounds}\n"
                    content += f"Data Type: {src.dtypes[0]}\n"
            
            else:
                content = f"[Geospatial File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024 / 1024:.2f} MB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Geospatial file processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process geospatial file: {attachment.filename}]"
        
        return attachment
    
    async def _process_molecular(self, attachment: Attachment) -> Attachment:
        try:
            ext = attachment.filename.lower().split('.')[-1]
            
            if ext == 'pdb':
                with open(attachment.path, 'r') as f:
                    pdb_content = f.read()
                
                atoms = []
                chains = set()
                residues = set()
                
                for line in pdb_content.splitlines():
                    if line.startswith('ATOM') or line.startswith('HETATM'):
                        atom_info = {
                            'name': line[12:16].strip(),
                            'residue': line[17:20].strip(),
                            'chain': line[21],
                            'residue_number': line[22:26].strip(),
                            'x': float(line[30:38]),
                            'y': float(line[38:46]),
                            'z': float(line[46:54])
                        }
                        atoms.append(atom_info)
                        chains.add(atom_info['chain'])
                        residues.add(f"{atom_info['residue']}_{atom_info['residue_number']}")
                
                attachment.metadata.update({
                    "format": "PDB",
                    "atoms": len(atoms),
                    "chains": list(chains),
                    "residues": len(residues)
                })
                
                content = f"[PDB File: {attachment.filename}]\n"
                content += f"Atoms: {len(atoms):,}\n"
                content += f"Chains: {', '.join(sorted(chains))}\n"
                content += f"Residues: {len(residues)}\n"
            
            elif ext in ['mol', 'mol2', 'sdf']:
                mol = Chem.MolFromMolFile(attachment.path) if ext == 'mol' else Chem.MolFromMol2File(attachment.path)
                
                if mol:
                    attachment.metadata.update({
                        "format": ext.upper(),
                        "atoms": mol.GetNumAtoms(),
                        "bonds": mol.GetNumBonds(),
                        "molecular_weight": Chem.Descriptors.MolWt(mol),
                        "formula": Chem.rdMolDescriptors.CalcMolFormula(mol)
                    })
                    
                    content = f"[Molecule File: {attachment.filename}]\n"
                    content += f"Format: {ext.upper()}\n"
                    content += f"Atoms: {mol.GetNumAtoms()}\n"
                    content += f"Bonds: {mol.GetNumBonds()}\n"
                    content += f"Formula: {Chem.rdMolDescriptors.CalcMolFormula(mol)}\n"
                    content += f"Molecular Weight: {Chem.Descriptors.MolWt(mol):.2f}\n"
            
            elif ext == 'fasta':
                sequences = []
                with open(attachment.path, 'r') as f:
                    for record in Bio.SeqIO.parse(f, "fasta"):
                        sequences.append({
                            "id": record.id,
                            "description": record.description,
                            "length": len(record.seq),
                            "sequence": str(record.seq)[:100]
                        })
                
                attachment.metadata.update({
                    "format": "FASTA",
                    "sequences": len(sequences),
                    "total_length": sum(s["length"] for s in sequences)
                })
                
                content = f"[FASTA File: {attachment.filename}]\n"
                content += f"Sequences: {len(sequences)}\n"
                content += f"Total Length: {sum(s['length'] for s in sequences):,}\n\n"
                
                for seq in sequences[:10]:
                    content += f">{seq['id']} {seq['description']}\n"
                    content += f"Length: {seq['length']:,}\n"
                    content += f"Preview: {seq['sequence']}...\n\n"
            
            else:
                content = f"[Molecular File: {attachment.filename}]\n"
                content += f"Format: {ext.upper()}\n"
                content += f"Size: {attachment.size / 1024:.1f} KB\n"
            
            attachment.extracted_text = content[:500000]
            attachment.processed = True
            
        except Exception as e:
            attachment.error = f"Molecular file processing error: {str(e)}"
            attachment.extracted_text = f"[Failed to process molecular file: {attachment.filename}]"
        
        return attachment
    
    async def _analyze_image(self, attachment: Attachment) -> Attachment:
        try:
            img = Image.open(attachment.path)
            img_array = np.array(img)
            
            histogram = {}
            if len(img_array.shape) == 3:
                for i, channel in enumerate(['red', 'green', 'blue']):
                    if i < img_array.shape[2]:
                        hist, _ = np.histogram(img_array[:,:,i], bins=256, range=(0, 256))
                        histogram[channel] = hist.tolist()
            else:
                hist, _ = np.histogram(img_array, bins=256, range=(0, 256))
                histogram['gray'] = hist.tolist()
            
            brightness = np.mean(img_array)
            contrast = np.std(img_array)
            
            edges = cv2.Canny(cv2.cvtColor(np.array(img), cv2.COLOR_RGB2GRAY), 100, 200)
            edge_density = np.sum(edges > 0) / edges.size
            
            attachment.analysis = {
                "histogram": histogram,
                "brightness": float(brightness),
                "contrast": float(contrast),
                "edge_density": float(edge_density),
                "dominant_colors": self._get_dominant_colors(img_array)
            }
            
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_document(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.extracted_text:
                text = attachment.extracted_text
                
                word_count = len(text.split())
                char_count = len(text)
                sentence_count = len(re.split(r'[.!?]+', text))
                
                readability_scores = {
                    "flesch_reading_ease": textstat.flesch_reading_ease(text),
                    "flesch_kincaid_grade": textstat.flesch_kincaid_grade(text),
                    "gunning_fog": textstat.gunning_fog(text),
                    "automated_readability_index": textstat.automated_readability_index(text)
                }
                
                language = 'en'
                try:
                    from langdetect import detect
                    language = detect(text[:1000])
                except:
                    pass
                
                attachment.analysis = {
                    "word_count": word_count,
                    "character_count": char_count,
                    "sentence_count": sentence_count,
                    "average_word_length": char_count / word_count if word_count > 0 else 0,
                    "average_sentence_length": word_count / sentence_count if sentence_count > 0 else 0,
                    "readability": readability_scores,
                    "language": language
                }
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_spreadsheet(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("sample_data"):
                df = pd.DataFrame(attachment.metadata["sample_data"])
                
                data_quality = {
                    "completeness": (1 - df.isnull().sum().sum() / (df.shape[0] * df.shape[1])) * 100,
                    "unique_values": {col: df[col].nunique() for col in df.columns},
                    "data_types": {col: str(df[col].dtype) for col in df.columns}
                }
                
                patterns = {}
                for col in df.select_dtypes(include=['object']).columns[:10]:
                    value_counts = df[col].value_counts()
                    if len(value_counts) < 50:
                        patterns[col] = value_counts.to_dict()
                
                attachment.analysis = {
                    "data_quality": data_quality,
                    "patterns": patterns
                }
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_code(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("complexity"):
                metrics = attachment.metadata["complexity"]
                
                code_quality = {
                    "maintainability_index": self._calculate_maintainability_index(metrics),
                    "technical_debt_ratio": self._calculate_technical_debt(metrics),
                    "code_smell_count": self._detect_code_smells(attachment.extracted_text)
                }
                
                attachment.analysis = {
                    "quality_metrics": code_quality,
                    "recommendations": self._generate_code_recommendations(metrics, code_quality)
                }
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_audio(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("mfcc_means"):
                audio_features = {
                    "tempo_category": self._categorize_tempo(attachment.metadata.get("tempo", 0)),
                    "energy_level": self._calculate_energy_level(attachment.metadata),
                    "spectral_characteristics": self._analyze_spectral_features(attachment.metadata)
                }
                
                attachment.analysis = audio_features
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_video(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("sample_frames"):
                video_analysis = {
                    "quality_assessment": self._assess_video_quality(attachment.metadata),
                    "content_type": self._detect_video_content_type(attachment.metadata),
                    "compression_efficiency": self._calculate_compression_efficiency(attachment)
                }
                
                attachment.analysis = video_analysis
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_scientific(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("variables") or attachment.metadata.get("hdus"):
                scientific_analysis = {
                    "data_dimensions": self._analyze_data_dimensions(attachment.metadata),
                    "data_types": self._analyze_data_types(attachment.metadata),
                    "metadata_completeness": self._assess_metadata_completeness(attachment.metadata)
                }
                
                attachment.analysis = scientific_analysis
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_medical(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("study_info"):
                medical_analysis = {
                    "modality_specific": self._analyze_modality(attachment.metadata),
                    "data_integrity": self._check_medical_data_integrity(attachment.metadata),
                    "compliance_check": self._check_dicom_compliance(attachment.metadata)
                }
                
                attachment.analysis = medical_analysis
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_geospatial(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("bounds") or attachment.metadata.get("features"):
                geo_analysis = {
                    "spatial_extent": self._analyze_spatial_extent(attachment.metadata),
                    "coordinate_system": self._analyze_coordinate_system(attachment.metadata),
                    "data_density": self._calculate_spatial_density(attachment.metadata)
                }
                
                attachment.analysis = geo_analysis
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    async def _analyze_molecular(self, attachment: Attachment) -> Attachment:
        try:
            if attachment.metadata.get("atoms") or attachment.metadata.get("sequences"):
                molecular_analysis = {
                    "structure_complexity": self._analyze_molecular_complexity(attachment.metadata),
                    "composition": self._analyze_molecular_composition(attachment.metadata),
                    "properties": self._predict_molecular_properties(attachment.metadata)
                }
                
                attachment.analysis = molecular_analysis
        
        except Exception as e:
            attachment.analysis = {"error": str(e)}
        
        return attachment
    
    def _detect_encoding(self, file_path: str) -> str:
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)
            
            result = chardet.detect(raw_data)
            encoding = result['encoding']
            
            if encoding and result['confidence'] > 0.7:
                return encoding
            
            return 'utf-8'
        except:
            return 'utf-8'
    
    def _analyze_json_structure(self, data: Any, max_depth: int = 5, current_depth: int = 0) -> Dict:
        if current_depth >= max_depth:
            return {"type": "truncated", "reason": "max_depth_exceeded"}
        
        if isinstance(data, dict):
            structure = {
                "type": "object",
                "keys": list(data.keys())[:100],
                "size": len(data),
                "sample_values": {}
            }
            
            for key in list(data.keys())[:10]:
                structure["sample_values"][key] = self._analyze_json_structure(
                    data[key], max_depth, current_depth + 1
                )
            
            return structure
            
        elif isinstance(data, list):
            structure = {
                "type": "array",
                "length": len(data),
                "item_types": list(set(type(item).__name__ for item in data[:100]))
            }
            
            if data:
                structure["sample_structure"] = self._analyze_json_structure(
                    data[0], max_depth, current_depth + 1
                )
            
            return structure
            
        else:
            return {
                "type": type(data).__name__,
                "value": str(data)[:100] if not isinstance(data, (int, float, bool, type(None))) else data
            }
    
    def _calculate_code_complexity(self, code: str, language: str) -> Dict[str, Any]:
        metrics = {
            "lines_of_code": len([line for line in code.splitlines() if line.strip() and not line.strip().startswith(('#', '//', '/*'))]),
            "cyclomatic_complexity": 1,
            "cognitive_complexity": 0,
            "nesting_depth": 0,
            "function_count": 0,
            "class_count": 0
        }
        
        control_flow_keywords = ['if', 'elif', 'else', 'for', 'while', 'switch', 'case', 'catch', 'except']
        for keyword in control_flow_keywords:
            metrics["cyclomatic_complexity"] += len(re.findall(r'\b' + keyword + r'\b', code))
        
        logical_operators = ['&&', '||', 'and', 'or']
        for op in logical_operators:
            metrics["cyclomatic_complexity"] += len(re.findall(re.escape(op), code))
        
        lines = code.splitlines()
        current_depth = 0
        max_depth = 0
        
        for line in lines:
            stripped = line.lstrip()
            if stripped:
                indent = len(line) - len(stripped)
                depth = indent // 4 if language in ['python'] else line.count('{') - line.count('}')
                current_depth = max(0, current_depth + depth)
                max_depth = max(max_depth, current_depth)
        
        metrics["nesting_depth"] = max_depth
        
        if language in ['python', 'py']:
            metrics["function_count"] = len(re.findall(r'^def\s+\w+', code, re.MULTILINE))
            metrics["class_count"] = len(re.findall(r'^class\s+\w+', code, re.MULTILINE))
        elif language in ['javascript', 'js', 'typescript', 'ts']:
            metrics["function_count"] = len(re.findall(r'function\s+\w+|const\s+\w+\s*=\s*(?:\([^)]*\)\s*=>|function)', code))
            metrics["class_count"] = len(re.findall(r'class\s+\w+', code))
        elif language in ['java', 'c++', 'cpp', 'c#', 'cs']:
            metrics["function_count"] = len(re.findall(r'(?:public|private|protected)?\s*(?:static)?\s*\w+\s+\w+\s*\([^)]*\)\s*{', code))
            metrics["class_count"] = len(re.findall(r'class\s+\w+', code))
        
        return metrics
    
    def _explore_h5_structure(self, group, max_depth: int = 5, current_depth: int = 0) -> Dict:
        if current_depth >= max_depth:
            return {"type": "truncated", "children": {}}
        
        structure = {
            "type": "group",
            "attrs": dict(group.attrs),
            "children": {}
        }
        
        for key in list(group.keys())[:100]:
            item = group[key]
            if isinstance(item, h5py.Group):
                structure["children"][key] = self._explore_h5_structure(
                    item, max_depth, current_depth + 1
                )
            elif isinstance(item, h5py.Dataset):
                structure["children"][key] = {
                    "type": "dataset",
                    "shape": item.shape,
                    "dtype": str(item.dtype),
                    "attrs": dict(item.attrs)
                }
        
        return structure
    
    def _format_h5_structure(self, structure: Dict, indent: int = 0) -> str:
        lines = []
        indent_str = "  " * indent
        
        if structure["type"] == "group":
            if structure.get("attrs"):
                lines.append(f"{indent_str}Attributes: {structure['attrs']}")
            
            for name, child in structure.get("children", {}).items():
                if child["type"] == "group":
                    lines.append(f"{indent_str}{name}/")
                    lines.append(self._format_h5_structure(child, indent + 1))
                else:
                    lines.append(f"{indent_str}{name}: {child['shape']} {child['dtype']}")
        
        return "\n".join(lines)
    
    def _count_h5_datasets(self, group) -> int:
        count = 0
        for item in group.values():
            if isinstance(item, h5py.Dataset):
                count += 1
            elif isinstance(item, h5py.Group):
                count += self._count_h5_datasets(item)
        return count
    
    def _count_h5_groups(self, group) -> int:
        count = 1
        for item in group.values():
            if isinstance(item, h5py.Group):
                count += self._count_h5_groups(item)
        return count
    
    def _get_dominant_colors(self, img_array: np.ndarray, n_colors: int = 5) -> List[Dict]:
        if len(img_array.shape) == 3:
            pixels = img_array.reshape(-1, img_array.shape[2])
        else:
            pixels = img_array.reshape(-1, 1)
        
        from sklearn.cluster import KMeans
        
        sample_size = min(10000, len(pixels))
        sample_indices = np.random.choice(len(pixels), sample_size, replace=False)
        sample_pixels = pixels[sample_indices]
        
        kmeans = KMeans(n_clusters=n_colors, random_state=42, n_init=10)
        kmeans.fit(sample_pixels)
        
        colors = kmeans.cluster_centers_
        labels = kmeans.labels_
        
        color_counts = np.bincount(labels)
        color_percentages = color_counts / len(labels)
        
        dominant_colors = []
        for i, (color, percentage) in enumerate(zip(colors, color_percentages)):
            color_dict = {
                "percentage": float(percentage * 100),
                "rgb": [int(c) for c in color]
            }
            if len(color) >= 3:
                color_dict["hex"] = '#{:02x}{:02x}{:02x}'.format(int(color[0]), int(color[1]), int(color[2]))
            dominant_colors.append(color_dict)
        
        return sorted(dominant_colors, key=lambda x: x['percentage'], reverse=True)
    
    def _calculate_maintainability_index(self, metrics: Dict) -> float:
        loc = metrics.get("lines_of_code", 1)
        complexity = metrics.get("cyclomatic_complexity", 1)
        
        mi = 171 - 5.2 * np.log(loc) - 0.23 * complexity - 16.2 * np.log(loc)
        return max(0, min(100, mi))
    
    def _calculate_technical_debt(self, metrics: Dict) -> float:
        complexity = metrics.get("cyclomatic_complexity", 1)
        nesting = metrics.get("nesting_depth", 0)
        loc = metrics.get("lines_of_code", 1)
        
        debt_ratio = (complexity / loc) * 10 + (nesting / 5) * 5
        return min(100, debt_ratio)
    
    def _detect_code_smells(self, code: str) -> int:
        smells = 0
        
        lines = code.splitlines()
        for line in lines:
            if len(line) > 120:
                smells += 1
            
            if line.count('(') > 5:
                smells += 1
            
            if re.search(r'(TODO|FIXME|HACK|XXX)', line, re.IGNORECASE):
                smells += 1
        
        functions = re.findall(r'def\s+\w+\s*\([^)]*\):|function\s+\w+\s*\([^)]*\)', code)
        for func in functions:
            param_count = len(re.findall(r'\w+(?:\s*[:=][^,)]+)?(?:\s*,|\s*\))', func))
            if param_count > 5:
                smells += 1
        
        duplicate_lines = {}
        for i, line in enumerate(lines):
            if len(line.strip()) > 10:
                if line in duplicate_lines:
                    duplicate_lines[line].append(i)
                else:
                    duplicate_lines[line] = [i]
        
        for line, occurrences in duplicate_lines.items():
            if len(occurrences) > 3:
                smells += len(occurrences) - 3
        
        return smells
    
    def _generate_code_recommendations(self, metrics: Dict, quality: Dict) -> List[str]:
        recommendations = []
        
        if metrics.get("cyclomatic_complexity", 0) > 10:
            recommendations.append("Consider refactoring complex functions to reduce cyclomatic complexity")
        
        if metrics.get("nesting_depth", 0) > 4:
            recommendations.append("Reduce nesting depth by extracting nested logic into separate functions")
        
        if quality.get("maintainability_index", 100) < 50:
            recommendations.append("Low maintainability index - consider breaking down into smaller, more focused modules")
        
        if quality.get("technical_debt_ratio", 0) > 30:
            recommendations.append("High technical debt detected - prioritize refactoring efforts")
        
        if quality.get("code_smell_count", 0) > 10:
            recommendations.append("Multiple code smells detected - review and address TODO comments and long lines")
        
        return recommendations
    
    def _categorize_tempo(self, tempo: float) -> str:
        if tempo < 60:
            return "Very Slow (Largo)"
        elif tempo < 76:
            return "Slow (Adagio)"
        elif tempo < 108:
            return "Moderate (Andante)"
        elif tempo < 120:
            return "Moderately Fast (Moderato)"
        elif tempo < 168:
            return "Fast (Allegro)"
        elif tempo < 200:
            return "Very Fast (Presto)"
        else:
            return "Extremely Fast (Prestissimo)"
    
    def _calculate_energy_level(self, metadata: Dict) -> str:
        spectral_centroid = metadata.get("spectral_centroid_mean", 0)
        zero_crossing = metadata.get("zero_crossing_rate_mean", 0)
        
        energy_score = (spectral_centroid / 4000) + (zero_crossing * 100)
        
        if energy_score < 0.3:
            return "Low Energy"
        elif energy_score < 0.6:
            return "Medium Energy"
        else:
            return "High Energy"
    
    def _analyze_spectral_features(self, metadata: Dict) -> Dict:
        return {
            "brightness": "Bright" if metadata.get("spectral_centroid_mean", 0) > 2000 else "Dark",
            "texture": "Smooth" if metadata.get("zero_crossing_rate_mean", 0) < 0.05 else "Rough",
            "harmonicity": self._estimate_harmonicity(metadata.get("mfcc_means", []))
        }
    
    def _estimate_harmonicity(self, mfcc_means: List[float]) -> str:
        if not mfcc_means:
            return "Unknown"
        
        variance = np.var(mfcc_means)
        if variance < 10:
            return "Highly Harmonic"
        elif variance < 50:
            return "Moderately Harmonic"
        else:
            return "Inharmonic/Noisy"
    
    def _assess_video_quality(self, metadata: Dict) -> Dict:
        width = metadata.get("width", 0)
        height = metadata.get("height", 0)
        fps = metadata.get("fps", 0)
        
        resolution_score = min(width * height / (1920 * 1080), 2.0)
        fps_score = min(fps / 30, 2.0)
        
        quality_score = (resolution_score + fps_score) / 2
        
        if quality_score > 1.5:
            quality = "High Quality"
        elif quality_score > 0.8:
            quality = "Medium Quality"
        else:
            quality = "Low Quality"
        
        return {
            "overall": quality,
            "resolution": f"{width}x{height}",
            "resolution_category": self._categorize_resolution(width, height),
            "fps_category": self._categorize_fps(fps)
        }
    
    def _categorize_resolution(self, width: int, height: int) -> str:
        pixels = width * height
        
        if pixels >= 3840 * 2160:
            return "4K Ultra HD"
        elif pixels >= 1920 * 1080:
            return "Full HD"
        elif pixels >= 1280 * 720:
            return "HD"
        elif pixels >= 854 * 480:
            return "SD"
        else:
            return "Low Resolution"
    
    def _categorize_fps(self, fps: float) -> str:
        if fps >= 120:
            return "High Frame Rate"
        elif fps >= 60:
            return "Smooth"
        elif fps >= 24:
            return "Standard"
        else:
            return "Low Frame Rate"
    
    def _detect_video_content_type(self, metadata: Dict) -> str:
        return "General Video Content"
    
    def _calculate_compression_efficiency(self, attachment: Attachment) -> Dict:
        if attachment.metadata.get("duration", 0) > 0:
            bitrate = (attachment.size * 8) / attachment.metadata["duration"] / 1000
            
            if bitrate < 1000:
                efficiency = "Highly Compressed"
            elif bitrate < 5000:
                efficiency = "Well Compressed"
            elif bitrate < 10000:
                efficiency = "Moderately Compressed"
            else:
                efficiency = "Low Compression"
            
            return {
                "bitrate_kbps": bitrate,
                "efficiency": efficiency
            }
        
        return {"efficiency": "Unknown"}
    
    def _analyze_data_dimensions(self, metadata: Dict) -> Dict:
        dimensions = {}
        
        if "variables" in metadata:
            for var, info in metadata["variables"].items():
                dimensions[var] = info.get("shape", [])
        
        elif "hdus" in metadata:
            for hdu in metadata["hdus"]:
                if "shape" in hdu:
                    dimensions[f"HDU_{hdu['index']}"] = hdu["shape"]
        
        return dimensions
    
    def _analyze_data_types(self, metadata: Dict) -> Dict:
        types = {}
        
        if "variables" in metadata:
            for var, info in metadata["variables"].items():
                types[var] = info.get("dtype", "unknown")
        
        elif "hdus" in metadata:
            for hdu in metadata["hdus"]:
                if "dtype" in hdu:
                    types[f"HDU_{hdu['index']}"] = hdu["dtype"]
        
        return types
    
    def _assess_metadata_completeness(self, metadata: Dict) -> float:
        required_fields = ["format", "size", "variables", "dimensions", "attributes"]
        present_fields = sum(1 for field in required_fields if field in metadata)
        
        return (present_fields / len(required_fields)) * 100
    
    def _analyze_modality(self, metadata: Dict) -> Dict:
        modality = metadata.get("study_info", {}).get("modality", "Unknown")
        
        modality_info = {
            "CT": {"type": "Computed Tomography", "uses": "X-rays", "contrast": "Often used"},
            "MR": {"type": "Magnetic Resonance", "uses": "Magnetic fields", "contrast": "Sometimes used"},
            "US": {"type": "Ultrasound", "uses": "Sound waves", "contrast": "Rarely used"},
            "XR": {"type": "X-Ray", "uses": "X-rays", "contrast": "Not used"},
            "PET": {"type": "Positron Emission Tomography", "uses": "Radioactive tracers", "contrast": "Always used"}
        }
        
        return modality_info.get(modality, {"type": modality, "uses": "Unknown", "contrast": "Unknown"})
    
    def _check_medical_data_integrity(self, metadata: Dict) -> Dict:
        integrity = {
            "patient_info_complete": all(k in metadata.get("patient_info", {}) for k in ["name", "id", "birth_date"]),
            "study_info_complete": all(k in metadata.get("study_info", {}) for k in ["description", "date", "modality"]),
            "image_info_complete": "shape" in metadata.get("image_info", {})
        }
        
        integrity["overall"] = all(integrity.values())
        
        return integrity
    
    def _check_dicom_compliance(self, metadata: Dict) -> Dict:
        return {
            "has_required_tags": True,
            "valid_modality": metadata.get("study_info", {}).get("modality", "") in ["CT", "MR", "US", "XR", "PET", "NM", "MG", "DX"],
            "has_patient_info": bool(metadata.get("patient_info")),
            "has_study_info": bool(metadata.get("study_info"))
        }
    
    def _analyze_spatial_extent(self, metadata: Dict) -> Dict:
        if "bounds" in metadata:
            bounds = metadata["bounds"]
            if isinstance(bounds, (list, tuple)) and len(bounds) >= 4:
                return {
                    "min_x": bounds[0],
                    "min_y": bounds[1],
                    "max_x": bounds[2],
                    "max_y": bounds[3],
                    "width": bounds[2] - bounds[0],
                    "height": bounds[3] - bounds[1]
                }
        
        return {"extent": "Unknown"}
    
    def _analyze_coordinate_system(self, metadata: Dict) -> Dict:
        crs = metadata.get("crs", "Unknown")
        
        common_crs = {
            "EPSG:4326": "WGS 84 (Geographic)",
            "EPSG:3857": "Web Mercator",
            "EPSG:32633": "UTM Zone 33N",
            "EPSG:2154": "Lambert-93 (France)"
        }
        
        return {
            "crs": str(crs),
            "description": common_crs.get(str(crs), "Custom or Unknown CRS"),
            "is_geographic": "4326" in str(crs) or "WGS" in str(crs)
        }
    
    def _calculate_spatial_density(self, metadata: Dict) -> Dict:
        features = metadata.get("features", 0)
        
        if "bounds" in metadata and features > 0:
            bounds = metadata["bounds"]
            if isinstance(bounds, (list, tuple)) and len(bounds) >= 4:
                area = (bounds[2] - bounds[0]) * (bounds[3] - bounds[1])
                if area > 0:
                    density = features / area
                    return {
                        "features_per_unit": density,
                        "density_category": "High" if density > 100 else "Medium" if density > 10 else "Low"
                    }
        
        return {"density": "Unknown"}
    
    def _analyze_molecular_complexity(self, metadata: Dict) -> Dict:
        atoms = metadata.get("atoms", 0)
        bonds = metadata.get("bonds", 0)
        
        if atoms > 0:
            complexity_score = atoms + (bonds * 1.5)
            
            if complexity_score < 50:
                complexity = "Simple"
            elif complexity_score < 200:
                complexity = "Moderate"
            elif complexity_score < 1000:
                complexity = "Complex"
            else:
                complexity = "Very Complex"
            
            return {
                "complexity": complexity,
                "atom_count": atoms,
                "bond_count": bonds,
                "complexity_score": complexity_score
            }
        
        return {"complexity": "Unknown"}
    
    def _analyze_molecular_composition(self, metadata: Dict) -> Dict:
        if "formula" in metadata:
            formula = metadata["formula"]
            
            elements = re.findall(r'([A-Z][a-z]?)(\d*)', formula)
            composition = {}
            
            for element, count in elements:
                count = int(count) if count else 1
                composition[element] = count
            
            return {
                "elements": composition,
                "unique_elements": len(composition),
                "total_atoms": sum(composition.values())
            }
        
        return {"composition": "Unknown"}
    
    def _predict_molecular_properties(self, metadata: Dict) -> Dict:
        properties = {}
        
        if "molecular_weight" in metadata:
            mw = metadata["molecular_weight"]
            properties["molecular_weight"] = mw
            properties["size_category"] = "Small" if mw < 500 else "Medium" if mw < 1000 else "Large"
        
        if "formula" in metadata:
            formula = metadata["formula"]
            properties["organic"] = "C" in formula
            properties["contains_nitrogen"] = "N" in formula
            properties["contains_oxygen"] = "O" in formula
            properties["contains_sulfur"] = "S" in formula
        
        return properties

session_manager = UltraSessionManager()
key_manager = UltraGeminiKeyManager()
client_manager = UltraStreamingClient(key_manager)
prompt_engineer = UltraPromptEngineer()
file_processor = UltraFileProcessor()

ULTIMATE_HTML_TEMPLATE = r"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Jack's AI Ultra - Ultimate Edition</title>
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=JetBrains+Mono:wght@400;500;600;700&family=Space+Grotesk:wght@400;500;600;700&display=swap" rel="stylesheet">
    <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.1/css/all.min.css" rel="stylesheet">
    <link href="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/themes/prism-tomorrow.min.css" rel="stylesheet">
    <link href="https://cdnjs.cloudflare.com/ajax/libs/animate.css/4.1.1/animate.min.css" rel="stylesheet">
    <style>
        :root {
            --primary: #6366f1;
            --primary-dark: #4f46e5;
            --primary-light: #818cf8;
            --secondary: #8b5cf6;
            --accent: #ec4899;
            --success: #10b981;
            --warning: #f59e0b;
            --danger: #ef4444;
            --info: #3b82f6;
            
            --bg-primary: #0a0a0a;
            --bg-secondary: #111111;
            --bg-tertiary: #1a1a1a;
            --bg-card: rgba(255, 255, 255, 0.02);
            --bg-hover: rgba(255, 255, 255, 0.05);
            --bg-active: rgba(99, 102, 241, 0.1);
            
            --text-primary: #ffffff;
            --text-secondary: #a0a0a0;
            --text-muted: #666666;
            
            --border-color: rgba(255, 255, 255, 0.1);
            --border-hover: rgba(255, 255, 255, 0.2);
            --border-active: var(--primary);
            
            --shadow-sm: 0 2px 8px rgba(0, 0, 0, 0.1);
            --shadow-md: 0 4px 20px rgba(0, 0, 0, 0.15);
            --shadow-lg: 0 10px 40px rgba(0, 0, 0, 0.2);
            --shadow-xl: 0 20px 60px rgba(0, 0, 0, 0.3);
            --shadow-glow: 0 0 40px rgba(99, 102, 241, 0.3);
            --shadow-glow-intense: 0 0 80px rgba(99, 102, 241, 0.5);
            
            --radius-sm: 8px;
            --radius-md: 12px;
            --radius-lg: 16px;
            --radius-xl: 24px;
            --radius-2xl: 32px;
            
            --transition-fast: 150ms cubic-bezier(0.4, 0, 0.2, 1);
            --transition-base: 300ms cubic-bezier(0.4, 0, 0.2, 1);
            --transition-slow: 500ms cubic-bezier(0.4, 0, 0.2, 1);
            --transition-spring: 600ms cubic-bezier(0.68, -0.55, 0.265, 1.55);
            
            --header-height: 70px;
            --sidebar-width: 360px;
            --input-height: 56px;
            
            --z-background: -1;
            --z-base: 1;
            --z-dropdown: 100;
            --z-sticky: 200;
            --z-fixed: 300;
            --z-modal-backdrop: 400;
            --z-modal: 500;
            --z-popover: 600;
            --z-tooltip: 700;
            --z-notification: 800;
            --z-top: 999;
        }
        
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }
        
        ::selection {
            background: var(--primary);
            color: white;
        }
        
        ::-webkit-scrollbar {
            width: 10px;
            height: 10px;
        }
        
        ::-webkit-scrollbar-track {
            background: var(--bg-secondary);
            border-radius: 5px;
        }
        
        ::-webkit-scrollbar-thumb {
            background: var(--border-color);
            border-radius: 5px;
            transition: background var(--transition-fast);
        }
        
        ::-webkit-scrollbar-thumb:hover {
            background: var(--border-hover);
        }
        
        ::-webkit-scrollbar-corner {
            background: var(--bg-secondary);
        }
        
        body {
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
            background: var(--bg-primary);
            color: var(--text-primary);
            min-height: 100vh;
            overflow: hidden;
            position: relative;
            -webkit-font-smoothing: antialiased;
            -moz-osx-font-smoothing: grayscale;
            font-feature-settings: "cv11", "ss01", "ss03";
        }
        
        .background-canvas {
            position: fixed;
            inset: 0;
            z-index: var(--z-background);
            opacity: 0.6;
        }
        
        .particle-field {
            position: absolute;
            inset: 0;
            overflow: hidden;
        }
        
        .particle {
            position: absolute;
            width: 2px;
            height: 2px;
            background: var(--primary);
            border-radius: 50%;
            opacity: 0;
            animation: particle-float 20s infinite;
        }
        
        @keyframes particle-float {
            0% {
                transform: translate(0, 100vh) scale(0);
                opacity: 0;
            }
            10% {
                opacity: 0.8;
                transform: translate(10px, 90vh) scale(1);
            }
            90% {
                opacity: 0.8;
                transform: translate(-10px, -90vh) scale(1);
            }
            100% {
                transform: translate(0, -100vh) scale(0);
                opacity: 0;
            }
        }
        
        .gradient-mesh {
            position: absolute;
            inset: 0;
            background: 
                radial-gradient(ellipse at 20% 30%, rgba(99, 102, 241, 0.15) 0%, transparent 50%),
                radial-gradient(ellipse at 80% 70%, rgba(139, 92, 246, 0.15) 0%, transparent 50%),
                radial-gradient(ellipse at 50% 50%, rgba(236, 72, 153, 0.1) 0%, transparent 70%);
            filter: blur(40px);
            animation: gradient-shift 30s ease-in-out infinite;
        }
        
        @keyframes gradient-shift {
            0%, 100% { transform: rotate(0deg) scale(1); }
            25% { transform: rotate(90deg) scale(1.1); }
            50% { transform: rotate(180deg) scale(1); }
            75% { transform: rotate(270deg) scale(1.1); }
        }
        
        .grid-overlay {
            position: absolute;
            inset: 0;
            background-image: 
                linear-gradient(rgba(255, 255, 255, 0.02) 1px, transparent 1px),
                linear-gradient(90deg, rgba(255, 255, 255, 0.02) 1px, transparent 1px);
            background-size: 50px 50px;
            animation: grid-move 60s linear infinite;
        }
        
        @keyframes grid-move {
            0% { transform: translate(0, 0); }
            100% { transform: translate(50px, 50px); }
        }
        
        .app-container {
            position: relative;
            z-index: var(--z-base);
            display: flex;
            height: 100vh;
            background: rgba(10, 10, 10, 0.8);
            backdrop-filter: blur(100px) saturate(150%);
            -webkit-backdrop-filter: blur(100px) saturate(150%);
        }
        
        .sidebar {
            width: var(--sidebar-width);
            background: rgba(17, 17, 17, 0.95);
            backdrop-filter: blur(20px);
            -webkit-backdrop-filter: blur(20px);
            border-right: 1px solid var(--border-color);
            display: flex;
            flex-direction: column;
            transition: transform var(--transition-base);
            position: relative;
            z-index: var(--z-sticky);
            box-shadow: 4px 0 24px rgba(0, 0, 0, 0.2);
        }
        
        .sidebar.collapsed {
            transform: translateX(-100%);
        }
        
        .sidebar-header {
            padding: 28px 24px;
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            position: relative;
            overflow: hidden;
            box-shadow: 0 4px 24px rgba(99, 102, 241, 0.2);
        }
        
        .sidebar-header::before {
            content: '';
            position: absolute;
            inset: 0;
            background: 
                radial-gradient(circle at 30% 50%, rgba(255, 255, 255, 0.2), transparent 50%),
                radial-gradient(circle at 70% 80%, rgba(255, 255, 255, 0.1), transparent 50%);
            animation: shimmer 4s ease-in-out infinite;
        }
        
        .sidebar-header::after {
            content: '';
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            height: 1px;
            background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.5), transparent);
            animation: scan 3s linear infinite;
        }
        
        @keyframes shimmer {
            0%, 100% { opacity: 0.5; }
            50% { opacity: 1; }
        }
        
        @keyframes scan {
            0% { transform: translateX(-100%); }
            100% { transform: translateX(100%); }
        }
        
        .logo-container {
            display: flex;
            align-items: center;
            gap: 16px;
            position: relative;
            z-index: 1;
        }
        
        .logo-icon {
            width: 56px;
            height: 56px;
            background: rgba(255, 255, 255, 0.15);
            backdrop-filter: blur(10px);
            -webkit-backdrop-filter: blur(10px);
            border-radius: var(--radius-lg);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 28px;
            color: white;
            box-shadow: 
                0 8px 32px rgba(0, 0, 0, 0.2),
                inset 0 1px 0 rgba(255, 255, 255, 0.2);
            animation: pulse 3s ease-in-out infinite;
            position: relative;
            overflow: hidden;
        }
        
        .logo-icon::before {
            content: '';
            position: absolute;
            inset: -2px;
            background: linear-gradient(45deg, var(--primary), var(--secondary), var(--accent));
            border-radius: var(--radius-lg);
            opacity: 0;
            animation: logo-glow 3s ease-in-out infinite;
            z-index: -1;
            filter: blur(10px);
        }
        
        @keyframes logo-glow {
            0%, 100% { opacity: 0; transform: scale(0.9); }
            50% { opacity: 0.6; transform: scale(1.1); }
        }
        
        @keyframes pulse {
            0%, 100% { transform: scale(1); }
            50% { transform: scale(1.05); }
        }
        
        .logo-text h1 {
            font-family: 'Space Grotesk', sans-serif;
            font-size: 24px;
            font-weight: 700;
            color: white;
            letter-spacing: -0.5px;
            text-shadow: 0 2px 10px rgba(0, 0, 0, 0.3);
        }
        
        .logo-text p {
            font-size: 13px;
            color: rgba(255, 255, 255, 0.9);
            margin-top: 2px;
            font-weight: 500;
            letter-spacing: 2px;
            text-transform: uppercase;
        }
        
        .sidebar-nav {
            padding: 20px 16px;
            flex: 0 0 auto;
        }
        
        .nav-section {
            margin-bottom: 24px;
        }
        
        .nav-section-title {
            font-size: 11px;
            font-weight: 600;
            color: var(--text-muted);
            text-transform: uppercase;
            letter-spacing: 1px;
            margin-bottom: 12px;
            padding: 0 8px;
        }
        
        .nav-item {
            display: flex;
            align-items: center;
            gap: 12px;
            padding: 12px 16px;
            margin-bottom: 4px;
            background: transparent;
            border: 1px solid transparent;
            border-radius: var(--radius-md);
            color: var(--text-secondary);
            text-decoration: none;
            font-size: 14px;
            font-weight: 500;
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
            overflow: hidden;
        }
        
        .nav-item::before {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(90deg, transparent, rgba(99, 102, 241, 0.1), transparent);
            transform: translateX(-100%);
            transition: transform 0.6s;
        }
        
        .nav-item:hover {
            background: var(--bg-hover);
            color: var(--text-primary);
            border-color: var(--border-color);
            transform: translateX(4px);
        }
        
        .nav-item:hover::before {
            transform: translateX(100%);
        }
        
        .nav-item.active {
            background: var(--bg-active);
            color: var(--primary);
            border-color: var(--primary);
            box-shadow: 0 0 20px rgba(99, 102, 241, 0.2);
        }
        
        .nav-item i {
            width: 20px;
            text-align: center;
            font-size: 16px;
        }
        
        .nav-badge {
            margin-left: auto;
            padding: 2px 8px;
            background: var(--primary);
            color: white;
            font-size: 11px;
            font-weight: 600;
            border-radius: 10px;
        }
        
        .chat-sessions {
            flex: 1;
            overflow-y: auto;
            padding: 16px;
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) transparent;
        }
        
        .sessions-header {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 16px;
            padding: 0 4px;
        }
        
        .sessions-title {
            font-size: 13px;
            font-weight: 600;
            color: var(--text-secondary);
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .sessions-actions {
            display: flex;
            gap: 8px;
        }
        
        .btn-icon {
            width: 32px;
            height: 32px;
            border-radius: var(--radius-sm);
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            color: var(--text-secondary);
            display: flex;
            align-items: center;
            justify-content: center;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .btn-icon:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
            transform: scale(1.05);
        }
        
        .session-item {
            padding: 16px;
            margin-bottom: 8px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
            overflow: hidden;
        }
        
        .session-item::before {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.05), transparent);
            transform: translateX(-100%);
            transition: transform 0.6s;
        }
        
        .session-item:hover {
            background: var(--bg-hover);
            border-color: var(--border-hover);
            transform: translateX(4px);
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.1);
        }
        
        .session-item:hover::before {
            transform: translateX(100%);
        }
        
        .session-item.active {
            background: var(--bg-active);
            border-color: var(--primary);
            box-shadow: 0 0 0 1px var(--primary), 0 4px 20px rgba(99, 102, 241, 0.2);
        }
        
        .session-header {
            display: flex;
            align-items: start;
            justify-content: space-between;
            margin-bottom: 8px;
        }
        
        .session-title {
            font-weight: 600;
            font-size: 14px;
            color: var(--text-primary);
            line-height: 1.4;
            flex: 1;
            margin-right: 8px;
        }
        
        .session-mode {
            padding: 2px 8px;
            background: var(--bg-hover);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            font-size: 11px;
            font-weight: 500;
            color: var(--text-secondary);
            text-transform: capitalize;
        }
        
        .session-preview {
            font-size: 13px;
            color: var(--text-secondary);
            line-height: 1.5;
            overflow: hidden;
            display: -webkit-box;
            -webkit-line-clamp: 2;
            -webkit-box-orient: vertical;
            margin-bottom: 8px;
        }
        
        .session-meta {
            display: flex;
            align-items: center;
            gap: 12px;
            font-size: 11px;
            color: var(--text-muted);
        }
        
        .session-meta-item {
            display: flex;
            align-items: center;
            gap: 4px;
        }
        
        .session-meta-item i {
            font-size: 10px;
        }
        
        .sidebar-footer {
            padding: 20px;
            border-top: 1px solid var(--border-color);
            background: rgba(17, 17, 17, 0.5);
        }
        
        .user-profile {
            display: flex;
            align-items: center;
            gap: 12px;
            padding: 12px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            margin-bottom: 16px;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .user-profile:hover {
            background: var(--bg-hover);
            border-color: var(--border-hover);
        }
        
        .user-avatar {
            width: 40px;
            height: 40px;
            border-radius: var(--radius-md);
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 18px;
            color: white;
            font-weight: 600;
        }
        
        .user-info {
            flex: 1;
        }
        
        .user-name {
            font-size: 14px;
            font-weight: 600;
            color: var(--text-primary);
        }
        
        .user-status {
            font-size: 12px;
            color: var(--text-secondary);
        }
        
        .btn-new-chat {
            width: 100%;
            padding: 14px;
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            color: white;
            border: none;
            border-radius: var(--radius-md);
            font-size: 14px;
            font-weight: 600;
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
            overflow: hidden;
            box-shadow: 0 4px 20px rgba(99, 102, 241, 0.3);
        }
        
        .btn-new-chat::before {
            content: '';
            position: absolute;
            top: 50%;
            left: 50%;
            width: 0;
            height: 0;
            background: rgba(255, 255, 255, 0.3);
            border-radius: 50%;
            transform: translate(-50%, -50%);
            transition: width 0.4s, height 0.4s;
        }
        
        .btn-new-chat:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 30px rgba(99, 102, 241, 0.4);
        }
        
        .btn-new-chat:hover::before {
            width: 300px;
            height: 300px;
        }
        
        .btn-new-chat:active {
            transform: translateY(0);
        }
        
        .main-content {
            flex: 1;
            display: flex;
            flex-direction: column;
            background: rgba(10, 10, 10, 0.5);
            position: relative;
            overflow: hidden;
        }
        
        .chat-header {
            height: var(--header-height);
            padding: 0 24px;
            background: rgba(17, 17, 17, 0.95);
            backdrop-filter: blur(20px);
            -webkit-backdrop-filter: blur(20px);
            border-bottom: 1px solid var(--border-color);
            display: flex;
            align-items: center;
            justify-content: space-between;
            position: relative;
            z-index: var(--z-sticky);
            box-shadow: 0 4px 24px rgba(0, 0, 0, 0.1);
        }
        
        .header-left {
            display: flex;
            align-items: center;
            gap: 20px;
        }
        
        .menu-toggle {
            width: 40px;
            height: 40px;
            border-radius: var(--radius-md);
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            color: var(--text-primary);
            display: flex;
            align-items: center;
            justify-content: center;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .menu-toggle:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
            transform: scale(1.05);
        }
        
        .mode-selector {
            display: flex;
            gap: 8px;
            padding: 4px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
        }
        
        .mode-btn {
            padding: 8px 16px;
            background: transparent;
            border: none;
            border-radius: var(--radius-sm);
            color: var(--text-secondary);
            font-size: 13px;
            font-weight: 500;
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 6px;
            position: relative;
            overflow: hidden;
        }
        
        .mode-btn::before {
            content: '';
            position: absolute;
            inset: 0;
            background: var(--primary);
            transform: scaleX(0);
            transform-origin: left;
            transition: transform var(--transition-fast);
            z-index: -1;
        }
        
        .mode-btn:hover {
            color: var(--text-primary);
        }
        
        .mode-btn.active {
            color: white;
        }
        
        .mode-btn.active::before {
            transform: scaleX(1);
        }
        
        .header-center {
            flex: 1;
            display: flex;
            justify-content: center;
        }
        
        .session-info {
            display: flex;
            align-items: center;
            gap: 16px;
            padding: 8px 16px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
        }
        
        .session-info-item {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 0 12px;
            border-right: 1px solid var(--border-color);
        }
        
        .session-info-item:last-child {
            border-right: none;
        }
        
        .session-info-label {
            font-size: 11px;
            color: var(--text-muted);
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .session-info-value {
            font-size: 14px;
            font-weight: 600;
            color: var(--text-primary);
        }
        
        .header-right {
            display: flex;
            align-items: center;
            gap: 12px;
        }
        
        .search-bar {
            position: relative;
            width: 240px;
        }
        
        .search-input {
            width: 100%;
            padding: 8px 16px 8px 40px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            color: var(--text-primary);
            font-size: 14px;
            transition: all var(--transition-fast);
        }
        
        .search-input:focus {
            outline: none;
            border-color: var(--primary);
            box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.1);
        }
        
        .search-icon {
            position: absolute;
            left: 14px;
            top: 50%;
            transform: translateY(-50%);
            color: var(--text-muted);
            font-size: 14px;
            pointer-events: none;
        }
        
        .header-actions {
            display: flex;
            gap: 8px;
        }
        
        .header-btn {
            width: 40px;
            height: 40px;
            border-radius: var(--radius-md);
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            color: var(--text-secondary);
            display: flex;
            align-items: center;
            justify-content: center;
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
        }
        
        .header-btn:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
            transform: scale(1.05);
        }
        
        .header-btn.has-indicator::after {
            content: '';
            position: absolute;
            top: 8px;
            right: 8px;
            width: 8px;
            height: 8px;
            background: var(--danger);
            border-radius: 50%;
            border: 2px solid var(--bg-secondary);
        }
        
        .advanced-toggle {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
            overflow: hidden;
        }
        
        .advanced-toggle::before {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(135deg, var(--danger), var(--warning));
            opacity: 0;
            transition: opacity var(--transition-fast);
        }
        
        .advanced-toggle.active::before {
            opacity: 1;
        }
        
        .advanced-toggle.active {
            border-color: transparent;
            animation: glow 2s ease-in-out infinite;
        }
        
        @keyframes glow {
            0%, 100% { box-shadow: 0 0 20px rgba(239, 68, 68, 0.5); }
            50% { box-shadow: 0 0 30px rgba(239, 68, 68, 0.8), 0 0 40px rgba(245, 158, 11, 0.6); }
        }
        
        .advanced-toggle-content {
            position: relative;
            z-index: 1;
            display: flex;
            align-items: center;
            gap: 8px;
        }
        
        .advanced-icon {
            font-size: 16px;
            transition: all var(--transition-fast);
        }
        
        .advanced-toggle.active .advanced-icon {
            color: white;
            animation: rotate 2s linear infinite;
        }
        
        @keyframes rotate {
            from { transform: rotate(0deg); }
            to { transform: rotate(360deg); }
        }
        
        .advanced-toggle-label {
            font-size: 13px;
            font-weight: 600;
            color: var(--text-primary);
            transition: color var(--transition-fast);
        }
        
        .advanced-toggle.active .advanced-toggle-label {
            color: white;
        }
        
        .toggle-switch {
            width: 40px;
            height: 20px;
            background: var(--bg-tertiary);
            border-radius: 20px;
            position: relative;
            transition: all var(--transition-fast);
        }
        
        .advanced-toggle.active .toggle-switch {
            background: rgba(255, 255, 255, 0.3);
        }
        
        .toggle-switch::after {
            content: '';
            position: absolute;
            top: 2px;
            left: 2px;
            width: 16px;
            height: 16px;
            background: var(--text-secondary);
            border-radius: 50%;
            transition: all var(--transition-fast);
            box-shadow: 0 2px 4px rgba(0, 0, 0, 0.2);
        }
        
        .advanced-toggle.active .toggle-switch::after {
            background: white;
            transform: translateX(20px);
            box-shadow: 0 2px 8px rgba(0, 0, 0, 0.3);
        }
        
        .messages-container {
            flex: 1;
            display: flex;
            flex-direction: column;
            position: relative;
            overflow: hidden;
        }
        
        .messages-area {
            flex: 1;
            overflow-y: auto;
            padding: 24px;
            scroll-behavior: smooth;
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) transparent;
        }
        
        .welcome-screen {
            display: flex;
            flex-direction: column;
            align-items: center;
            justify-content: center;
            height: 100%;
            padding: 48px;
            text-align: center;
            animation: fadeIn 0.6s ease-out;
        }
        
        @keyframes fadeIn {
            from {
                opacity: 0;
                transform: translateY(20px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
        
        .welcome-icon {
            width: 120px;
            height: 120px;
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            border-radius: var(--radius-2xl);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 60px;
            color: white;
            margin-bottom: 32px;
            box-shadow: 0 20px 60px rgba(99, 102, 241, 0.3);
            animation: float 6s ease-in-out infinite;
        }
        
        @keyframes float {
            0%, 100% { transform: translateY(0) rotate(0deg); }
            25% { transform: translateY(-10px) rotate(-5deg); }
            75% { transform: translateY(-10px) rotate(5deg); }
        }
        
        .welcome-title {
            font-size: 48px;
            font-weight: 800;
            color: var(--text-primary);
            margin-bottom: 16px;
            font-family: 'Space Grotesk', sans-serif;
            letter-spacing: -1px;
        }
        
        .welcome-subtitle {
            font-size: 20px;
            color: var(--text-secondary);
            margin-bottom: 48px;
            line-height: 1.6;
            max-width: 600px;
        }
        
        .feature-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 20px;
            width: 100%;
            max-width: 800px;
            margin-bottom: 48px;
        }
        
        .feature-card {
            padding: 24px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            transition: all var(--transition-fast);
            cursor: pointer;
            position: relative;
            overflow: hidden;
        }
        
        .feature-card::before {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            opacity: 0;
            transition: opacity var(--transition-fast);
        }
        
        .feature-card:hover {
            border-color: var(--primary);
            transform: translateY(-4px);
            box-shadow: 0 12px 40px rgba(99, 102, 241, 0.2);
        }
        
        .feature-card:hover::before {
            opacity: 0.1;
        }
        
        .feature-icon {
            width: 48px;
            height: 48px;
            background: var(--bg-hover);
            border-radius: var(--radius-md);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 24px;
            color: var(--primary);
            margin-bottom: 16px;
        }
        
        .feature-title {
            font-size: 16px;
            font-weight: 600;
            color: var(--text-primary);
            margin-bottom: 8px;
        }
        
        .feature-description {
            font-size: 13px;
            color: var(--text-secondary);
            line-height: 1.5;
        }
        
        .quick-actions {
            display: flex;
            gap: 16px;
            flex-wrap: wrap;
            justify-content: center;
        }
        
        .quick-action {
            padding: 12px 24px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            color: var(--text-secondary);
            font-size: 14px;
            font-weight: 500;
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 8px;
        }
        
        .quick-action:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
            transform: scale(1.05);
        }
        
        .message {
            margin-bottom: 24px;
            display: flex;
            align-items: flex-start;
            gap: 16px;
            animation: messageSlide 0.3s ease-out;
            position: relative;
        }
        
        @keyframes messageSlide {
            from {
                opacity: 0;
                transform: translateY(20px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
        
        .message.user {
            flex-direction: row-reverse;
        }
        
        .message-avatar {
            width: 44px;
            height: 44px;
            border-radius: var(--radius-md);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 20px;
            flex-shrink: 0;
            position: relative;
            overflow: hidden;
            box-shadow: 0 4px 20px rgba(0, 0, 0, 0.1);
        }
        
        .message.user .message-avatar {
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            color: white;
        }
        
        .message.assistant .message-avatar {
            background: linear-gradient(135deg, var(--secondary), var(--accent));
            color: white;
        }
        
        .message-avatar::after {
            content: '';
            position: absolute;
            inset: 0;
            background: radial-gradient(circle, rgba(255, 255, 255, 0.2), transparent);
            animation: ripple 3s ease-in-out infinite;
        }
        
        @keyframes ripple {
            0%, 100% { transform: scale(0.8); opacity: 1; }
            50% { transform: scale(1.2); opacity: 0; }
        }
        
        .message-content-wrapper {
            max-width: 70%;
            display: flex;
            flex-direction: column;
            gap: 8px;
        }
        
        .message-header {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 0 4px;
        }
        
        .message-author {
            font-weight: 600;
            font-size: 13px;
            color: var(--text-primary);
        }
        
        .message-badges {
            display: flex;
            gap: 6px;
        }
        
        .message-badge {
            padding: 2px 6px;
            background: var(--bg-hover);
            border-radius: var(--radius-sm);
            font-size: 10px;
            font-weight: 500;
            color: var(--text-secondary);
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .message-time {
            font-size: 11px;
            color: var(--text-muted);
            margin-left: auto;
        }
        
        .message-bubble {
            padding: 16px 20px;
            border-radius: var(--radius-lg);
            line-height: 1.6;
            font-size: 15px;
            position: relative;
            word-wrap: break-word;
            box-shadow: 0 2px 10px rgba(0, 0, 0, 0.05);
        }
        
        .message.user .message-bubble {
            background: linear-gradient(135deg, var(--primary), var(--primary-dark));
            color: white;
            border-bottom-right-radius: 4px;
        }
        
        .message.assistant .message-bubble {
            background: var(--bg-secondary);
            border: 1px solid var(--border-color);
            color: var(--text-primary);
            border-bottom-left-radius: 4px;
        }
        
        .message-bubble h1, .message-bubble h2, .message-bubble h3 {
            margin-top: 20px;
            margin-bottom: 12px;
            font-weight: 600;
            color: var(--text-primary);
        }
        
        .message.user .message-bubble h1,
        .message.user .message-bubble h2,
        .message.user .message-bubble h3 {
            color: white;
        }
        
        .message-bubble h1 { font-size: 24px; }
        .message-bubble h2 { font-size: 20px; }
        .message-bubble h3 { font-size: 18px; }
        
        .message-bubble p {
            margin-bottom: 12px;
        }
        
        .message-bubble ul, .message-bubble ol {
            margin: 12px 0;
            padding-left: 24px;
        }
        
        .message-bubble li {
            margin-bottom: 6px;
        }
        
        .message-bubble pre {
            background: var(--bg-primary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            padding: 16px;
            margin: 12px 0;
            overflow-x: auto;
            font-family: 'JetBrains Mono', monospace;
            font-size: 13px;
            line-height: 1.5;
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) transparent;
            box-shadow: inset 0 2px 4px rgba(0, 0, 0, 0.1);
        }
        
        .message.user .message-bubble pre {
            background: rgba(0, 0, 0, 0.2);
            border-color: rgba(255, 255, 255, 0.1);
        }
        
        .message-bubble code {
            background: rgba(139, 92, 246, 0.1);
            padding: 2px 6px;
            border-radius: 4px;
            font-family: 'JetBrains Mono', monospace;
            font-size: 13px;
            color: var(--primary-light);
        }
        
        .message.user .message-bubble code {
            background: rgba(255, 255, 255, 0.2);
            color: white;
        }
        
        .message-bubble blockquote {
            margin: 16px 0;
            padding: 12px 20px;
            border-left: 4px solid var(--primary);
            background: var(--bg-hover);
            border-radius: var(--radius-sm);
            font-style: italic;
            color: var(--text-secondary);
        }
        
        .message.user .message-bubble blockquote {
            border-left-color: rgba(255, 255, 255, 0.5);
            background: rgba(255, 255, 255, 0.1);
            color: rgba(255, 255, 255, 0.9);
        }
        
        .message-bubble table {
            width: 100%;
            margin: 16px 0;
            border-collapse: collapse;
            border-radius: var(--radius-md);
            overflow: hidden;
            box-shadow: 0 2px 10px rgba(0, 0, 0, 0.05);
        }
        
        .message-bubble th,
        .message-bubble td {
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid var(--border-color);
        }
        
        .message-bubble th {
            background: var(--bg-hover);
            font-weight: 600;
            color: var(--text-primary);
        }
        
        .message.user .message-bubble th {
            background: rgba(255, 255, 255, 0.1);
            color: white;
        }
        
        .message.user .message-bubble td {
            border-bottom-color: rgba(255, 255, 255, 0.1);
        }
        
        .message-bubble tr:last-child td {
            border-bottom: none;
        }
        
        .message-bubble a {
            color: var(--primary-light);
            text-decoration: none;
            border-bottom: 1px solid transparent;
            transition: border-color var(--transition-fast);
        }
        
        .message-bubble a:hover {
            border-bottom-color: var(--primary-light);
        }
        
        .message.user .message-bubble a {
            color: white;
            border-bottom-color: rgba(255, 255, 255, 0.3);
        }
        
        .message.user .message-bubble a:hover {
            border-bottom-color: white;
        }
        
        .message-attachments {
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 12px;
        }
        
        .attachment-item {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 12px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            font-size: 13px;
            color: var(--text-secondary);
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .attachment-item:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
        }
        
        .attachment-icon {
            font-size: 16px;
        }
        
        .attachment-name {
            max-width: 150px;
            overflow: hidden;
            text-overflow: ellipsis;
            white-space: nowrap;
        }
        
        .attachment-size {
            font-size: 11px;
            color: var(--text-muted);
        }
        
        .streaming-indicator {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            padding: 4px 8px;
            background: var(--bg-hover);
            border-radius: var(--radius-sm);
            font-size: 11px;
            color: var(--text-secondary);
            animation: pulse 2s ease-in-out infinite;
        }
        
        .streaming-dots {
            display: flex;
            gap: 3px;
        }
        
        .streaming-dot {
            width: 4px;
            height: 4px;
            background: var(--primary);
            border-radius: 50%;
            animation: streamingDot 1.4s ease-in-out infinite;
        }
        
        .streaming-dot:nth-child(2) {
            animation-delay: 0.2s;
        }
        
        .streaming-dot:nth-child(3) {
            animation-delay: 0.4s;
        }
        
        @keyframes streamingDot {
            0%, 80%, 100% {
                transform: scale(0.8);
                opacity: 0.5;
            }
            40% {
                transform: scale(1.2);
                opacity: 1;
            }
        }
        
        .message-actions {
            display: flex;
            gap: 8px;
            padding: 0 4px;
            opacity: 0;
            transition: opacity var(--transition-fast);
        }
        
        .message:hover .message-actions {
            opacity: 1;
        }
        
        .message-action {
            padding: 6px 12px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            font-size: 12px;
            color: var(--text-secondary);
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 4px;
        }
        
        .message-action:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            color: var(--primary);
            transform: scale(1.05);
        }
        
        .message-reactions {
            display: flex;
            gap: 6px;
            margin-top: 8px;
            flex-wrap: wrap;
        }
        
        .reaction {
            display: flex;
            align-items: center;
            gap: 4px;
            padding: 4px 8px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            font-size: 12px;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .reaction:hover {
            background: var(--bg-hover);
            border-color: var(--primary);
            transform: scale(1.05);
        }
        
        .reaction.active {
            background: var(--bg-active);
            border-color: var(--primary);
        }
        
        .reaction-emoji {
            font-size: 16px;
        }
        
        .reaction-count {
            color: var(--text-secondary);
            font-weight: 500;
        }
        
        .typing-indicator {
            display: none;
            align-items: center;
            gap: 16px;
            padding: 0 24px;
            margin-bottom: 24px;
            animation: fadeIn 0.3s;
        }
        
        .typing-indicator.active {
            display: flex;
        }
        
        .typing-avatar {
            width: 44px;
            height: 44px;
            border-radius: var(--radius-md);
            background: linear-gradient(135deg, var(--secondary), var(--accent));
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 20px;
            color: white;
            flex-shrink: 0;
        }
        
        .typing-dots {
            display: flex;
            gap: 6px;
            padding: 16px 20px;
            background: var(--bg-secondary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            border-bottom-left-radius: 4px;
        }
        
        .typing-dot {
            width: 8px;
            height: 8px;
            background: var(--primary);
            border-radius: 50%;
            animation: typing 1.4s infinite;
        }
        
        .typing-dot:nth-child(2) {
            animation-delay: 0.2s;
        }
        
        .typing-dot:nth-child(3) {
            animation-delay: 0.4s;
        }
        
        @keyframes typing {
            0%, 60%, 100% {
                transform: translateY(0);
                opacity: 0.7;
            }
            30% {
                transform: translateY(-10px);
                opacity: 1;
            }
        }
        
        .input-section {
            padding: 24px;
            background: rgba(17, 17, 17, 0.95);
            backdrop-filter: blur(20px);
            -webkit-backdrop-filter: blur(20px);
            border-top: 1px solid var(--border-color);
            position: relative;
            z-index: var(--z-sticky);
            box-shadow: 0 -4px 24px rgba(0, 0, 0, 0.1);
        }
        
        .input-features {
            display: flex;
            gap: 12px;
            margin-bottom: 16px;
            flex-wrap: wrap;
            align-items: center;
        }
        
        .feature-group {
            display: flex;
            gap: 8px;
            padding: 4px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
        }
        
        .feature-btn {
            padding: 8px 12px;
            background: transparent;
            border: none;
            border-radius: var(--radius-sm);
            font-size: 13px;
            color: var(--text-secondary);
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 6px;
            position: relative;
        }
        
        .feature-btn:hover {
            background: var(--bg-hover);
            color: var(--text-primary);
        }
        
        .feature-btn.active {
            background: var(--primary);
            color: white;
        }
        
        .feature-btn-tooltip {
            position: absolute;
            bottom: 100%;
            left: 50%;
            transform: translateX(-50%);
            padding: 6px 12px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            font-size: 11px;
            color: var(--text-primary);
            white-space: nowrap;
            opacity: 0;
            pointer-events: none;
            transition: all var(--transition-fast);
            margin-bottom: 8px;
        }
        
        .feature-btn:hover .feature-btn-tooltip {
            opacity: 1;
        }
        
        .input-divider {
            width: 1px;
            height: 32px;
            background: var(--border-color);
            margin: 0 8px;
        }
        
        .token-display {
            margin-left: auto;
            display: flex;
            align-items: center;
            gap: 12px;
            padding: 8px 16px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
        }
        
        .token-info {
            display: flex;
            flex-direction: column;
            align-items: flex-end;
            gap: 2px;
        }
        
        .token-label {
            font-size: 11px;
            color: var(--text-muted);
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .token-value {
            font-size: 14px;
            font-weight: 600;
            color: var(--text-primary);
        }
        
        .token-bar {
            width: 120px;
            height: 6px;
            background: var(--bg-tertiary);
            border-radius: 3px;
            overflow: hidden;
            position: relative;
        }
        
        .token-fill {
            height: 100%;
            background: linear-gradient(90deg, var(--primary), var(--secondary));
            transition: width 0.5s ease;
            position: relative;
            overflow: hidden;
        }
        
        .token-fill::after {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.3), transparent);
            animation: shimmer-move 2s infinite;
        }
        
        @keyframes shimmer-move {
            0% { transform: translateX(-100%); }
            100% { transform: translateX(100%); }
        }
        
        .file-upload-area {
            display: none;
            margin-bottom: 16px;
            padding: 32px;
            background: var(--bg-card);
            border: 2px dashed var(--border-color);
            border-radius: var(--radius-lg);
            text-align: center;
            transition: all var(--transition-fast);
            position: relative;
            overflow: hidden;
        }
        
        .file-upload-area.active {
            display: block;
            animation: slideDown 0.3s ease;
        }
        
        @keyframes slideDown {
            from {
                opacity: 0;
                transform: translateY(-20px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
        
        .file-upload-area.dragover {
            background: rgba(99, 102, 241, 0.05);
            border-color: var(--primary);
        }
        
        .file-upload-icon {
            font-size: 48px;
            color: var(--text-muted);
            margin-bottom: 16px;
        }
        
        .file-upload-text {
            font-size: 16px;
            color: var(--text-secondary);
            margin-bottom: 8px;
            font-weight: 500;
        }
        
        .file-upload-hint {
            font-size: 13px;
            color: var(--text-muted);
        }
        
        .file-list {
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 24px;
        }
        
        .file-item {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            background: var(--primary);
            color: white;
            border-radius: var(--radius-md);
            font-size: 13px;
            animation: fadeIn 0.3s;
            position: relative;
            overflow: hidden;
        }
        
        .file-item::before {
            content: '';
            position: absolute;
            inset: 0;
            background: rgba(255, 255, 255, 0.1);
            transform: scaleX(var(--progress, 0));
            transform-origin: left;
            transition: transform 0.3s;
        }
        
        .file-item-content {
            position: relative;
            z-index: 1;
            display: flex;
            align-items: center;
            gap: 8px;
        }
        
        .file-item .remove-file {
            cursor: pointer;
            opacity: 0.8;
            transition: opacity var(--transition-fast);
            margin-left: 4px;
        }
        
        .file-item .remove-file:hover {
            opacity: 1;
        }
        
        .input-container {
            display: flex;
            gap: 12px;
            align-items: flex-end;
            position: relative;
        }
        
        .input-wrapper {
            flex: 1;
            position: relative;
        }
        
        .input-box {
            width: 100%;
            min-height: var(--input-height);
            max-height: 200px;
            padding: 16px 20px;
            padding-right: 140px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            color: var(--text-primary);
            font-size: 15px;
            font-family: inherit;
            resize: none;
            transition: all var(--transition-fast);
            line-height: 1.5;
            scrollbar-width: thin;
            scrollbar-color: var(--border-color) transparent;
        }
        
        .input-box:focus {
            outline: none;
            border-color: var(--primary);
            background: var(--bg-hover);
            box-shadow: 0 0 0 3px rgba(99, 102, 241, 0.1);
        }
        
        .input-box::placeholder {
            color: var(--text-muted);
        }
        
        .input-actions {
            position: absolute;
            right: 12px;
            top: 50%;
            transform: translateY(-50%);
            display: flex;
            gap: 8px;
            align-items: center;
        }
        
        .input-action {
            width: 32px;
            height: 32px;
            border-radius: var(--radius-sm);
            background: var(--bg-hover);
            border: 1px solid var(--border-color);
            color: var(--text-secondary);
            display: flex;
            align-items: center;
            justify-content: center;
            cursor: pointer;
            transition: all var(--transition-fast);
            position: relative;
        }
        
        .input-action:hover {
            background: var(--primary);
            border-color: transparent;
            color: white;
            transform: scale(1.1);
        }
        
        .char-counter {
            font-size: 11px;
            color: var(--text-muted);
            margin-right: 8px;
        }
        
        .send-button {
            padding: 0 28px;
            height: var(--input-height);
            background: linear-gradient(135deg, var(--primary), var(--secondary));
            color: white;
            border: none;
            border-radius: var(--radius-lg);
            font-size: 15px;
            font-weight: 600;
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 10px;
            position: relative;
            overflow: hidden;
            box-shadow: 0 4px 20px rgba(99, 102, 241, 0.3);
        }
        
        .send-button::before {
            content: '';
            position: absolute;
            inset: 0;
            background: linear-gradient(135deg, rgba(255, 255, 255, 0.2), transparent);
            opacity: 0;
            transition: opacity var(--transition-fast);
        }
        
        .send-button:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 30px rgba(99, 102, 241, 0.4);
        }
        
        .send-button:hover::before {
            opacity: 1;
        }
        
        .send-button:active {
            transform: translateY(0);
        }
        
        .send-button:disabled {
            opacity: 0.5;
            cursor: not-allowed;
            transform: none;
        }
        
        .send-button.loading {
            color: transparent;
        }
        
        .send-button.loading::after {
            content: '';
            position: absolute;
            width: 20px;
            height: 20px;
            top: 50%;
            left: 50%;
            margin-left: -10px;
            margin-top: -10px;
            border: 2px solid white;
            border-radius: 50%;
            border-top-color: transparent;
            animation: spin 0.8s linear infinite;
        }
        
        @keyframes spin {
            to { transform: rotate(360deg); }
        }
        
        .context-menu {
            position: fixed;
            background: var(--bg-secondary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            padding: 8px;
            min-width: 200px;
            box-shadow: var(--shadow-lg);
            z-index: var(--z-dropdown);
            display: none;
        }
        
        .context-menu.active {
            display: block;
        }
        
        .context-menu-item {
            padding: 10px 16px;
            border-radius: var(--radius-sm);
            color: var(--text-secondary);
            cursor: pointer;
            transition: all var(--transition-fast);
            display: flex;
            align-items: center;
            gap: 12px;
            font-size: 14px;
        }
        
        .context-menu-item:hover {
            background: var(--bg-hover);
            color: var(--text-primary);
        }
        
        .context-menu-item.danger {
            color: var(--danger);
        }
        
        .context-menu-divider {
            height: 1px;
            background: var(--border-color);
            margin: 8px 0;
        }
        
        .modal-overlay {
            position: fixed;
            inset: 0;
            background: rgba(0, 0, 0, 0.8);
            backdrop-filter: blur(10px);
            -webkit-backdrop-filter: blur(10px);
            display: none;
            align-items: center;
            justify-content: center;
            z-index: var(--z-modal-backdrop);
            padding: 20px;
            opacity: 0;
            transition: opacity var(--transition-base);
        }
        
        .modal-overlay.active {
            display: flex;
            opacity: 1;
        }
        
        .modal-content {
            background: var(--bg-secondary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-xl);
            padding: 32px;
            max-width: 600px;
            width: 100%;
            max-height: 90vh;
            overflow-y: auto;
            box-shadow: var(--shadow-xl);
            transform: scale(0.9);
            transition: transform var(--transition-spring);
            position: relative;
        }
        
        .modal-overlay.active .modal-content {
            transform: scale(1);
        }
        
        .modal-header {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 24px;
        }
        
        .modal-title {
            font-size: 24px;
            font-weight: 700;
            color: var(--text-primary);
        }
        
        .modal-close {
            width: 36px;
            height: 36px;
            border-radius: var(--radius-md);
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            color: var(--text-secondary);
            display: flex;
            align-items: center;
            justify-content: center;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .modal-close:hover {
            background: var(--bg-hover);
            border-color: var(--danger);
            color: var(--danger);
            transform: scale(1.1);
        }
        
        .modal-body {
            color: var(--text-secondary);
            line-height: 1.6;
        }
        
        .settings-section {
            margin-bottom: 32px;
        }
        
        .settings-title {
            font-size: 16px;
            font-weight: 600;
            color: var(--text-primary);
            margin-bottom: 16px;
        }
        
        .setting-item {
            display: flex;
            align-items: center;
            justify-content: space-between;
            padding: 16px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            margin-bottom: 12px;
        }
        
        .setting-info {
            flex: 1;
        }
        
        .setting-label {
            font-size: 14px;
            font-weight: 500;
            color: var(--text-primary);
            margin-bottom: 4px;
        }
        
        .setting-description {
            font-size: 13px;
            color: var(--text-secondary);
        }
        
        .setting-control {
            display: flex;
            align-items: center;
            gap: 12px;
        }
        
        .switch {
            width: 48px;
            height: 24px;
            background: var(--bg-tertiary);
            border-radius: 24px;
            position: relative;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .switch.active {
            background: var(--primary);
        }
        
        .switch::after {
            content: '';
            position: absolute;
            top: 3px;
            left: 3px;
            width: 18px;
            height: 18px;
            background: white;
            border-radius: 50%;
            transition: all var(--transition-fast);
            box-shadow: 0 2px 4px rgba(0, 0, 0, 0.2);
        }
        
        .switch.active::after {
            transform: translateX(24px);
        }
        
        .notification-container {
            position: fixed;
            top: 20px;
            right: 20px;
            z-index: var(--z-notification);
            pointer-events: none;
        }
        
        .notification {
            padding: 16px 20px;
            border-radius: var(--radius-md);
            color: white;
            font-size: 14px;
            font-weight: 500;
            box-shadow: var(--shadow-lg);
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 12px;
            pointer-events: auto;
            animation: slideInRight 0.3s ease-out;
            position: relative;
            overflow: hidden;
            min-width: 300px;
        }
        
        @keyframes slideInRight {
            from {
                opacity: 0;
                transform: translateX(100%);
            }
            to {
                opacity: 1;
                transform: translateX(0);
            }
        }
        
        .notification::before {
            content: '';
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            height: 3px;
            background: rgba(255, 255, 255, 0.3);
            animation: progress 5s linear forwards;
        }
        
        @keyframes progress {
            from { transform: scaleX(1); }
            to { transform: scaleX(0); }
        }
        
        .notification.success {
            background: linear-gradient(135deg, var(--success), #22c55e);
        }
        
        .notification.error {
            background: linear-gradient(135deg, var(--danger), #f87171);
        }
        
        .notification.warning {
            background: linear-gradient(135deg, var(--warning), #fbbf24);
        }
        
        .notification.info {
            background: linear-gradient(135deg, var(--info), #60a5fa);
        }
        
        .notification-icon {
            font-size: 20px;
        }
        
        .notification-content {
            flex: 1;
        }
        
        .notification-title {
            font-weight: 600;
            margin-bottom: 2px;
        }
        
        .notification-message {
            font-size: 13px;
            opacity: 0.9;
        }
        
        .notification-close {
            cursor: pointer;
            opacity: 0.8;
            transition: opacity var(--transition-fast);
            font-size: 16px;
        }
        
        .notification-close:hover {
            opacity: 1;
        }
        
        .loading-overlay {
            position: fixed;
            inset: 0;
            background: rgba(0, 0, 0, 0.9);
            backdrop-filter: blur(10px);
            -webkit-backdrop-filter: blur(10px);
            display: none;
            align-items: center;
            justify-content: center;
            z-index: var(--z-top);
        }
        
        .loading-overlay.active {
            display: flex;
        }
        
        .loading-content {
            text-align: center;
        }
        
        .loading-spinner {
            width: 80px;
            height: 80px;
            position: relative;
            margin: 0 auto 24px;
        }
        
        .loading-spinner::before,
        .loading-spinner::after {
            content: '';
            position: absolute;
            inset: 0;
            border-radius: 50%;
            border: 3px solid transparent;
        }
        
        .loading-spinner::before {
            border-top-color: var(--primary);
            animation: spin 1s linear infinite;
        }
        
        .loading-spinner::after {
            border-bottom-color: var(--secondary);
            animation: spin 1s linear infinite reverse;
            inset: 10px;
        }
        
        .loading-text {
            font-size: 18px;
            color: var(--text-primary);
            font-weight: 500;
        }
        
        .tooltip {
            position: absolute;
            padding: 8px 12px;
            background: var(--bg-tertiary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            font-size: 12px;
            color: var(--text-primary);
            white-space: nowrap;
            pointer-events: none;
            z-index: var(--z-tooltip);
            opacity: 0;
            transition: all var(--transition-fast);
        }
        
        .tooltip.active {
            opacity: 1;
        }
        
        .emoji-picker {
            position: absolute;
            bottom: 100%;
            right: 0;
            width: 320px;
            height: 400px;
            background: var(--bg-secondary);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-lg);
            box-shadow: var(--shadow-lg);
            display: none;
            flex-direction: column;
            margin-bottom: 8px;
        }
        
        .emoji-picker.active {
            display: flex;
        }
        
        .emoji-search {
            padding: 12px;
            border-bottom: 1px solid var(--border-color);
        }
        
        .emoji-search input {
            width: 100%;
            padding: 8px 12px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-md);
            color: var(--text-primary);
            font-size: 14px;
        }
        
        .emoji-categories {
            display: flex;
            gap: 4px;
            padding: 8px;
            border-bottom: 1px solid var(--border-color);
        }
        
        .emoji-category {
            flex: 1;
            padding: 8px;
            background: var(--bg-card);
            border: 1px solid var(--border-color);
            border-radius: var(--radius-sm);
            text-align: center;
            cursor: pointer;
            transition: all var(--transition-fast);
        }
        
        .emoji-category:hover,
        .emoji-category.active {
            background: var(--bg-hover);
            border-color: var(--primary);
        }
        
        .emoji-grid {
            flex: 1;
            overflow-y: auto;
            padding: 8px;
            display: grid;
            grid-template-columns: repeat(8, 1fr);
            gap: 4px;
        }
        
        .emoji-item {
            width: 32px;
            height: 32px;
            display: flex;
            align-items: center;
            justify-content: center;
            border-radius: var(--radius-sm);
            cursor: pointer;
            transition: all var(--transition-fast);
            font-size: 20px;
        }
        
        .emoji-item:hover {
            background: var(--bg-hover);
            transform: scale(1.2);
        }
        
        @media (max-width: 1024px) {
            .sidebar {
                position: fixed;
                z-index: var(--z-modal);
                height: 100vh;
                box-shadow: 4px 0 24px rgba(0, 0, 0, 0.3);
            }
            
            .header-center {
                display: none;
            }
            
            .search-bar {
                width: 180px;
            }
        }
        
        @media (max-width: 768px) {
            :root {
                --sidebar-width: 300px;
            }
            
            .mode-selector {
                display: none;
            }
            
            .header-right {
                gap: 8px;
            }
            
            .advanced-toggle {
                padding: 8px;
            }
            
            .advanced-toggle-label {
                display: none;
            }
            
            .message-content-wrapper {
                max-width: 85%;
            }
            
            .input-features {
                gap: 8px;
            }
            
            .token-display {
                display: none;
            }
            
            .feature-group {
                flex-wrap: wrap;
            }
            
            .search-bar {
                display: none;
            }
        }
        
        @media (max-width: 480px) {
            .messages-area {
                padding: 16px;
            }
            
            .input-section {
                padding: 16px;
            }
            
            .input-box {
                padding-right: 100px;
                font-size: 16px;
            }
            
            .send-button {
                padding: 0 20px;
            }
            
            .send-button span {
                display: none;
            }
            
            .welcome-title {
                font-size: 32px;
            }
            
            .welcome-subtitle {
                font-size: 16px;
            }
            
            .feature-grid {
                grid-template-columns: 1fr;
            }
            
            .modal-content {
                padding: 24px;
            }
        }
        
        @media (prefers-reduced-motion: reduce) {
            * {
                animation-duration: 0.01ms !important;
                animation-iteration-count: 1 !important;
                transition-duration: 0.01ms !important;
            }
        }
        
        @media (prefers-color-scheme: light) {
            :root {
                --bg-primary: #ffffff;
                --bg-secondary: #f9fafb;
                --bg-tertiary: #f3f4f6;
                --bg-card: rgba(0, 0, 0, 0.02);
                --bg-hover: rgba(0, 0, 0, 0.05);
                
                --text-primary: #111827;
                --text-secondary: #6b7280;
                --text-muted: #9ca3af;
                
                --border-color: rgba(0, 0, 0, 0.1);
                --border-hover: rgba(0, 0, 0, 0.2);
            }
        }
        
        .sr-only {
            position: absolute;
            width: 1px;
            height: 1px;
            padding: 0;
            margin: -1px;
            overflow: hidden;
            clip: rect(0, 0, 0, 0);
            white-space: nowrap;
            border-width: 0;
        }
        
        .no-scroll {
            overflow: hidden;
        }
        
        .fade-enter {
            opacity: 0;
        }
        
        .fade-enter-active {
            opacity: 1;
            transition: opacity var(--transition-base);
        }
        
        .fade-exit {
            opacity: 1;
        }
        
        .fade-exit-active {
            opacity: 0;
            transition: opacity var(--transition-base);
        }
        
        .slide-enter {
            transform: translateX(-100%);
        }
        
        .slide-enter-active {
            transform: translateX(0);
            transition: transform var(--transition-base);
        }
        
        .slide-exit {
            transform: translateX(0);
        }
        
        .slide-exit-active {
            transform: translateX(-100%);
            transition: transform var(--transition-base);
        }
    </style>
</head>
<body>
    <canvas class="background-canvas" id="backgroundCanvas"></canvas>
    
    <div class="particle-field" id="particleField"></div>
    
    <div class="gradient-mesh"></div>
    
    <div class="grid-overlay"></div>
    
    <div class="loading-overlay" id="loadingOverlay">
        <div class="loading-content">
            <div class="loading-spinner"></div>
            <div class="loading-text">Initializing AI Ultra...</div>
        </div>
    </div>
    
    <div class="notification-container" id="notificationContainer"></div>
    
    <div class="app-container">
        <aside class="sidebar" id="sidebar">
            <div class="sidebar-header">
                <div class="logo-container">
                    <div class="logo-icon">
                        <i class="fas fa-brain"></i>
                    </div>
                    <div class="logo-text">
                        <h1>Jack's AI Ultra</h1>
                        <p>Ultimate Edition</p>
                    </div>
                </div>
            </div>
            
            <nav class="sidebar-nav">
                <div class="nav-section">
                    <div class="nav-section-title">Navigation</div>
                    <a href="#" class="nav-item active" data-page="chat">
                        <i class="fas fa-comments"></i>
                        <span>Chat</span>
                        <span class="nav-badge">3</span>
                    </a>
                    <a href="#" class="nav-item" data-page="explore">
                        <i class="fas fa-compass"></i>
                        <span>Explore</span>
                    </a>
                    <a href="#" class="nav-item" data-page="library">
                        <i class="fas fa-book"></i>
                        <span>Library</span>
                    </a>
                    <a href="#" class="nav-item" data-page="analytics">
                        <i class="fas fa-chart-bar"></i>
                        <span>Analytics</span>
                    </a>
                </div>
                
                <div class="nav-section">
                    <div class="nav-section-title">Tools</div>
                    <a href="#" class="nav-item" data-page="code">
                        <i class="fas fa-code"></i>
                        <span>Code Editor</span>
                    </a>
                    <a href="#" class="nav-item" data-page="canvas">
                        <i class="fas fa-paint-brush"></i>
                        <span>Canvas</span>
                    </a>
                    <a href="#" class="nav-item" data-page="data">
                        <i class="fas fa-database"></i>
                        <span>Data Studio</span>
                    </a>
                </div>
            </nav>
            
            <div class="chat-sessions" id="chatSessions">
                <div class="sessions-header">
                    <h3 class="sessions-title">Recent Chats</h3>
                    <div class="sessions-actions">
                        <button class="btn-icon" id="btnSearchSessions">
                            <i class="fas fa-search"></i>
                        </button>
                        <button class="btn-icon" id="btnFilterSessions">
                            <i class="fas fa-filter"></i>
                        </button>
                    </div>
                </div>
            </div>
            
            <div class="sidebar-footer">
                <div class="user-profile" id="userProfile">
                    <div class="user-avatar">U</div>
                    <div class="user-info">
                        <div class="user-name">User</div>
                        <div class="user-status">Premium</div>
                    </div>
                    <i class="fas fa-chevron-right"></i>
                </div>
                
                <button class="btn-new-chat" id="btnNewChat">
                    <i class="fas fa-plus"></i> New Chat
                </button>
            </div>
        </aside>
        
        <main class="main-content">
            <header class="chat-header">
                <div class="header-left">
                    <button class="menu-toggle" id="menuToggle">
                        <i class="fas fa-bars"></i>
                    </button>
                    
                    <div class="mode-selector" id="modeSelector">
                        <button class="mode-btn active" data-mode="balanced">
                            <i class="fas fa-balance-scale"></i>
                            <span>Balanced</span>
                        </button>
                        <button class="mode-btn" data-mode="creative">
                            <i class="fas fa-palette"></i>
                            <span>Creative</span>
                        </button>
                        <button class="mode-btn" data-mode="precise">
                            <i class="fas fa-microscope"></i>
                            <span>Precise</span>
                        </button>
                        <button class="mode-btn" data-mode="code">
                            <i class="fas fa-code"></i>
                            <span>Code</span>
                        </button>
                        <button class="mode-btn" data-mode="research">
                            <i class="fas fa-book"></i>
                            <span>Research</span>
                        </button>
                    </div>
                </div>
                
                <div class="header-center">
                    <div class="session-info">
                        <div class="session-info-item">
                            <div>
                                <div class="session-info-label">Messages</div>
                                <div class="session-info-value" id="messageCount">0</div>
                            </div>
                        </div>
                        <div class="session-info-item">
                            <div>
                                <div class="session-info-label">Tokens</div>
                                <div class="session-info-value" id="tokenCount">0</div>
                            </div>
                        </div>
                        <div class="session-info-item">
                            <div>
                                <div class="session-info-label">Model</div>
                                <div class="session-info-value">Gemini 2.0</div>
                            </div>
                        </div>
                    </div>
                </div>
                
                <div class="header-right">
                    <div class="search-bar">
                        <i class="fas fa-search search-icon"></i>
                        <input type="text" class="search-input" placeholder="Search messages..." id="searchInput">
                    </div>
                    
                    <div class="header-actions">
                        <button class="header-btn" id="btnShare">
                            <i class="fas fa-share-alt"></i>
                        </button>
                        <button class="header-btn" id="btnExport">
                            <i class="fas fa-download"></i>
                        </button>
                        <button class="header-btn has-indicator" id="btnNotifications">
                            <i class="fas fa-bell"></i>
                        </button>
                        <button class="header-btn" id="btnSettings">
                            <i class="fas fa-cog"></i>
                        </button>
                    </div>
                    
                    <div class="advanced-toggle" id="advancedToggle">
                        <div class="advanced-toggle-content">
                            <i class="fas fa-bolt advanced-icon"></i>
                            <span class="advanced-toggle-label">Advanced</span>
                            <div class="toggle-switch"></div>
                        </div>
                    </div>
                </div>
            </header>
            
            <div class="messages-container">
                <div class="messages-area" id="messagesArea">
                    <div class="welcome-screen" id="welcomeScreen">
                        <div class="welcome-icon">
                            <i class="fas fa-rocket"></i>
                        </div>
                        <h1 class="welcome-title">Welcome to AI Ultra</h1>
                        <p class="welcome-subtitle">
                            Experience the pinnacle of AI assistance with unparalleled capabilities, 
                            extended context, and cutting-edge intelligence.
                        </p>
                        
                        <div class="feature-grid">
                            <div class="feature-card">
                                <div class="feature-icon">
                                    <i class="fas fa-brain"></i>
                                </div>
                                <h3 class="feature-title">Advanced Intelligence</h3>
                                <p class="feature-description">
                                    Powered by the latest AI models with enhanced reasoning
                                </p>
                            </div>
                            <div class="feature-card">
                                <div class="feature-icon">
                                    <i class="fas fa-infinity"></i>
                                </div>
                                <h3 class="feature-title">Extended Context</h3>
                                <p class="feature-description">
                                    Handle massive documents and long conversations effortlessly
                                </p>
                            </div>
                            <div class="feature-card">
                                <div class="feature-icon">
                                    <i class="fas fa-magic"></i>
                                </div>
                                <h3 class="feature-title">Multi-Modal</h3>
                                <p class="feature-description">
                                    Process text, images, code, data, and more
                                </p>
                            </div>
                            <div class="feature-card">
                                <div class="feature-icon">
                                    <i class="fas fa-bolt"></i>
                                </div>
                                <h3 class="feature-title">Real-Time</h3>
                                <p class="feature-description">
                                    Stream responses with minimal latency
                                </p>
                            </div>
                        </div>
                        
                        <div class="quick-actions">
                            <button class="quick-action" data-prompt="Explain quantum computing">
                                <i class="fas fa-atom"></i>
                                <span>Explain quantum computing</span>
                            </button>
                            <button class="quick-action" data-prompt="Write a Python web scraper">
                                <i class="fas fa-spider"></i>
                                <span>Write a Python web scraper</span>
                            </button>
                            <button class="quick-action" data-prompt="Analyze market trends">
                                <i class="fas fa-chart-line"></i>
                                <span>Analyze market trends</span>
                            </button>
                            <button class="quick-action" data-prompt="Create a business plan">
                                <i class="fas fa-briefcase"></i>
                                <span>Create a business plan</span>
                            </button>
                        </div>
                    </div>
                </div>
                
                <div class="typing-indicator" id="typingIndicator">
                    <div class="typing-avatar">
                        <i class="fas fa-robot"></i>
                    </div>
                    <div class="typing-dots">
                        <span class="typing-dot"></span>
                        <span class="typing-dot"></span>
                        <span class="typing-dot"></span>
                    </div>
                </div>
            </div>
            
            <div class="input-section">
                <div class="input-features">
                    <div class="feature-group">
                        <button class="feature-btn" id="btnFiles" data-feature="files">
                            <i class="fas fa-paperclip"></i>
                            <span>Files</span>
                            <div class="feature-btn-tooltip">Attach files up to 2GB</div>
                        </button>
                        <button class="feature-btn" id="btnImage" data-feature="image">
                            <i class="fas fa-image"></i>
                            <span>Image</span>
                            <div class="feature-btn-tooltip">Add images</div>
                        </button>
                        <button class="feature-btn" id="btnVoice" data-feature="voice">
                            <i class="fas fa-microphone"></i>
                            <span>Voice</span>
                            <div class="feature-btn-tooltip">Voice input</div>
                        </button>
                    </div>
                    
                    <div class="input-divider"></div>
                    
                    <div class="feature-group">
                        <button class="feature-btn" id="btnTemplates" data-feature="templates">
                            <i class="fas fa-file-alt"></i>
                            <span>Templates</span>
                        </button>
                        <button class="feature-btn" id="btnSnippets" data-feature="snippets">
                            <i class="fas fa-puzzle-piece"></i>
                            <span>Snippets</span>
                        </button>
                        <button class="feature-btn" id="btnContext" data-feature="context">
                            <i class="fas fa-layer-group"></i>
                            <span>Context</span>
                        </button>
                    </div>
                    
                    <div class="token-display">
                        <div class="token-info">
                            <div class="token-label">Tokens</div>
                            <div class="token-value" id="tokenValue">0 / 2M</div>
                        </div>
                        <div class="token-bar">
                            <div class="token-fill" id="tokenFill" style="width: 0%"></div>
                        </div>
                    </div>
                </div>
                
                <div class="file-upload-area" id="fileUploadArea">
                    <i class="fas fa-cloud-upload-alt file-upload-icon"></i>
                    <p class="file-upload-text">Drop files here or click to browse</p>
                    <p class="file-upload-hint">Supports all file types up to 2GB</p>
                    <div class="file-list" id="fileList"></div>
                </div>
                
                <div class="input-container">
                    <div class="input-wrapper">
                        <textarea 
                            class="input-box" 
                            id="messageInput" 
                            placeholder="Ask anything... (Shift+Enter for new line)"
                            rows="1"
                        ></textarea>
                        <div class="input-actions">
                            <span class="char-counter" id="charCounter">0</span>
                            <button class="input-action" id="btnClear" title="Clear input">
                                <i class="fas fa-times"></i>
                            </button>
                            <button class="input-action" id="btnEmoji" title="Add emoji">
                                <i class="fas fa-smile"></i>
                            </button>
                            <button class="input-action" id="btnFormat" title="Format text">
                                <i class="fas fa-font"></i>
                            </button>
                            
                            <div class="emoji-picker" id="emojiPicker">
                                <div class="emoji-search">
                                    <input type="text" placeholder="Search emojis...">
                                </div>
                                <div class="emoji-categories">
                                    <div class="emoji-category active" data-category="people">😀</div>
                                    <div class="emoji-category" data-category="nature">🌿</div>
                                    <div class="emoji-category" data-category="food">🍔</div>
                                    <div class="emoji-category" data-category="activity">⚽</div>
                                    <div class="emoji-category" data-category="travel">✈️</div>
                                    <div class="emoji-category" data-category="objects">💡</div>
                                    <div class="emoji-category" data-category="symbols">❤️</div>
                                    <div class="emoji-category" data-category="flags">🏁</div>
                                </div>
                                <div class="emoji-grid" id="emojiGrid"></div>
                            </div>
                        </div>
                    </div>
                    <button class="send-button" id="sendButton">
                        <i class="fas fa-paper-plane"></i>
                        <span>Send</span>
                    </button>
                </div>
            </div>
        </main>
    </div>
    
    <div class="context-menu" id="contextMenu">
        <div class="context-menu-item" data-action="copy">
            <i class="fas fa-copy"></i>
            <span>Copy</span>
        </div>
        <div class="context-menu-item" data-action="edit">
            <i class="fas fa-edit"></i>
            <span>Edit</span>
        </div>
        <div class="context-menu-item" data-action="regenerate">
            <i class="fas fa-redo"></i>
            <span>Regenerate</span>
        </div>
        <div class="context-menu-divider"></div>
        <div class="context-menu-item" data-action="pin">
            <i class="fas fa-thumbtack"></i>
            <span>Pin</span>
        </div>
        <div class="context-menu-item" data-action="branch">
            <i class="fas fa-code-branch"></i>
            <span>Create branch</span>
        </div>
        <div class="context-menu-divider"></div>
        <div class="context-menu-item danger" data-action="delete">
            <i class="fas fa-trash"></i>
            <span>Delete</span>
        </div>
    </div>
    
    <div class="modal-overlay" id="modalOverlay">
        <div class="modal-content" id="modalContent">
            <div class="modal-header">
                <h2 class="modal-title" id="modalTitle">Settings</h2>
                <button class="modal-close" id="modalClose">
                    <i class="fas fa-times"></i>
                </button>
            </div>
            <div class="modal-body" id="modalBody">
                <div class="settings-section">
                    <h3 class="settings-title">Appearance</h3>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Dark Mode</div>
                            <div class="setting-description">Use dark theme for better visibility</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch active" data-setting="darkMode"></div>
                        </div>
                    </div>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Animations</div>
                            <div class="setting-description">Enable smooth animations and transitions</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch active" data-setting="animations"></div>
                        </div>
                    </div>
                </div>
                
                <div class="settings-section">
                    <h3 class="settings-title">Chat</h3>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Streaming</div>
                            <div class="setting-description">Stream responses in real-time</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch active" data-setting="streaming"></div>
                        </div>
                    </div>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Auto-save</div>
                            <div class="setting-description">Automatically save conversations</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch active" data-setting="autoSave"></div>
                        </div>
                    </div>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Notifications</div>
                            <div class="setting-description">Show notifications for important events</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch active" data-setting="notifications"></div>
                        </div>
                    </div>
                </div>
                
                <div class="settings-section">
                    <h3 class="settings-title">Advanced</h3>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Developer Mode</div>
                            <div class="setting-description">Show advanced options and debugging info</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch" data-setting="developerMode"></div>
                        </div>
                    </div>
                    <div class="setting-item">
                        <div class="setting-info">
                            <div class="setting-label">Experimental Features</div>
                            <div class="setting-description">Enable cutting-edge experimental features</div>
                        </div>
                        <div class="setting-control">
                            <div class="switch" data-setting="experimental"></div>
                        </div>
                    </div>
                </div>
            </div>
        </div>
    </div>
    
    <div class="tooltip" id="tooltip"></div>
    
    <input type="file" id="fileInput" multiple style="display: none;">
    
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/prism.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-python.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-javascript.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-java.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-cpp.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-csharp.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-go.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-rust.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-typescript.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-jsx.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-tsx.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-sql.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-bash.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-yaml.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-json.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/components/prism-markdown.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/marked/11.1.1/marked.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/dompurify/3.0.6/purify.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
    
    <script>
        const AppState = {
            sessionId: null,
            currentMode: 'balanced',
            messages: [],
            tokenUsage: 0,
            maxTokens: 2000000,
            attachedFiles: [],
            isTyping: false,
            isStreaming: false,
            currentStreamController: null,
            advancedMode: false,
            voiceRecognition: null,
            activeStream: null,
            settings: {
                darkMode: true,
                animations: true,
                streaming: true,
                notifications: true,
                soundEffects: false,
                autoSave: true,
                developerMode: false,
                experimental: false
            },
            ui: {
                sidebarCollapsed: false,
                currentPage: 'chat',
                searchQuery: '',
                selectedMessages: new Set(),
                contextMenuTarget: null,
                emojiPickerOpen: false,
                modalOpen: false
            },
            cache: {
                sessions: new Map(),
                messages: new Map(),
                templates: new Map()
            },
            performance: {
                lastResponseTime: 0,
                averageResponseTime: 0,
                totalMessages: 0
            },
            features: {
                voiceEnabled: false,
                filesEnabled: false,
                templatesEnabled: false,
                contextEnabled: false
            }
        };
        
        const API_ENDPOINTS = {
            base: '/api',
            session: {
                create: '/api/session/create',
                get: '/api/session/:id',
                update: '/api/session/update',
                delete: '/api/session/delete',
                list: '/api/sessions',
                search: '/api/sessions/search'
            },
            chat: '/api/chat',
            export: '/api/export/:format',
            models: '/api/models',
            stats: '/api/stats',
            upload: '/api/upload',
            settings: '/api/settings',
            share: '/api/share'
        };
        
        const EMOJIS = {
            people: ['😀', '😃', '😄', '😁', '😆', '😅', '🤣', '😂', '🙂', '🙃', '😉', '😊', '😇', '🥰', '😍', '🤩', '😘', '😗', '☺️', '😚', '😙', '🥲', '😋', '😛', '😜', '🤪', '😝', '🤑', '🤗', '🤭', '🤫', '🤔', '🤐', '🤨', '😐', '😑', '😶', '😏', '😒', '🙄', '😬', '🤥', '😌', '😔', '😪', '🤤', '😴', '😷', '🤒', '🤕', '🤢', '🤮', '🤧', '🥵', '🥶', '🥴', '😵', '🤯', '🤠', '🥳', '🥸', '😎', '🤓', '🧐', '😕', '😟', '🙁', '☹️', '😮', '😯', '😲', '😳', '🥺', '😦', '😧', '😨', '😰', '😥', '😢', '😭', '😱', '😖', '😣', '😞', '😓', '😩', '😫', '🥱', '😤', '😡', '😠', '🤬', '😈', '👿', '💀', '☠️', '💩', '🤡', '👹', '👺', '👻', '👽', '👾', '🤖'],
            nature: ['🌿', '🌵', '🌲', '🌳', '🌴', '🌱', '🌾', '🌺', '🌻', '🌹', '🌷', '🌼', '🌸', '💐', '🍄', '🌰', '🦀', '🦞', '🦐', '🦑', '🌍', '🌎', '🌏', '🌕', '🌖', '🌗', '🌘', '🌑', '🌒', '🌓', '🌔', '🌚', '🌝', '🌛', '🌜', '☀️', '🌤', '⛅', '🌥', '🌦', '🌧', '⛈', '🌩', '🌨', '❄️', '☃️', '⛄', '🌬', '💨', '🌪', '🌫', '🌊', '💧', '💦', '☔'],
            food: ['🍔', '🍕', '🌭', '🥪', '🌮', '🌯', '🥙', '🧆', '🍖', '🍗', '🥓', '🍟', '🍿', '🧈', '🧂', '🥚', '🍳', '🥞', '🧇', '🥐', '🍞', '🥖', '🥨', '🧀', '🥗', '🥣', '🍲', '🍝', '🍜', '🍛', '🍚', '🍙', '🍘', '🥟', '🍤', '🍣', '🍥', '🍡', '🦪', '🍦', '🍧', '🍨', '🍩', '🍪', '🎂', '🍰', '🧁', '🥧', '🍫', '🍬', '🍭', '🍯', '🍼', '🥛', '☕', '🍵', '🍶', '🍾', '🍷', '🍸', '🍹', '🍺', '🍻', '🥂', '🥃', '🥤', '🧋', '🧃', '🧉'],
            activity: ['⚽', '🏀', '🏈', '⚾', '🥎', '🎾', '🏐', '🏉', '🥏', '🎱', '🪀', '🏓', '🏸', '🏒', '🏑', '🥍', '🏏', '🪃', '🥅', '⛳', '🪁', '🏹', '🎣', '🤿', '🥊', '🥋', '🎽', '🛹', '🛼', '🛷', '⛸', '🥌', '🎿', '⛷', '🏂', '🪂', '🏋️', '🤸', '🤾', '🏌️', '🏇', '🧘', '🏄', '🏊', '🤽', '🚣', '🧗', '🚴', '🚵', '🎪', '🎭', '🎨', '🎬', '🎤', '🎧', '🎼', '🎹', '🥁', '🪘', '🎷', '🎺', '🪗', '🎸', '🪕', '🎻', '🎲', '♟', '🎯', '🎳', '🎮', '🎰', '🧩'],
            travel: ['✈️', '🚀', '🛸', '🚁', '🛶', '⛵', '🚤', '🛥', '🛳', '⛴', '🚢', '⚓', '🚂', '🚆', '🚇', '🚊', '🚝', '🚞', '🚋', '🚃', '🚎', '🚌', '🚍', '🚙', '🚘', '🚗', '🚕', '🚖', '🚛', '🚚', '🚒', '🚑', '🚓', '🏎', '🚜', '🛵', '🏍', '🛺', '🚲', '🛴', '🚏', '🛣', '🛤', '⛽', '🚨', '🚥', '🚦', '🚧', '🏰', '🏯', '🏟', '🎡', '🎢', '🎠', '⛲', '⛱', '🏖', '🏝', '🏜', '🌋', '⛰', '🏔', '🗻', '🏕', '⛺', '🏠', '🏡', '🏘', '🏚', '🏗', '🏭', '🏢', '🏬', '🏣', '🏤', '🏥', '🏦', '🏨', '🏪', '🏫', '🏩', '💒', '🏛', '⛪', '🕌', '🕍', '🛕', '🕋', '⛩', '🛤', '🛣', '🗾', '🎑', '🏞', '🌅', '🌄', '🌠', '🎇', '🎆', '🌇', '🌆', '🏙', '🌃', '🌌', '🌉', '🌁'],
            objects: ['💡', '🔦', '🏮', '🪔', '📱', '💻', '⌨️', '🖥', '🖨', '🖱', '🖲', '💽', '💾', '💿', '📀', '📼', '📷', '📸', '📹', '🎥', '📽', '🎞', '☎️', '📞', '📟', '📠', '📺', '📻', '🎙', '🎚', '🎛', '🧭', '⏱', '⏲', '⏰', '🕰', '⌛', '⏳', '📡', '🔋', '🔌', '💡', '🔦', '🕯', '🪔', '🧯', '🛢', '💸', '💵', '💴', '💶', '💷', '🪙', '💰', '💳', '💎', '⚖️', '🪜', '🧰', '🪛', '🔧', '🔨', '⚒', '🛠', '⛏', '🪚', '🔩', '⚙️', '🪤', '🧱', '⛓', '🧲', '🔫', '💣', '🧨', '🪓', '🔪', '🗡', '⚔️', '🛡', '🚬', '⚰️', '🪦', '⚱️', '🏺', '🔮', '📿', '🧿', '💈', '⚗️', '🔭', '🔬', '🕳', '🩹', '🩺', '💊', '💉', '🩸', '🧬', '🦠', '🧫', '🧪', '🌡', '🧹', '🪠', '🧺', '🧻', '🚽', '🚰', '🚿', '🛁', '🛀', '🧼', '🪥', '🪒', '🧽', '🪣', '🧴', '🛎', '🔑', '🗝', '🚪', '🪑', '🛋', '🛏', '🛌', '🧸', '🪆', '🖼', '🪞', '🪟', '🛍', '🛒', '🎁', '🎈', '🎏', '🎀', '🪄', '🪅', '🎊', '🎉', '🎎', '🏮', '🎐', '🧧', '✉️', '📩', '📨', '📧', '💌', '📥', '📤', '📦', '🏷', '🪧', '📪', '📫', '📬', '📭', '📮', '📯', '📜', '📃', '📄', '📑', '🧾', '📊', '📈', '📉', '🗒', '🗓', '📆', '📅', '🗑', '📇', '🗃', '🗳', '🗄', '📋', '📁', '📂', '🗂', '🗞', '📰', '📓', '📔', '📒', '📕', '📗', '📘', '📙', '📚', '📖', '🔖', '🧷', '🔗', '📎', '🖇', '📐', '📏', '🧮', '📌', '📍', '✂️', '🖊', '🖋', '✒️', '🖌', '🖍', '📝', '✏️', '🔍', '🔎', '🔏', '🔐', '🔒', '🔓'],
            symbols: ['❤️', '🧡', '💛', '💚', '💙', '💜', '🖤', '🤍', '🤎', '💔', '❣️', '💕', '💞', '💓', '💗', '💖', '💘', '💝', '💟', '☮️', '✝️', '☪️', '🕉', '☸️', '✡️', '🔯', '🕎', '☯️', '☦️', '🛐', '⛎', '♈', '♉', '♊', '♋', '♌', '♍', '♎', '♏', '♐', '♑', '♒', '♓', '🆔', '⚛️', '🉑', '☢️', '☣️', '📴', '📳', '🈶', '🈚', '🈸', '🈺', '🈷️', '✴️', '🆚', '💮', '🉐', '㊙️', '㊗️', '🈴', '🈵', '🈹', '🈲', '🅰️', '🅱️', '🆎', '🆑', '🅾️', '🆘', '❌', '⭕', '🛑', '⛔', '📛', '🚫', '💯', '💢', '♨️', '🚷', '🚯', '🚳', '🚱', '🔞', '📵', '🚭', '❗', '❕', '❓', '❔', '‼️', '⁉️', '🔅', '🔆', '〽️', '⚠️', '🚸', '🔱', '⚜️', '🔰', '♻️', '✅', '🈯', '💹', '❇️', '✳️', '❎', '🌐', '💠', 'Ⓜ️', '🌀', '💤', '🏧', '🚾', '♿', '🅿️', '🛗', '🈳', '🈂️', '🛂', '🛃', '🛄', '🛅', '🚹', '🚺', '🚼', '⚧', '🚻', '🚮', '🎦', '📶', '🈁', '🔣', 'ℹ️', '🔤', '🔡', '🔠', '🆖', '🆗', '🆙', '🆒', '🆕', '🆓', '0️⃣', '1️⃣', '2️⃣', '3️⃣', '4️⃣', '5️⃣', '6️⃣', '7️⃣', '8️⃣', '9️⃣', '🔟', '🔢', '#️⃣', '*️⃣', '⏏️', '▶️', '⏸', '⏯', '⏹', '⏺', '⏭', '⏮', '⏩', '⏪', '⏫', '⏬', '◀️', '🔼', '🔽', '➡️', '⬅️', '⬆️', '⬇️', '↗️', '↘️', '↙️', '↖️', '↕️', '↔️', '↪️', '↩️', '⤴️', '⤵️', '🔀', '🔁', '🔂', '🔄', '🔃', '🎵', '🎶', '➕', '➖', '➗', '✖️', '♾', '💲', '💱', '™️', '©️', '®️', '〰️', '➰', '➿', '🔚', '🔙', '🔛', '🔝', '🔜', '✔️', '☑️', '🔘', '🔴', '🟠', '🟡', '🟢', '🔵', '🟣', '⚫', '⚪', '🟤', '🔺', '🔻', '🔸', '🔹', '🔶', '🔷', '🔳', '🔲', '▪️', '▫️', '◾', '◽', '◼️', '◻️', '🟥', '🟧', '🟨', '🟩', '🟦', '🟪', '⬛', '⬜', '🟫', '🔈', '🔇', '🔉', '🔊', '🔔', '🔕', '📣', '📢', '👁‍🗨', '💬', '💭', '🗯', '♠️', '♣️', '♥️', '♦️', '🃏', '🎴', '🀄', '🕐', '🕑', '🕒', '🕓', '🕔', '🕕', '🕖', '🕗', '🕘', '🕙', '🕚', '🕛', '🕜', '🕝', '🕞', '🕟', '🕠', '🕡', '🕢', '🕣', '🕤', '🕥', '🕦', '🕧'],
            flags: ['🏁', '🚩', '🎌', '🏴', '🏳️', '🏳️‍🌈', '🏳️‍⚧️', '🏴‍☠️', '🇦🇫', '🇦🇽', '🇦🇱', '🇩🇿', '🇦🇸', '🇦🇩', '🇦🇴', '🇦🇮', '🇦🇶', '🇦🇬', '🇦🇷', '🇦🇲', '🇦🇼', '🇦🇺', '🇦🇹', '🇦🇿', '🇧🇸', '🇧🇭', '🇧🇩', '🇧🇧', '🇧🇾', '🇧🇪', '🇧🇿', '🇧🇯', '🇧🇲', '🇧🇹', '🇧🇴', '🇧🇦', '🇧🇼', '🇧🇷', '🇮🇴', '🇻🇬', '🇧🇳', '🇧🇬', '🇧🇫', '🇧🇮', '🇰🇭', '🇨🇲', '🇨🇦', '🇮🇨', '🇨🇻', '🇧🇶', '🇰🇾', '🇨🇫', '🇹🇩', '🇨🇱', '🇨🇳', '🇨🇽', '🇨🇨', '🇨🇴', '🇰🇲', '🇨🇬', '🇨🇩', '🇨🇰', '🇨🇷', '🇨🇮', '🇭🇷', '🇨🇺', '🇨🇼', '🇨🇾', '🇨🇿', '🇩🇰', '🇩🇯', '🇩🇲', '🇩🇴', '🇪🇨', '🇪🇬', '🇸🇻', '🇬🇶', '🇪🇷', '🇪🇪', '🇪🇹', '🇪🇺', '🇫🇰', '🇫🇴', '🇫🇯', '🇫🇮', '🇫🇷', '🇬🇫', '🇵🇫', '🇹🇫', '🇬🇦', '🇬🇲', '🇬🇪', '🇩🇪', '🇬🇭', '🇬🇮', '🇬🇷', '🇬🇱', '🇬🇩', '🇬🇵', '🇬🇺', '🇬🇹', '🇬🇬', '🇬🇳', '🇬🇼', '🇬🇾', '🇭🇹', '🇭🇳', '🇭🇰', '🇭🇺', '🇮🇸', '🇮🇳', '🇮🇩', '🇮🇷', '🇮🇶', '🇮🇪', '🇮🇲', '🇮🇱', '🇮🇹', '🇯🇲', '🇯🇵', '🎌', '🇯🇪', '🇯🇴', '🇰🇿', '🇰🇪', '🇰🇮', '🇽🇰', '🇰🇼', '🇰🇬', '🇱🇦', '🇱🇻', '🇱🇧', '🇱🇸', '🇱🇷', '🇱🇾', '🇱🇮', '🇱🇹', '🇱🇺', '🇲🇴', '🇲🇰', '🇲🇬', '🇲🇼', '🇲🇾', '🇲🇻', '🇲🇱', '🇲🇹', '🇲🇭', '🇲🇶', '🇲🇷', '🇲🇺', '🇾🇹', '🇲🇽', '🇫🇲', '🇲🇩', '🇲🇨', '🇲🇳', '🇲🇪', '🇲🇸', '🇲🇦', '🇲🇿', '🇲🇲', '🇳🇦', '🇳🇷', '🇳🇵', '🇳🇱', '🇳🇨', '🇳🇿', '🇳🇮', '🇳🇪', '🇳🇬', '🇳🇺', '🇳🇫', '🇰🇵', '🇲🇵', '🇳🇴', '🇴🇲', '🇵🇰', '🇵🇼', '🇵🇸', '🇵🇦', '🇵🇬', '🇵🇾', '🇵🇪', '🇵🇭', '🇵🇳', '🇵🇱', '🇵🇹', '🇵🇷', '🇶🇦', '🇷🇪', '🇷🇴', '🇷🇺', '🇷🇼', '🇼🇸', '🇸🇲', '🇸🇦', '🇸🇳', '🇷🇸', '🇸🇨', '🇸🇱', '🇸🇬', '🇸🇽', '🇸🇰', '🇸🇮', '🇬🇸', '🇸🇧', '🇸🇴', '🇿🇦', '🇰🇷', '🇸🇸', '🇪🇸', '🇱🇰', '🇧🇱', '🇸🇭', '🇰🇳', '🇱🇨', '🇵🇲', '🇻🇨', '🇸🇩', '🇸🇷', '🇸🇿', '🇸🇪', '🇨🇭', '🇸🇾', '🇹🇼', '🇹🇯', '🇹🇿', '🇹🇭', '🇹🇱', '🇹🇬', '🇹🇰', '🇹🇴', '🇹🇹', '🇹🇳', '🇹🇷', '🇹🇲', '🇹🇨', '🇹🇻', '🇻🇮', '🇺🇬', '🇺🇦', '🇦🇪', '🇬🇧', '🏴󐁧󐁢󐁥󐁮󐁧󐁿', '🏴󐁧󐁢󐁳󐁣󐁴󐁿', '🏴󐁧󐁢󐁷󐁬󐁳󐁿', '🇺🇳', '🇺🇸', '🇺🇾', '🇺🇿', '🇻🇺', '🇻🇦', '🇻🇪', '🇻🇳', '🇼🇫', '🇪🇭', '🇾🇪', '🇿🇲', '🇿🇼']
        };
        
        class AIUltraApp {
            constructor() {
                this.state = AppState;
                this.api = new APIClient();
                this.ui = new UIManager();
                this.chat = new ChatManager();
                this.session = new SessionManager();
                this.voice = new VoiceManager();
                this.file = new FileManager();
                this.settings = new SettingsManager();
                this.performance = new PerformanceMonitor();
                this.analytics = new AnalyticsTracker();
                this.particles = new ParticleSystem();
                this.shortcuts = new ShortcutManager();
                this.init();
            }
            
            async init() {
                try {
                    this.showLoading(true);
                    
                    await this.initializeBackground();
                    await this.loadSettings();
                    await this.initializeSession();
                    await this.loadModels();
                    
                    this.setupEventListeners();
                    this.setupWebSocket();
                    this.setupServiceWorker();
                    
                    this.ui.initialize();
                    this.shortcuts.initialize();
                    this.analytics.track('app_initialized');
                    
                    this.showLoading(false);
                    this.showNotification('AI Ultra initialized successfully', 'success');
                } catch (error) {
                    console.error('Initialization error:', error);
                    this.showLoading(false);
                    this.showNotification('Failed to initialize application', 'error');
                }
            }
            
            async initializeBackground() {
                this.particles.create();
                this.setupCanvasAnimation();
            }
            
            setupCanvasAnimation() {
                const canvas = document.getElementById('backgroundCanvas');
                const ctx = canvas.getContext('2d');
                
                const resize = () => {
                    canvas.width = window.innerWidth;
                    canvas.height = window.innerHeight;
                };
                
                resize();
                window.addEventListener('resize', resize);
                
                const nodes = [];
                const nodeCount = 50;
                const connectionDistance = 150;
                
                for (let i = 0; i < nodeCount; i++) {
                    nodes.push({
                        x: Math.random() * canvas.width,
                        y: Math.random() * canvas.height,
                        vx: (Math.random() - 0.5) * 0.5,
                        vy: (Math.random() - 0.5) * 0.5,
                        radius: Math.random() * 2 + 1
                    });
                }
                
                const animate = () => {
                    ctx.clearRect(0, 0, canvas.width, canvas.height);
                    
                    nodes.forEach((node, i) => {
                        node.x += node.vx;
                        node.y += node.vy;
                        
                        if (node.x < 0 || node.x > canvas.width) node.vx *= -1;
                        if (node.y < 0 || node.y > canvas.height) node.vy *= -1;
                        
                        ctx.beginPath();
                        ctx.arc(node.x, node.y, node.radius, 0, Math.PI * 2);
                        ctx.fillStyle = 'rgba(99, 102, 241, 0.5)';
                        ctx.fill();
                        
                        for (let j = i + 1; j < nodes.length; j++) {
                            const dx = nodes[j].x - node.x;
                            const dy = nodes[j].y - node.y;
                            const distance = Math.sqrt(dx * dx + dy * dy);
                            
                            if (distance < connectionDistance) {
                                ctx.beginPath();
                                ctx.moveTo(node.x, node.y);
                                ctx.lineTo(nodes[j].x, nodes[j].y);
                                ctx.strokeStyle = `rgba(99, 102, 241, ${0.2 * (1 - distance / connectionDistance)})`;
                                ctx.stroke();
                            }
                        }
                    });
                    
                    requestAnimationFrame(animate);
                };
                
                if (this.state.settings.animations) {
                    animate();
                }
            }
            
            async loadSettings() {
                const savedSettings = localStorage.getItem('ai_ultra_settings');
                if (savedSettings) {
                    this.state.settings = { ...this.state.settings, ...JSON.parse(savedSettings) };
                }
                
                this.applySettings();
            }
            
            applySettings() {
                document.documentElement.setAttribute('data-theme', this.state.settings.darkMode ? 'dark' : 'light');
                
                if (!this.state.settings.animations) {
                    document.documentElement.style.setProperty('--transition-fast', '0ms');
                    document.documentElement.style.setProperty('--transition-base', '0ms');
                    document.documentElement.style.setProperty('--transition-slow', '0ms');
                }
            }
            
            async initializeSession() {
                const sessionId = localStorage.getItem('ai_ultra_session_id');
                
                if (sessionId) {
                    try {
                        const session = await this.api.getSession(sessionId);
                        this.state.sessionId = sessionId;
                        this.session.loadSession(session);
                    } catch (error) {
                        await this.createNewSession();
                    }
                } else {
                    await this.createNewSession();
                }
            }
            
            async createNewSession() {
                const session = await this.api.createSession({
                    mode: this.state.currentMode
                });
                
                this.state.sessionId = session.sessionId;
                localStorage.setItem('ai_ultra_session_id', session.sessionId);
                this.session.loadSession(session);
            }
            
            async loadModels() {
                try {
                    const models = await this.api.getModels();
                    this.state.availableModels = models;
                } catch (error) {
                    console.error('Failed to load models:', error);
                }
            }
            
            setupEventListeners() {
                document.getElementById('menuToggle').addEventListener('click', () => this.ui.toggleSidebar());
                document.getElementById('btnNewChat').addEventListener('click', () => this.startNewChat());
                document.getElementById('sendButton').addEventListener('click', () => this.sendMessage());
                document.getElementById('messageInput').addEventListener('keydown', (e) => this.handleInputKeydown(e));
                document.getElementById('advancedToggle').addEventListener('click', () => this.toggleAdvancedMode());
                document.getElementById('btnSettings').addEventListener('click', () => this.ui.openSettings());
                document.getElementById('btnExport').addEventListener('click', () => this.exportChat());
                document.getElementById('btnShare').addEventListener('click', () => this.shareChat());
                document.getElementById('btnNotifications').addEventListener('click', () => this.ui.showNotificationCenter());
                
                document.querySelectorAll('.mode-btn').forEach(btn => {
                    btn.addEventListener('click', () => this.changeMode(btn.dataset.mode));
                });
                
                document.querySelectorAll('.feature-btn').forEach(btn => {
                    btn.addEventListener('click', () => this.toggleFeature(btn.dataset.feature));
                });
                
                document.querySelectorAll('.quick-action').forEach(btn => {
                    btn.addEventListener('click', () => {
                        const prompt = btn.dataset.prompt;
                        document.getElementById('messageInput').value = prompt;
                        this.sendMessage();
                    });
                });
                
                this.setupDragAndDrop();
                this.setupContextMenu();
                this.setupSearch();
                this.setupEmojis();
                this.setupModalHandlers();
                this.setupVoiceInput();
            }
            
            setupDragAndDrop() {
                const uploadArea = document.getElementById('fileUploadArea');
                const messagesArea = document.getElementById('messagesArea');
                
                ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {
                    uploadArea.addEventListener(eventName, preventDefaults, false);
                    messagesArea.addEventListener(eventName, preventDefaults, false);
                    document.body.addEventListener(eventName, preventDefaults, false);
                });
                
                function preventDefaults(e) {
                    e.preventDefault();
                    e.stopPropagation();
                }
                
                ['dragenter', 'dragover'].forEach(eventName => {
                    uploadArea.addEventListener(eventName, () => {
                        uploadArea.classList.add('dragover');
                    }, false);
                });
                
                ['dragleave', 'drop'].forEach(eventName => {
                    uploadArea.addEventListener(eventName, () => {
                        uploadArea.classList.remove('dragover');
                    }, false);
                });
                
                uploadArea.addEventListener('drop', (e) => {
                    const files = e.dataTransfer.files;
                    this.handleFiles(files);
                }, false);
                
                messagesArea.addEventListener('drop', (e) => {
                    const files = e.dataTransfer.files;
                    if (files.length > 0) {
                        this.showFileUploadArea();
                        this.handleFiles(files);
                    }
                }, false);
                
                uploadArea.addEventListener('click', () => {
                    document.getElementById('fileInput').click();
                });
                
                document.getElementById('fileInput').addEventListener('change', (e) => {
                    this.handleFiles(e.target.files);
                });
            }
            
            setupContextMenu() {
                const contextMenu = document.getElementById('contextMenu');
                
                document.addEventListener('contextmenu', (e) => {
                    const message = e.target.closest('.message');
                    if (message) {
                        e.preventDefault();
                        this.state.ui.contextMenuTarget = message;
                        this.ui.showContextMenu(e.pageX, e.pageY);
                    }
                });
                
                document.addEventListener('click', (e) => {
                    if (!e.target.closest('#contextMenu')) {
                        this.ui.hideContextMenu();
                    }
                });
                
                contextMenu.querySelectorAll('.context-menu-item').forEach(item => {
                    item.addEventListener('click', () => {
                        const action = item.dataset.action;
                        this.handleContextMenuAction(action);
                        this.ui.hideContextMenu();
                    });
                });
            }
            
            setupSearch() {
                const searchInput = document.getElementById('searchInput');
                let searchTimeout;
                
                searchInput.addEventListener('input', (e) => {
                    clearTimeout(searchTimeout);
                    searchTimeout = setTimeout(() => {
                        this.searchMessages(e.target.value);
                    }, 300);
                });
                
                searchInput.addEventListener('keydown', (e) => {
                    if (e.key === 'Escape') {
                        e.target.value = '';
                        this.searchMessages('');
                        e.target.blur();
                    }
                });
            }
            
            setupEmojis() {
                const emojiBtn = document.getElementById('btnEmoji');
                const emojiPicker = document.getElementById('emojiPicker');
                const emojiGrid = document.getElementById('emojiGrid');
                
                emojiBtn.addEventListener('click', () => {
                    this.state.ui.emojiPickerOpen = !this.state.ui.emojiPickerOpen;
                    emojiPicker.classList.toggle('active', this.state.ui.emojiPickerOpen);
                    
                    if (this.state.ui.emojiPickerOpen && emojiGrid.children.length === 0) {
                        this.loadEmojis('people');
                    }
                });
                
                document.querySelectorAll('.emoji-category').forEach(cat => {
                    cat.addEventListener('click', () => {
                        document.querySelectorAll('.emoji-category').forEach(c => c.classList.remove('active'));
                        cat.classList.add('active');
                        this.loadEmojis(cat.dataset.category);
                    });
                });
                
                document.addEventListener('click', (e) => {
                    if (!e.target.closest('#btnEmoji') && !e.target.closest('#emojiPicker')) {
                        this.state.ui.emojiPickerOpen = false;
                        emojiPicker.classList.remove('active');
                    }
                });
            }
            
            loadEmojis(category) {
                const emojiGrid = document.getElementById('emojiGrid');
                emojiGrid.innerHTML = '';
                
                const emojis = EMOJIS[category] || [];
                emojis.forEach(emoji => {
                    const emojiItem = document.createElement('div');
                    emojiItem.className = 'emoji-item';
                    emojiItem.textContent = emoji;
                    emojiItem.addEventListener('click', () => {
                        this.insertEmoji(emoji);
                    });
                    emojiGrid.appendChild(emojiItem);
                });
            }
            
            insertEmoji(emoji) {
                const input = document.getElementById('messageInput');
                const start = input.selectionStart;
                const end = input.selectionEnd;
                const text = input.value;
                
                input.value = text.substring(0, start) + emoji + text.substring(end);
                input.selectionStart = input.selectionEnd = start + emoji.length;
                input.focus();
                
                this.updateCharCounter();
            }
            
            setupModalHandlers() {
                const modalOverlay = document.getElementById('modalOverlay');
                const modalClose = document.getElementById('modalClose');
                
                modalClose.addEventListener('click', () => {
                    this.ui.closeModal();
                });
                
                modalOverlay.addEventListener('click', (e) => {
                    if (e.target === modalOverlay) {
                        this.ui.closeModal();
                    }
                });
                
                document.querySelectorAll('.switch').forEach(switchEl => {
                    switchEl.addEventListener('click', () => {
                        const setting = switchEl.dataset.setting;
                        this.toggleSetting(setting);
                    });
                });
            }
            
            setupVoiceInput() {
                if ('webkitSpeechRecognition' in window || 'SpeechRecognition' in window) {
                    const SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
                    this.state.voiceRecognition = new SpeechRecognition();
                    this.state.voiceRecognition.continuous = true;
                    this.state.voiceRecognition.interimResults = true;
                    this.state.voiceRecognition.lang = 'en-US';
                    
                    this.state.voiceRecognition.onresult = (event) => {
                        let transcript = '';
                        for (let i = event.resultIndex; i < event.results.length; i++) {
                            if (event.results[i].isFinal) {
                                transcript += event.results[i][0].transcript + ' ';
                            }
                        }
                        
                        if (transcript) {
                            const input = document.getElementById('messageInput');
                            input.value += transcript;
                            this.updateCharCounter();
                        }
                    };
                    
                    this.state.voiceRecognition.onerror = (event) => {
                        console.error('Voice recognition error:', event.error);
                        this.showNotification('Voice recognition error', 'error');
                        this.stopVoiceRecording();
                    };
                }
            }
            
            setupWebSocket() {
                
            }
            
            setupServiceWorker() {
                if ('serviceWorker' in navigator) {
                    
                }
            }
            
            toggleSidebar() {
                this.state.ui.sidebarCollapsed = !this.state.ui.sidebarCollapsed;
                document.getElementById('sidebar').classList.toggle('collapsed', this.state.ui.sidebarCollapsed);
            }
            
            async startNewChat() {
                if (this.state.messages.length > 0) {
                    const confirmed = await this.ui.showConfirmDialog(
                        'Start New Chat?',
                        'Your current conversation will be saved. Do you want to continue?'
                    );
                    
                    if (!confirmed) return;
                }
                
                await this.createNewSession();
                this.clearMessages();
                this.showWelcomeScreen();
                this.showNotification('New chat started', 'success');
                this.analytics.track('new_chat_started');
            }
            
            clearMessages() {
                this.state.messages = [];
                const messagesArea = document.getElementById('messagesArea');
                messagesArea.innerHTML = '';
                this.updateMessageCount();
                this.updateTokenUsage(0);
            }
            
            showWelcomeScreen() {
                const messagesArea = document.getElementById('messagesArea');
                messagesArea.innerHTML = document.getElementById('welcomeScreen').outerHTML;
                
                document.querySelectorAll('.quick-action').forEach(btn => {
                    btn.addEventListener('click', () => {
                        const prompt = btn.dataset.prompt;
                        document.getElementById('messageInput').value = prompt;
                        this.sendMessage();
                    });
                });
            }
            
            async sendMessage() {
                const input = document.getElementById('messageInput');
                const message = input.value.trim();
                
                if (!message && this.state.attachedFiles.length === 0) return;
                
                if (this.state.isStreaming) {
                    this.showNotification('Please wait for the current response', 'warning');
                    return;
                }
                
                const sendButton = document.getElementById('sendButton');
                sendButton.disabled = true;
                sendButton.classList.add('loading');
                
                input.value = '';
                this.updateCharCounter();
                this.autoResizeTextarea();
                
                if (document.getElementById('welcomeScreen')) {
                    document.getElementById('welcomeScreen').remove();
                }
                
                const userMessage = this.addMessage('user', message);
                
                const formData = new FormData();
                formData.append('message', message);
                formData.append('sessionId', this.state.sessionId);
                formData.append('mode', this.state.currentMode);
                formData.append('advancedMode', this.state.advancedMode);
                formData.append('stream', 'true');
                
                this.state.attachedFiles.forEach(file => {
                    formData.append('files', file);
                });
                
                this.showTypingIndicator();
                
                const startTime = performance.now();
                
                try {
                    const response = await fetch(API_ENDPOINTS.chat, {
                        method: 'POST',
                        body: formData
                    });
                    
                    if (!response.ok) {
                        throw new Error(`HTTP error! status: ${response.status}`);
                    }
                    
                    if (response.body) {
                        await this.handleStreamingResponse(response.body);
                    } else {
                        const data = await response.json();
                        this.hideTypingIndicator();
                        this.addMessage('assistant', data.response);
                        this.updateTokenUsage(data.tokenUsage);
                    }
                    
                    const endTime = performance.now();
                    const responseTime = endTime - startTime;
                    this.updatePerformanceMetrics(responseTime);
                    
                } catch (error) {
                    console.error('Error sending message:', error);
                    this.hideTypingIndicator();
                    this.showNotification('Failed to send message', 'error');
                    this.addMessage('assistant', 'I apologize, but I encountered an error. Please try again.');
                } finally {
                    sendButton.disabled = false;
                    sendButton.classList.remove('loading');
                    this.clearFiles();
                }
                
                this.analytics.track('message_sent', {
                    mode: this.state.currentMode,
                    hasFiles: this.state.attachedFiles.length > 0,
                    advancedMode: this.state.advancedMode
                });
            }
            
            async handleStreamingResponse(stream) {
                this.state.isStreaming = true;
                const reader = stream.getReader();
                const decoder = new TextDecoder();
                let assistantMessage = '';
                let messageElement = null;
                
                this.hideTypingIndicator();
                
                try {
                    while (true) {
                        const { done, value } = await reader.read();
                        if (done) break;
                        
                        const chunk = decoder.decode(value);
                        const lines = chunk.split('\n');
                        
                        for (const line of lines) {
                            if (line.startsWith('data: ')) {
                                const data = line.slice(6);
                                if (data === '[DONE]') {
                                    this.state.isStreaming = false;
                                    if (messageElement) {
                                        this.finalizeStreamingMessage(messageElement);
                                    }
                                } else {
                                    try {
                                        const parsed = JSON.parse(data);
                                        if (parsed.content) {
                                            assistantMessage += parsed.content;
                                            
                                            if (!messageElement) {
                                                messageElement = this.addMessage('assistant', assistantMessage, false);
                                            } else {
                                                this.updateStreamingMessage(messageElement, assistantMessage);
                                            }
                                        }
                                        
                                        if (parsed.tokenUsage) {
                                            this.updateTokenUsage(parsed.tokenUsage);
                                        }
                                    } catch (e) {
                                        console.error('Error parsing stream data:', e);
                                    }
                                }
                            }
                        }
                    }
                } catch (error) {
                    console.error('Stream reading error:', error);
                    this.showNotification('Stream interrupted', 'error');
                } finally {
                    this.state.isStreaming = false;
                    if (messageElement) {
                        this.finalizeStreamingMessage(messageElement);
                    }
                }
                
                if (assistantMessage) {
                    this.state.messages.push({
                        role: 'assistant',
                        content: assistantMessage,
                        timestamp: new Date(),
                        model: 'gemini-2.0-flash-exp'
                    });
                    await this.saveSession();
                }
            }
            
            addMessage(role, content, save = true) {
                const messageId = this.generateId();
                const messageElement = this.createMessageElement(role, content, messageId);
                
                document.getElementById('messagesArea').appendChild(messageElement);
                this.scrollToBottom();
                
                if (save && !this.state.isStreaming) {
                    this.state.messages.push({
                        id: messageId,
                        role: role,
                        content: content,
                        timestamp: new Date(),
                        reactions: {},
                        attachments: [...this.state.attachedFiles]
                    });
                    this.updateMessageCount();
                    this.saveSession();
                }
                
                return messageElement;
            }
            
            createMessageElement(role, content, messageId) {
                const messageDiv = document.createElement('div');
                messageDiv.className = `message ${role}`;
                messageDiv.id = messageId;
                messageDiv.setAttribute('data-message-id', messageId);
                
                const time = new Date().toLocaleTimeString([], {
                    hour: '2-digit',
                    minute: '2-digit'
                });
                
                const avatar = role === 'user' ? 'U' : '<i class="fas fa-robot"></i>';
                const author = role === 'user' ? 'You' : "Jack's AI Ultra";
                
                let processedContent = content;
                if (role === 'assistant' && !this.state.isStreaming) {
                    processedContent = this.processMessageContent(content);
                }
                
                messageDiv.innerHTML = `
                    <div class="message-avatar">
                        ${avatar}
                    </div>
                    <div class="message-content-wrapper">
                        <div class="message-header">
                            <span class="message-author">${author}</span>
                            ${role === 'assistant' ? `
                                <div class="message-badges">
                                    <span class="message-badge">Gemini 2.0</span>
                                    ${this.state.advancedMode ? '<span class="message-badge">Advanced</span>' : ''}
                                </div>
                            ` : ''}
                            <span class="message-time">${time}</span>
                        </div>
                        <div class="message-bubble">
                            ${processedContent}
                            ${this.state.isStreaming && role === 'assistant' ? '<span class="streaming-indicator"><span class="streaming-dots"><span class="streaming-dot"></span><span class="streaming-dot"></span><span class="streaming-dot"></span></span></span>' : ''}
                        </div>
                        ${this.state.attachedFiles.length > 0 && role === 'user' ? this.createAttachmentsHTML() : ''}
                        <div class="message-reactions"></div>
                        <div class="message-actions">
                            <button class="message-action" data-action="copy">
                                <i class="fas fa-copy"></i> Copy
                            </button>
                            ${role === 'assistant' ? `
                                <button class="message-action" data-action="regenerate">
                                    <i class="fas fa-redo"></i> Regenerate
                                </button>
                            ` : ''}
                            <button class="message-action" data-action="share">
                                <i class="fas fa-share"></i> Share
                            </button>
                        </div>
                    </div>
                `;
                
                messageDiv.querySelectorAll('.message-action').forEach(btn => {
                    btn.addEventListener('click', (e) => {
                        e.stopPropagation();
                        const action = btn.dataset.action;
                        this.handleMessageAction(action, messageId);
                    });
                });
                
                return messageDiv;
            }
            
            updateStreamingMessage(element, content) {
                const bubble = element.querySelector('.message-bubble');
                if (bubble) {
                    const processedContent = this.processMessageContent(content);
                    bubble.innerHTML = processedContent + '<span class="streaming-indicator"><span class="streaming-dots"><span class="streaming-dot"></span><span class="streaming-dot"></span><span class="streaming-dot"></span></span></span>';
                    
                    this.highlightCode(element);
                    this.scrollToBottom();
                }
            }
            
            finalizeStreamingMessage(element) {
                const bubble = element.querySelector('.message-bubble');
                const indicator = bubble.querySelector('.streaming-indicator');
                if (indicator) {
                    indicator.remove();
                }
                
                this.highlightCode(element);
            }
            
            processMessageContent(content) {
                const renderer = new marked.Renderer();
                
                renderer.code = (code, language) => {
                    const validLanguage = Prism.languages[language] ? language : 'plaintext';
                    const highlighted = Prism.highlight(code, Prism.languages[validLanguage] || Prism.languages.plaintext, validLanguage);
                    return `<pre><code class="language-${validLanguage}">${highlighted}</code></pre>`;
                };
                
                renderer.link = (href, title, text) => {
                    return `<a href="${href}" title="${title || ''}" target="_blank" rel="noopener noreferrer">${text}</a>`;
                };
                
                marked.setOptions({
                    renderer: renderer,
                    breaks: true,
                    gfm: true,
                    sanitize: false
                });
                
                let processed = marked.parse(content);
                
                processed = DOMPurify.sanitize(processed, {
                    ALLOWED_TAGS: ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'br', 'strong', 'em', 'u', 's', 'code', 'pre', 'blockquote', 'ul', 'ol', 'li', 'a', 'img', 'table', 'thead', 'tbody', 'tr', 'th', 'td', 'hr', 'sup', 'sub'],
                    ALLOWED_ATTR: ['href', 'title', 'target', 'rel', 'src', 'alt', 'class', 'id'],
                    ALLOWED_URI_REGEXP: /^(?:(?:https?|mailto|tel|sms):|[^a-z]|[a-z+.\-]+(?:[^a-z+.\-:]|$))/i
                });
                
                return processed;
            }
            
            highlightCode(element) {
                element.querySelectorAll('pre code').forEach(block => {
                    if (!block.classList.contains('highlighted')) {
                        Prism.highlightElement(block);
                        block.classList.add('highlighted');
                    }
                });
            }
            
            createAttachmentsHTML() {
                return `
                    <div class="message-attachments">
                        ${this.state.attachedFiles.map(file => `
                            <div class="attachment-item">
                                <i class="fas ${this.getFileIcon(file)} attachment-icon"></i>
                                <span class="attachment-name">${file.name}</span>
                                <span class="attachment-size">${this.formatFileSize(file.size)}</span>
                            </div>
                        `).join('')}
                    </div>
                `;
            }
            
            getFileIcon(file) {
                const ext = file.name.split('.').pop().toLowerCase();
                const iconMap = {
                    pdf: 'fa-file-pdf',
                    doc: 'fa-file-word',
                    docx: 'fa-file-word',
                    xls: 'fa-file-excel',
                    xlsx: 'fa-file-excel',
                    ppt: 'fa-file-powerpoint',
                    pptx: 'fa-file-powerpoint',
                    zip: 'fa-file-archive',
                    rar: 'fa-file-archive',
                    '7z': 'fa-file-archive',
                    jpg: 'fa-file-image',
                    jpeg: 'fa-file-image',
                    png: 'fa-file-image',
                    gif: 'fa-file-image',
                    mp4: 'fa-file-video',
                    avi: 'fa-file-video',
                    mov: 'fa-file-video',
                    mp3: 'fa-file-audio',
                    wav: 'fa-file-audio',
                    flac: 'fa-file-audio',
                    py: 'fa-file-code',
                    js: 'fa-file-code',
                    html: 'fa-file-code',
                    css: 'fa-file-code',
                    json: 'fa-file-code',
                    xml: 'fa-file-code',
                    csv: 'fa-file-csv',
                    txt: 'fa-file-alt'
                };
                
                return iconMap[ext] || 'fa-file';
            }
            
            formatFileSize(bytes) {
                if (bytes === 0) return '0 Bytes';
                
                const k = 1024;
                const sizes = ['Bytes', 'KB', 'MB', 'GB'];
                const i = Math.floor(Math.log(bytes) / Math.log(k));
                
                return parseFloat((bytes / Math.pow(k, i)).toFixed(2)) + ' ' + sizes[i];
            }
            
            handleMessageAction(action, messageId) {
                switch (action) {
                    case 'copy':
                        this.copyMessage(messageId);
                        break;
                    case 'regenerate':
                        this.regenerateMessage(messageId);
                        break;
                    case 'share':
                        this.shareMessage(messageId);
                        break;
                }
            }
            
            copyMessage(messageId) {
                const message = document.getElementById(messageId);
                if (!message) return;
                
                const bubble = message.querySelector('.message-bubble');
                const text = bubble.textContent;
                
                navigator.clipboard.writeText(text).then(() => {
                    this.showNotification('Message copied to clipboard', 'success');
                }).catch(err => {
                    console.error('Failed to copy:', err);
                    this.showNotification('Failed to copy message', 'error');
                });
            }
            
            async regenerateMessage(messageId) {
                const messageIndex = this.state.messages.findIndex(m => m.id === messageId);
                if (messageIndex === -1) return;
                
                const previousUserMessage = this.state.messages.slice(0, messageIndex).reverse().find(m => m.role === 'user');
                if (!previousUserMessage) return;
                
                const confirmed = await this.ui.showConfirmDialog(
                    'Regenerate Response?',
                    'This will generate a new response to your previous message.'
                );
                
                if (!confirmed) return;
                
                const messageElement = document.getElementById(messageId);
                messageElement.remove();
                
                this.state.messages = this.state.messages.slice(0, messageIndex);
                
                const input = document.getElementById('messageInput');
                input.value = previousUserMessage.content;
                await this.sendMessage();
            }
            
            shareMessage(messageId) {
                const message = this.state.messages.find(m => m.id === messageId);
                if (!message) return;
                
                const shareUrl = `${window.location.origin}/share/${messageId}`;
                
                if (navigator.share) {
                    navigator.share({
                        title: "Jack's AI Ultra - Shared Message",
                        text: message.content.substring(0, 200) + '...',
                        url: shareUrl
                    }).catch(err => console.log('Error sharing:', err));
                } else {
                    navigator.clipboard.writeText(shareUrl).then(() => {
                        this.showNotification('Share link copied to clipboard', 'success');
                    });
                }
            }
            
            handleContextMenuAction(action) {
                const target = this.state.ui.contextMenuTarget;
                if (!target) return;
                
                const messageId = target.getAttribute('data-message-id');
                
                switch (action) {
                    case 'copy':
                        this.copyMessage(messageId);
                        break;
                    case 'edit':
                        this.editMessage(messageId);
                        break;
                    case 'regenerate':
                        this.regenerateMessage(messageId);
                        break;
                    case 'pin':
                        this.pinMessage(messageId);
                        break;
                    case 'branch':
                        this.createBranch(messageId);
                        break;
                    case 'delete':
                        this.deleteMessage(messageId);
                        break;
                }
            }
            
            async editMessage(messageId) {
                const message = this.state.messages.find(m => m.id === messageId);
                if (!message || message.role !== 'user') return;
                
                const newContent = await this.ui.showInputDialog(
                    'Edit Message',
                    'Edit your message:',
                    message.content
                );
                
                if (newContent && newContent !== message.content) {
                    message.content = newContent;
                    message.edited = true;
                    message.editedAt = new Date();
                    
                    const messageElement = document.getElementById(messageId);
                    const bubble = messageElement.querySelector('.message-bubble');
                    bubble.innerHTML = this.processMessageContent(newContent);
                    
                    await this.saveSession();
                    this.showNotification('Message edited', 'success');
                }
            }
            
            pinMessage(messageId) {
                const message = this.state.messages.find(m => m.id === messageId);
                if (!message) return;
                
                message.pinned = !message.pinned;
                
                const messageElement = document.getElementById(messageId);
                messageElement.classList.toggle('pinned', message.pinned);
                
                this.saveSession();
                this.showNotification(message.pinned ? 'Message pinned' : 'Message unpinned', 'success');
            }
            
            async createBranch(messageId) {
                const messageIndex = this.state.messages.findIndex(m => m.id === messageId);
                if (messageIndex === -1) return;
                
                const branchName = await this.ui.showInputDialog(
                    'Create Branch',
                    'Enter a name for this conversation branch:'
                );
                
                if (!branchName) return;
                
                const branch = {
                    id: this.generateId(),
                    name: branchName,
                    messageId: messageId,
                    createdAt: new Date(),
                    messages: this.state.messages.slice(0, messageIndex + 1)
                };
                
                this.state.branches.push(branch);
                await this.saveSession();
                
                this.showNotification(`Branch "${branchName}" created`, 'success');
            }
            
            async deleteMessage(messageId) {
                const confirmed = await this.ui.showConfirmDialog(
                    'Delete Message?',
                    'This action cannot be undone.'
                );
                
                if (!confirmed) return;
                
                const messageIndex = this.state.messages.findIndex(m => m.id === messageId);
                if (messageIndex === -1) return;
                
                this.state.messages.splice(messageIndex, 1);
                
                const messageElement = document.getElementById(messageId);
                messageElement.remove();
                
                await this.saveSession();
                this.updateMessageCount();
                this.showNotification('Message deleted', 'success');
            }
            
            searchMessages(query) {
                this.state.ui.searchQuery = query.toLowerCase();
                
                const messages = document.querySelectorAll('.message');
                let foundCount = 0;
                
                messages.forEach(message => {
                    const content = message.textContent.toLowerCase();
                    const matches = query && content.includes(query);
                    
                    message.style.display = !query || matches ? 'flex' : 'none';
                    
                    if (matches) {
                        foundCount++;
                        this.highlightSearchTerms(message, query);
                    } else {
                        this.clearSearchHighlights(message);
                    }
                });
                
                if (query) {
                    this.showNotification(`Found ${foundCount} message${foundCount !== 1 ? 's' : ''}`, 'info');
                }
            }
            
            highlightSearchTerms(element, query) {
                
            }
            
            clearSearchHighlights(element) {
                
            }
            
            changeMode(mode) {
                this.state.currentMode = mode;
                
                document.querySelectorAll('.mode-btn').forEach(btn => {
                    btn.classList.toggle('active', btn.dataset.mode === mode);
                });
                
                const modeDescriptions = {
                    balanced: 'Balanced mode activated',
                    creative: 'Creative mode activated - Unleashing imagination',
                    precise: 'Precise mode activated - Maximum accuracy',
                    code: 'Code mode activated - Professional development',
                    research: 'Research mode activated - Deep analysis'
                };
                
                this.showNotification(modeDescriptions[mode], 'success');
                this.analytics.track('mode_changed', { mode });
            }
            
            toggleAdvancedMode() {
                this.state.advancedMode = !this.state.advancedMode;
                const toggle = document.getElementById('advancedToggle');
                toggle.classList.toggle('active', this.state.advancedMode);
                
                if (this.state.advancedMode) {
                    this.showNotification('🚀 ULTRA MODE ACTIVATED - Maximum capability unlocked!', 'warning');
                } else {
                    this.showNotification('Advanced mode deactivated', 'info');
                }
                
                this.analytics.track('advanced_mode_toggled', { active: this.state.advancedMode });
            }
            
            toggleFeature(feature) {
                const btn = document.querySelector(`[data-feature="${feature}"]`);
                const isActive = btn.classList.contains('active');
                
                switch (feature) {
                    case 'files':
                        this.toggleFileUpload();
                        break;
                    case 'image':
                        this.toggleImageUpload();
                        break;
                    case 'voice':
                        this.toggleVoiceInput();
                        break;
                    case 'templates':
                        this.showTemplates();
                        break;
                    case 'snippets':
                        this.showSnippets();
                        break;
                    case 'context':
                        this.showContextManager();
                        break;
                }
                
                btn.classList.toggle('active', !isActive);
                this.state.features[`${feature}Enabled`] = !isActive;
            }
            
            toggleFileUpload() {
                const uploadArea = document.getElementById('fileUploadArea');
                uploadArea.classList.toggle('active');
            }
            
            showFileUploadArea() {
                const uploadArea = document.getElementById('fileUploadArea');
                uploadArea.classList.add('active');
                document.querySelector('[data-feature="files"]').classList.add('active');
            }
            
            toggleImageUpload() {
                document.getElementById('fileInput').accept = 'image/*';
                document.getElementById('fileInput').click();
            }
            
            toggleVoiceInput() {
                if (!this.state.voiceRecognition) {
                    this.showNotification('Voice input not supported', 'error');
                    return;
                }
                
                if (this.state.features.voiceEnabled) {
                    this.stopVoiceRecording();
                } else {
                    this.startVoiceRecording();
                }
            }
            
            startVoiceRecording() {
                this.state.voiceRecognition.start();
                this.state.features.voiceEnabled = true;
                document.querySelector('[data-feature="voice"]').classList.add('active');
                this.showNotification('Voice recording started', 'info');
            }
            
            stopVoiceRecording() {
                this.state.voiceRecognition.stop();
                this.state.features.voiceEnabled = false;
                document.querySelector('[data-feature="voice"]').classList.remove('active');
                this.showNotification('Voice recording stopped', 'info');
            }
            
            showTemplates() {
                const templates = [
                    { name: 'Code Review', prompt: 'Please review this code for best practices, potential bugs, and optimization opportunities:' },
                    { name: 'Explain Concept', prompt: 'Can you explain [concept] in simple terms with examples?' },
                    { name: 'Debug Help', prompt: "I'm getting this error: [error]. Here's my code:" },
                    { name: 'Research Topic', prompt: 'I need a comprehensive analysis of [topic] including:' },
                    { name: 'Creative Writing', prompt: 'Write a creative story about [theme] that includes:' },
                    { name: 'Data Analysis', prompt: 'Analyze this data and provide insights:' },
                    { name: 'Business Strategy', prompt: 'Develop a strategy for [business goal]:' },
                    { name: 'Learning Plan', prompt: 'Create a learning plan for [skill/topic]:' }
                ];
                
                this.ui.showTemplateDialog(templates, (template) => {
                    document.getElementById('messageInput').value = template.prompt;
                    this.updateCharCounter();
                });
            }
            
            showSnippets() {
                const snippets = this.state.cache.get('snippets') || [];
                
                this.ui.showSnippetDialog(snippets, (snippet) => {
                    const input = document.getElementById('messageInput');
                    const start = input.selectionStart;
                    const end = input.selectionEnd;
                    const text = input.value;
                    
                    input.value = text.substring(0, start) + snippet.content + text.substring(end);
                    input.selectionStart = input.selectionEnd = start + snippet.content.length;
                    input.focus();
                    
                    this.updateCharCounter();
                });
            }
            
            showContextManager() {
                this.ui.showContextDialog({
                    currentContext: this.state.currentContext,
                    availableContexts: this.state.availableContexts
                }, (context) => {
                    this.state.currentContext = context;
                    this.showNotification('Context updated', 'success');
                });
            }
            
            handleFiles(files) {
                const fileList = document.getElementById('fileList');
                fileList.innerHTML = '';
                
                Array.from(files).forEach(file => {
                    if (file.size > 2 * 1024 * 1024 * 1024) {
                        this.showNotification(`File ${file.name} is too large (max 2GB)`, 'error');
                        return;
                    }
                    
                    this.state.attachedFiles.push(file);
                    
                    const fileItem = document.createElement('div');
                    fileItem.className = 'file-item';
                    fileItem.innerHTML = `
                        <div class="file-item-content">
                            <i class="fas ${this.getFileIcon(file)}"></i>
                            <span>${file.name}</span>
                            <i class="fas fa-times remove-file"></i>
                        </div>
                    `;
                    
                    fileItem.querySelector('.remove-file').addEventListener('click', () => {
                        this.removeFile(file, fileItem);
                    });
                    
                    fileList.appendChild(fileItem);
                    
                    this.uploadFile(file, fileItem);
                });
                
                document.getElementById('fileInput').value = '';
            }
            
            async uploadFile(file, fileItem) {
                try {
                    const formData = new FormData();
                    formData.append('file', file);
                    
                    const response = await fetch(API_ENDPOINTS.upload, {
                        method: 'POST',
                        body: formData
                    });
                    
                    if (!response.ok) throw new Error('Upload failed');
                    
                    const result = await response.json();
                    fileItem.setAttribute('data-file-id', result.fileId);
                    fileItem.style.setProperty('--progress', '1');
                    
                    this.showNotification(`${file.name} uploaded successfully`, 'success');
                } catch (error) {
                    console.error('Upload error:', error);
                    this.showNotification(`Failed to upload ${file.name}`, 'error');
                    this.removeFile(file, fileItem);
                }
            }
            
            removeFile(file, fileItem) {
                const index = this.state.attachedFiles.indexOf(file);
                if (index > -1) {
                    this.state.attachedFiles.splice(index, 1);
                }
                fileItem.remove();
                
                if (this.state.attachedFiles.length === 0) {
                    document.getElementById('fileUploadArea').classList.remove('active');
                    document.querySelector('[data-feature="files"]').classList.remove('active');
                }
            }
            
            clearFiles() {
                this.state.attachedFiles = [];
                document.getElementById('fileList').innerHTML = '';
                document.getElementById('fileUploadArea').classList.remove('active');
                document.querySelector('[data-feature="files"]').classList.remove('active');
            }
            
            handleInputKeydown(e) {
                if (e.key === 'Enter' && !e.shiftKey) {
                    e.preventDefault();
                    this.sendMessage();
                } else if (e.key === 'Escape') {
                    this.clearInput();
                }
            }
            
            clearInput() {
                document.getElementById('messageInput').value = '';
                this.updateCharCounter();
                this.autoResizeTextarea();
            }
            
            updateCharCounter() {
                const input = document.getElementById('messageInput');
                const counter = document.getElementById('charCounter');
                counter.textContent = input.value.length;
            }
            
            autoResizeTextarea() {
                const textarea = document.getElementById('messageInput');
                textarea.style.height = 'auto';
                textarea.style.height = Math.min(textarea.scrollHeight, 200) + 'px';
            }
            
            showTypingIndicator() {
                document.getElementById('typingIndicator').classList.add('active');
                this.scrollToBottom();
            }
            
            hideTypingIndicator() {
                document.getElementById('typingIndicator').classList.remove('active');
            }
            
            updateTokenUsage(usage) {
                this.state.tokenUsage = usage || 0;
                const percentage = (this.state.tokenUsage / this.state.maxTokens) * 100;
                
                document.getElementById('tokenFill').style.width = `${percentage}%`;
                document.getElementById('tokenValue').textContent = `${this.formatNumber(this.state.tokenUsage)} / ${this.formatNumber(this.state.maxTokens)}`;
                document.getElementById('tokenCount').textContent = this.formatNumber(this.state.tokenUsage);
                
                if (percentage > 80) {
                    this.showNotification('Approaching token limit', 'warning');
                }
            }
            
            updateMessageCount() {
                const count = this.state.messages.length;
                document.getElementById('messageCount').textContent = count;
            }
            
            updatePerformanceMetrics(responseTime) {
                this.state.performance.lastResponseTime = responseTime;
                this.state.performance.totalMessages++;
                
                const avgTime = this.state.performance.averageResponseTime;
                const totalMessages = this.state.performance.totalMessages;
                
                this.state.performance.averageResponseTime = 
                    (avgTime * (totalMessages - 1) + responseTime) / totalMessages;
            }
            
            formatNumber(num) {
                if (num >= 1000000) {
                    return (num / 1000000).toFixed(1) + 'M';
                } else if (num >= 1000) {
                    return (num / 1000).toFixed(1) + 'K';
                }
                return num.toString();
            }
            
            scrollToBottom() {
                const messagesArea = document.getElementById('messagesArea');
                messagesArea.scrollTop = messagesArea.scrollHeight;
            }
            
            generateId() {
                return 'msg-' + Date.now() + '-' + Math.random().toString(36).substr(2, 9);
            }
            
            async saveSession() {
                if (!this.state.settings.autoSave) return;
                
                try {
                    await this.api.updateSession({
                        sessionId: this.state.sessionId,
                        messages: this.state.messages,
                        tokenUsage: this.state.tokenUsage,
                        mode: this.state.currentMode,
                        metadata: {
                            advancedMode: this.state.advancedMode,
                            branches: this.state.branches
                        }
                    });
                } catch (error) {
                    console.error('Failed to save session:', error);
                }
            }
            
            async exportChat() {
                const format = await this.ui.showExportDialog();
                if (!format) return;
                
                try {
                    const response = await fetch(API_ENDPOINTS.export.replace(':format', format) + `?sessionId=${this.state.sessionId}`);
                    
                    if (!response.ok) throw new Error('Export failed');
                    
                    const blob = await response.blob();
                    const url = URL.createObjectURL(blob);
                    const a = document.createElement('a');
                    a.href = url;
                    a.download = `chat-export-${new Date().toISOString().split('T')[0]}.${format}`;
                    a.click();
                    URL.revokeObjectURL(url);
                    
                    this.showNotification('Chat exported successfully', 'success');
                    this.analytics.track('chat_exported', { format });
                } catch (error) {
                    console.error('Export error:', error);
                    this.showNotification('Failed to export chat', 'error');
                }
            }
            
            async shareChat() {
                try {
                    const response = await this.api.shareSession(this.state.sessionId);
                    const shareUrl = `${window.location.origin}/shared/${response.shareId}`;
                    
                    if (navigator.share) {
                        await navigator.share({
                            title: "Jack's AI Ultra - Shared Conversation",
                            text: 'Check out this AI conversation',
                            url: shareUrl
                        });
                    } else {
                        await navigator.clipboard.writeText(shareUrl);
                        this.showNotification('Share link copied to clipboard', 'success');
                    }
                    
                    this.analytics.track('chat_shared');
                } catch (error) {
                    console.error('Share error:', error);
                    this.showNotification('Failed to share chat', 'error');
                }
            }
            
            toggleSetting(setting) {
                this.state.settings[setting] = !this.state.settings[setting];
                
                const switchEl = document.querySelector(`[data-setting="${setting}"]`);
                switchEl.classList.toggle('active', this.state.settings[setting]);
                
                this.applySettings();
                this.saveSettings();
                
                this.showNotification(`${setting} ${this.state.settings[setting] ? 'enabled' : 'disabled'}`, 'info');
            }
            
            saveSettings() {
                localStorage.setItem('ai_ultra_settings', JSON.stringify(this.state.settings));
            }
            
            showLoading(show) {
                const overlay = document.getElementById('loadingOverlay');
                overlay.classList.toggle('active', show);
            }
            
            showNotification(message, type = 'info', duration = 5000) {
                if (!this.state.settings.notifications && type !== 'error') return;
                
                const container = document.getElementById('notificationContainer');
                const notification = document.createElement('div');
                notification.className = `notification ${type}`;
                
                const icons = {
                    success: 'check-circle',
                    error: 'exclamation-circle',
                    warning: 'exclamation-triangle',
                    info: 'info-circle'
                };
                
                notification.innerHTML = `
                    <i class="fas fa-${icons[type]} notification-icon"></i>
                    <div class="notification-content">
                        <div class="notification-message">${message}</div>
                    </div>
                    <i class="fas fa-times notification-close"></i>
                `;
                
                const closeBtn = notification.querySelector('.notification-close');
                closeBtn.addEventListener('click', () => {
                    notification.remove();
                });
                
                container.appendChild(notification);
                
                setTimeout(() => {
                    notification.style.opacity = '0';
                    setTimeout(() => notification.remove(), 300);
                }, duration);
            }
        }
        
        class APIClient {
            async request(url, options = {}) {
                const defaultOptions = {
                    headers: {
                        'Content-Type': 'application/json',
                        'X-Session-ID': AppState.sessionId
                    }
                };
                
                if (options.body && !(options.body instanceof FormData)) {
                    options.body = JSON.stringify(options.body);
                }
                
                const response = await fetch(url, { ...defaultOptions, ...options });
                
                if (!response.ok) {
                    throw new Error(`API error: ${response.status}`);
                }
                
                return response.json();
            }
            
            createSession(data) {
                return this.request(API_ENDPOINTS.session.create, {
                    method: 'POST',
                    body: data
                });
            }
            
            getSession(sessionId) {
                return this.request(API_ENDPOINTS.session.get.replace(':id', sessionId));
            }
            
            updateSession(data) {
                return this.request(API_ENDPOINTS.session.update, {
                    method: 'PUT',
                    body: data
                });
            }
            
            getModels() {
                return this.request(API_ENDPOINTS.models);
            }
            
            shareSession(sessionId) {
                return this.request(API_ENDPOINTS.share, {
                    method: 'POST',
                    body: { sessionId }
                });
            }
        }
        
        class UIManager {
            initialize() {
                this.setupTheme();
                this.loadSessions();
                this.updateUI();
            }
            
            setupTheme() {
                const theme = AppState.settings.darkMode ? 'dark' : 'light';
                document.documentElement.setAttribute('data-theme', theme);
            }
            
            async loadSessions() {
                try {
                    const sessions = await app.api.request(API_ENDPOINTS.session.list);
                    this.displaySessions(sessions);
                } catch (error) {
                    console.error('Failed to load sessions:', error);
                }
            }
            
            displaySessions(sessions) {
                const container = document.getElementById('chatSessions');
                const existingHeader = container.querySelector('.sessions-header');
                
                container.innerHTML = '';
                if (existingHeader) {
                    container.appendChild(existingHeader);
                }
                
                sessions.forEach(session => {
                    const sessionItem = this.createSessionItem(session);
                    container.appendChild(sessionItem);
                });
            }
            
            createSessionItem(session) {
                const div = document.createElement('div');
                div.className = 'session-item';
                div.setAttribute('data-session-id', session.id);
                
                if (session.id === AppState.sessionId) {
                    div.classList.add('active');
                }
                
                const lastMessage = session.messages[session.messages.length - 1];
                const preview = lastMessage ? lastMessage.content.substring(0, 100) + '...' : 'New conversation';
                const time = new Date(session.updated_at).toLocaleDateString();
                
                div.innerHTML = `
                    <div class="session-header">
                        <div class="session-title">${session.title || 'Untitled Chat'}</div>
                        <div class="session-mode">${session.mode}</div>
                    </div>
                    <div class="session-preview">${preview}</div>
                    <div class="session-meta">
                        <div class="session-meta-item">
                            <i class="fas fa-comments"></i>
                            <span>${session.messages.length}</span>
                        </div>
                        <div class="session-meta-item">
                            <i class="fas fa-clock"></i>
                            <span>${time}</span>
                        </div>
                    </div>
                `;
                
                div.addEventListener('click', () => {
                    app.session.switchSession(session.id);
                });
                
                return div;
            }
            
            updateUI() {
                this.updateMessageCount();
                this.updateTokenDisplay();
            }
            
            updateMessageCount() {
                const count = AppState.messages.length;
                document.getElementById('messageCount').textContent = count;
            }
            
            updateTokenDisplay() {
                const percentage = (AppState.tokenUsage / AppState.maxTokens) * 100;
                document.getElementById('tokenFill').style.width = `${percentage}%`;
                document.getElementById('tokenValue').textContent = 
                    `${app.formatNumber(AppState.tokenUsage)} / ${app.formatNumber(AppState.maxTokens)}`;
            }
            
            toggleSidebar() {
                AppState.ui.sidebarCollapsed = !AppState.ui.sidebarCollapsed;
                document.getElementById('sidebar').classList.toggle('collapsed', AppState.ui.sidebarCollapsed);
            }
            
            showContextMenu(x, y) {
                const menu = document.getElementById('contextMenu');
                menu.style.left = x + 'px';
                menu.style.top = y + 'px';
                menu.classList.add('active');
            }
            
            hideContextMenu() {
                document.getElementById('contextMenu').classList.remove('active');
                AppState.ui.contextMenuTarget = null;
            }
            
            openSettings() {
                this.showModal('Settings', this.createSettingsContent());
            }
            
            createSettingsContent() {
                return document.getElementById('modalBody').innerHTML;
            }
            
            showModal(title, content) {
                document.getElementById('modalTitle').textContent = title;
                
                if (typeof content === 'string') {
                    document.getElementById('modalBody').innerHTML = content;
                } else {
                    document.getElementById('modalBody').innerHTML = '';
                    document.getElementById('modalBody').appendChild(content);
                }
                
                document.getElementById('modalOverlay').classList.add('active');
                AppState.ui.modalOpen = true;
            }
            
            closeModal() {
                document.getElementById('modalOverlay').classList.remove('active');
                AppState.ui.modalOpen = false;
            }
            
            async showConfirmDialog(title, message) {
                return new Promise(resolve => {
                    const content = `
                        <p>${message}</p>
                        <div style="display: flex; gap: 12px; justify-content: flex-end; margin-top: 24px;">
                            <button class="btn-secondary" onclick="app.ui.closeModal(); window.__confirmResolve(false)">Cancel</button>
                            <button class="btn-primary" onclick="app.ui.closeModal(); window.__confirmResolve(true)">Confirm</button>
                        </div>
                    `;
                    
                    window.__confirmResolve = resolve;
                    this.showModal(title, content);
                });
            }
            
            async showInputDialog(title, message, defaultValue = '') {
                return new Promise(resolve => {
                    const content = `
                        <p>${message}</p>
                        <input type="text" id="dialogInput" class="dialog-input" value="${defaultValue}" style="width: 100%; margin-top: 16px;">
                        <div style="display: flex; gap: 12px; justify-content: flex-end; margin-top: 24px;">
                            <button class="btn-secondary" onclick="app.ui.closeModal(); window.__inputResolve(null)">Cancel</button>
                            <button class="btn-primary" onclick="app.ui.closeModal(); window.__inputResolve(document.getElementById('dialogInput').value)">OK</button>
                        </div>
                    `;
                    
                    window.__inputResolve = resolve;
                    this.showModal(title, content);
                    
                    setTimeout(() => {
                        const input = document.getElementById('dialogInput');
                        if (input) {
                            input.focus();
                            input.select();
                        }
                    }, 100);
                });
            }
            
            async showExportDialog() {
                return new Promise(resolve => {
                    const content = `
                        <p>Select export format:</p>
                        <div style="display: flex; flex-direction: column; gap: 12px; margin-top: 16px;">
                            <label class="export-option">
                                <input type="radio" name="exportFormat" value="markdown" checked>
                                <span>Markdown (.md)</span>
                            </label>
                            <label class="export-option">
                                <input type="radio" name="exportFormat" value="json">
                                <span>JSON (.json)</span>
                            </label>
                            <label class="export-option">
                                <input type="radio" name="exportFormat" value="pdf">
                                <span>PDF (.pdf)</span>
                            </label>
                            <label class="export-option">
                                <input type="radio" name="exportFormat" value="html">
                                <span>HTML (.html)</span>
                            </label>
                        </div>
                        <div style="display: flex; gap: 12px; justify-content: flex-end; margin-top: 24px;">
                            <button class="btn-secondary" onclick="app.ui.closeModal(); window.__exportResolve(null)">Cancel</button>
                            <button class="btn-primary" onclick="app.ui.closeModal(); window.__exportResolve(document.querySelector('input[name=exportFormat]:checked').value)">Export</button>
                        </div>
                    `;
                    
                    window.__exportResolve = resolve;
                    this.showModal('Export Chat', content);
                });
            }
            
            showTemplateDialog(templates, callback) {
                const content = document.createElement('div');
                content.className = 'template-list';
                
                templates.forEach(template => {
                    const item = document.createElement('div');
                    item.className = 'template-item';
                    item.innerHTML = `
                        <h4>${template.name}</h4>
                        <p>${template.prompt}</p>
                    `;
                    item.addEventListener('click', () => {
                        this.closeModal();
                        callback(template);
                    });
                    content.appendChild(item);
                });
                
                this.showModal('Templates', content);
            }
            
            showNotificationCenter() {
                app.showNotification('Notification center coming soon!', 'info');
            }
        }
        
        class ChatManager {
            
        }
        
        class SessionManager {
            loadSession(sessionData) {
                AppState.messages = sessionData.messages || [];
                AppState.tokenUsage = sessionData.tokenUsage || 0;
                AppState.currentMode = sessionData.mode || 'balanced';
                
                this.displayMessages();
                app.ui.updateUI();
            }
            
            displayMessages() {
                const messagesArea = document.getElementById('messagesArea');
                messagesArea.innerHTML = '';
                
                if (AppState.messages.length === 0) {
                    app.showWelcomeScreen();
                } else {
                    AppState.messages.forEach(msg => {
                        app.addMessage(msg.role, msg.content, false);
                    });
                }
            }
            
            async switchSession(sessionId) {
                if (sessionId === AppState.sessionId) return;
                
                await app.saveSession();
                
                AppState.sessionId = sessionId;
                localStorage.setItem('ai_ultra_session_id', sessionId);
                
                const session = await app.api.getSession(sessionId);
                this.loadSession(session);
                
                document.querySelectorAll('.session-item').forEach(item => {
                    item.classList.toggle('active', item.getAttribute('data-session-id') === sessionId);
                });
                
                app.showNotification('Switched to session', 'info');
            }
        }
        
        class VoiceManager {
            
        }
        
        class FileManager {
            
        }
        
        class SettingsManager {
            
        }
        
        class PerformanceMonitor {
            
        }
        
        class AnalyticsTracker {
            track(event, data = {}) {
                console.log('Analytics:', event, data);
            }
        }
        
        class ParticleSystem {
            create() {
                const particleField = document.getElementById('particleField');
                const particleCount = 50;
                
                for (let i = 0; i < particleCount; i++) {
                    const particle = document.createElement('div');
                    particle.className = 'particle';
                    particle.style.left = Math.random() * 100 + '%';
                    particle.style.animationDelay = Math.random() * 20 + 's';
                    particle.style.animationDuration = (15 + Math.random() * 10) + 's';
                    particleField.appendChild(particle);
                }
            }
        }
        
        class ShortcutManager {
            initialize() {
                document.addEventListener('keydown', (e) => {
                    if (e.ctrlKey || e.metaKey) {
                        switch (e.key) {
                            case 'k':
                                e.preventDefault();
                                document.getElementById('searchInput').focus();
                                break;
                            case 'n':
                                e.preventDefault();
                                app.startNewChat();
                                break;
                            case 's':
                                e.preventDefault();
                                app.saveSession();
                                app.showNotification('Session saved', 'success');
                                break;
                            case '/':
                                e.preventDefault();
                                document.getElementById('messageInput').focus();
                                break;
                        }
                    }
                });
            }
        }
        
        const app = new AIUltraApp();
        
        window.app = app;
    </script>
</body>
</html>"""

@app.route('/')
def index():
    return render_template_string(ULTIMATE_HTML_TEMPLATE)

@app.route('/api/session/create', methods=['POST'])
@limiter.limit("20 per minute")
def create_session():
    try:
        data = request.json
        mode = ChatMode(data.get('mode', 'balanced'))
        
        session = session_manager.create_session(
            user_id=data.get('userId'),
            mode=mode
        )
        
        return jsonify({
            'success': True,
            'sessionId': session.id,
            'mode': mode.value,
            'metadata': session.metadata
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/session/<session_id>', methods=['GET'])
@cache.cached(timeout=300)
def get_session(session_id):
    try:
        session = session_manager.get_session(session_id)
        if not session:
            return jsonify({'success': False, 'error': 'Session not found'}), 404
        
        return jsonify({
            'success': True,
            **session.to_dict()
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/session/update', methods=['PUT'])
@limiter.limit("50 per minute")
def update_session():
    try:
        data = request.json
        session_id = data.get('sessionId')
        
        session = session_manager.get_session(session_id)
        if not session:
            return jsonify({'success': False, 'error': 'Session not found'}), 404
        
        if 'metadata' in data:
            session.metadata.update(data['metadata'])
        
        if 'settings' in data:
            session.settings.update(data['settings'])
        
        session_manager.update_session(session)
        
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/sessions', methods=['GET'])
@limiter.limit("30 per minute")
def list_sessions():
    try:
        user_id = request.headers.get('X-User-ID')
        limit = int(request.args.get('limit', 100))
        offset = int(request.args.get('offset', 0))
        
        if user_id:
            sessions = session_manager.get_user_sessions(user_id, limit, offset)
        else:
            sessions = list(session_manager.sessions.values())[:limit]
        
        return jsonify({
            'success': True,
            'sessions': [s.to_dict() for s in sessions],
            'total': len(sessions)
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/sessions/search', methods=['POST'])
@limiter.limit("20 per minute")
def search_sessions():
    try:
        data = request.json
        query = data.get('query', '')
        user_id = data.get('userId')
        mode = ChatMode(data.get('mode')) if data.get('mode') else None
        tags = data.get('tags', [])
        date_from = datetime.fromisoformat(data['dateFrom']) if data.get('dateFrom') else None
        date_to = datetime.fromisoformat(data['dateTo']) if data.get('dateTo') else None
        limit = data.get('limit', 50)
        
        results = session_manager.search_sessions(
            query=query,
            user_id=user_id,
            mode=mode,
            tags=tags,
            date_from=date_from,
            date_to=date_to,
            limit=limit
        )
        
        return jsonify({
            'success': True,
            'results': [s.to_dict() for s in results],
            'count': len(results)
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
@limiter.limit("100 per minute")
async def chat():
    try:
        message = request.form.get('message', '')
        session_id = request.form.get('sessionId')
        mode = ChatMode(request.form.get('mode', 'balanced'))
        files = request.files.getlist('files')
        advanced_mode = request.form.get('advancedMode', 'false') == 'true'
        stream = request.form.get('stream', 'true') == 'true'
        
        session = session_manager.get_session(session_id)
        if not session:
            return jsonify({'success': False, 'error': 'Session not found'}), 404
        
        attachments = []
        if files:
            for file in files:
                attachment_id = str(uuid.uuid4())
                attachment = await file_processor.process_file(file, attachment_id)
                attachments.append(attachment)
        
        if advanced_mode:
            prompt_engineer.advanced_mode_active = True
        
        enhanced_message = prompt_engineer.enhance_prompt(
            message,
            mode,
            optimization_level=4 if advanced_mode else 3
        )
        
        full_prompt = enhanced_message
        if attachments:
            file_contexts = []
            for att in attachments:
                if att.processed and att.extracted_text:
                    file_contexts.append(f"\n[{att.filename}]:\n{att.extracted_text}")
            
            if file_contexts:
                full_prompt += "\n\n--- Attached Files ---" + "".join(file_contexts)
        
        system_prompt = prompt_engineer.create_system_prompt(mode)
        
        messages = [
            {"role": "system", "content": system_prompt}
        ]
        
        context_messages = session.messages[-30:]
        for msg in context_messages:
            if msg.role in [MessageRole.USER, MessageRole.ASSISTANT]:
                messages.append({
                    "role": msg.role.value,
                    "content": msg.content[:10000]
                })
        
        messages.append({"role": "user", "content": full_prompt})
        
        required_capabilities = []
        if advanced_mode:
            required_capabilities.extend(["extended_context", "advanced_reasoning"])
        if mode == ChatMode.CODE:
            required_capabilities.append("multimodal")
        
        client, api_key = client_manager.get_client(required_capabilities)
        if not client:
            return jsonify({
                'success': False,
                'error': 'No AI service available. Please try again later.'
            }), 503
        
        model = AIModel.GEMINI_2_FLASH_EXP.value
        
        start_time = time.time()
        
        try:
            response = client.chat.completions.create(
                model=model,
                messages=messages,
                max_tokens=8192 if not advanced_mode else 16384,
                temperature=0.9 if mode == ChatMode.CREATIVE else 0.7 if mode == ChatMode.BALANCED else 0.3,
                stream=stream,
                top_p=0.95,
                frequency_penalty=0.0,
                presence_penalty=0.0
            )
            
            if stream:
                def generate():
                    assistant_message = ""
                    token_count = 0
                    stream_sequence = 0
                    
                    try:
                        for chunk in response:
                            if chunk.choices[0].delta.content:
                                content = chunk.choices[0].delta.content
                                assistant_message += content
                                token_count += len(content) // 4
                                stream_sequence += 1
                                
                                stream_data = {
                                    'content': content,
                                    'sequence': stream_sequence,
                                    'timestamp': time.time()
                                }
                                
                                yield f"data: {json.dumps(stream_data)}\n\n"
                        
                        user_msg = Message(
                            id=str(uuid.uuid4()),
                            role=MessageRole.USER,
                            content=message,
                            timestamp=datetime.now(timezone.utc),
                            tokens=len(message) // 4,
                            attachments=attachments,
                            model_used=model
                        )
                        
                        processing_time = time.time() - start_time
                        
                        assistant_msg = Message(
                            id=str(uuid.uuid4()),
                            role=MessageRole.ASSISTANT,
                            content=assistant_message,
                            timestamp=datetime.now(timezone.utc),
                            tokens=token_count,
                            model_used=model,
                            processing_time=processing_time,
                            metadata={
                                'mode': mode.value,
                                'advanced': advanced_mode
                            }
                        )
                        
                        session.add_message(user_msg)
                        session.add_message(assistant_msg)
                        
                        if not session.title and len(session.messages) == 2:
                            session.title = message[:50] + ('...' if len(message) > 50 else '')
                        
                        session_manager.update_session(session)
                        
                        yield f"data: {json.dumps({'tokenUsage': session.total_tokens})}\n\n"
                        yield "data: [DONE]\n\n"
                        
                    except Exception as e:
                        yield f"data: {json.dumps({'error': str(e)})}\n\n"
                        yield "data: [DONE]\n\n"
                
                response_time = time.time() - start_time
                key_manager.mark_key_success(api_key, response_time)
                
                return Response(
                    stream_with_context(generate()),
                    mimetype='text/event-stream',
                    headers={
                        'Cache-Control': 'no-cache',
                        'X-Accel-Buffering': 'no',
                        'Connection': 'keep-alive'
                    }
                )
            
            else:
                ai_response = response.choices[0].message.content
                
                user_msg = Message(
                    id=str(uuid.uuid4()),
                    role=MessageRole.USER,
                    content=message,
                    timestamp=datetime.now(timezone.utc),
                    tokens=len(message) // 4,
                    attachments=attachments,
                    model_used=model
                )
                
                processing_time = time.time() - start_time
                
                assistant_msg = Message(
                    id=str(uuid.uuid4()),
                    role=MessageRole.ASSISTANT,
                    content=ai_response,
                    timestamp=datetime.now(timezone.utc),
                    tokens=len(ai_response) // 4,
                    model_used=model,
                    processing_time=processing_time,
                    metadata={
                        'mode': mode.value,
                        'advanced': advanced_mode
                    }
                )
                
                session.add_message(user_msg)
                session.add_message(assistant_msg)
                
                if not session.title and len(session.messages) == 2:
                    session.title = message[:50] + ('...' if len(message) > 50 else '')
                
                session_manager.update_session(session)
                
                response_time = time.time() - start_time
                key_manager.mark_key_success(api_key, response_time)
                
                return jsonify({
                    'success': True,
                    'response': ai_response,
                    'tokenUsage': session.total_tokens,
                    'processingTime': processing_time,
                    'model': model
                })
                
        except Exception as api_error:
            key_manager.mark_key_failure(api_key, str(api_error))
            raise api_error
            
    except Exception as e:
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/models', methods=['GET'])
@cache.cached(timeout=3600)
def get_available_models():
    models = [
        {
            'id': AIModel.GEMINI_2_FLASH_EXP.value,
            'name': 'Gemini 2.0 Flash Experimental',
            'provider': 'Google',
            'capabilities': ['chat', 'streaming', 'extended_context', 'multimodal'],
            'context_window': 2000000,
            'default': True
        },
        {
            'id': AIModel.GEMINI_2_FLASH_THINKING.value,
            'name': 'Gemini 2.0 Flash Thinking',
            'provider': 'Google',
            'capabilities': ['chat', 'streaming', 'advanced_reasoning'],
            'context_window': 2000000
        },
        {
            'id': AIModel.GEMINI_15_PRO.value,
            'name': 'Gemini 1.5 Pro',
            'provider': 'Google',
            'capabilities': ['chat', 'streaming', 'multimodal'],
            'context_window': 1000000
        },
        {
            'id': AIModel.GEMINI_15_FLASH.value,
            'name': 'Gemini 1.5 Flash',
            'provider': 'Google',
            'capabilities': ['chat', 'streaming'],
            'context_window': 1000000
        }
    ]
    
    return jsonify({
        'success': True,
        'models': models,
        'defaultModel': AIModel.GEMINI_2_FLASH_EXP.value
    })

@app.route('/api/upload', methods=['POST'])
@limiter.limit("50 per minute")
async def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'success': False, 'error': 'No file selected'}), 400
        
        attachment_id = str(uuid.uuid4())
        attachment = await file_processor.process_file(file, attachment_id)
        
        return jsonify({
            'success': True,
            'fileId': attachment.id,
            'filename': attachment.filename,
            'size': attachment.size,
            'type': attachment.file_type.value,
            'processed': attachment.processed,
            'metadata': attachment.metadata
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/export/<format_type>', methods=['GET'])
@limiter.limit("20 per minute")
def export_session(format_type):
    try:
        session_id = request.args.get('sessionId')
        session = session_manager.get_session(session_id)
        
        if not session:
            return jsonify({'success': False, 'error': 'Session not found'}), 404
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        if format_type == 'json':
            data = session.to_dict()
            
            return jsonify(data), 200, {
                'Content-Disposition': f'attachment; filename=ai_ultra_export_{timestamp}.json'
            }
            
        elif format_type == 'markdown':
            content = f"# Jack's AI Ultra - Ultimate Edition\n\n"
            content += f"**Export Date:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
            content += f"**Session ID:** {session.id}\n"
            content += f"**Mode:** {session.mode.value}\n"
            content += f"**Total Messages:** {len(session.messages)}\n"
            content += f"**Total Tokens:** {session.total_tokens:,}\n\n"
            
            if session.analytics:
                content += "## Analytics\n\n"
                content += f"- Average Response Time: {session.analytics.get('average_response_time', 0):.2f}ms\n"
                content += f"- Total Cost: ${session.analytics.get('total_cost', 0):.4f}\n"
                content += f"- Models Used: {', '.join(session.analytics.get('models_used', {}).keys())}\n\n"
            
            content += "## Conversation\n\n"
            
            for msg in session.messages:
                icon = "👤" if msg.role == MessageRole.USER else "🤖"
                content += f"### {icon} {msg.role.value.title()}\n"
                content += f"*{msg.timestamp.strftime('%Y-%m-%d %H:%M:%S')}*\n\n"
                content += f"{msg.content}\n\n"
                
                if msg.attachments:
                    content += "**Attachments:**\n"
                    for att in msg.attachments:
                        content += f"- {att.filename} ({att.file_type.value}, {att.size} bytes)\n"
                    content += "\n"
                
                if msg.reactions:
                    content += "**Reactions:** "
                    content += ", ".join([f"{emoji} ({count})" for emoji, count in msg.reactions.items()])
                    content += "\n\n"
                
                content += "---\n\n"
            
            return Response(content, mimetype='text/markdown', headers={
                'Content-Disposition': f'attachment; filename=ai_ultra_export_{timestamp}.md'
            })
            
        elif format_type == 'html':
            html_content = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="UTF-8">
                <title>AI Ultra Export - {session.id}</title>
                <style>
                    body {{ font-family: -apple-system, BlinkMacSystemFont, sans-serif; line-height: 1.6; max-width: 800px; margin: 0 auto; padding: 20px; }}
                    .message {{ margin: 20px 0; padding: 15px; border-radius: 10px; }}
                    .user {{ background: #e3f2fd; }}
                    .assistant {{ background: #f5f5f5; }}
                    .meta {{ font-size: 0.9em; color: #666; }}
                    pre {{ background: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }}
                    code {{ background: #f4f4f4; padding: 2px 4px; border-radius: 3px; }}
                </style>
            </head>
            <body>
                <h1>AI Ultra Conversation Export</h1>
                <p><strong>Date:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
                <p><strong>Session:</strong> {session.id}</p>
                <p><strong>Messages:</strong> {len(session.messages)}</p>
                <hr>
            """
            
            for msg in session.messages:
                role_class = 'user' if msg.role == MessageRole.USER else 'assistant'
                html_content += f"""
                <div class="message {role_class}">
                    <div class="meta">
                        <strong>{msg.role.value.title()}</strong> - {msg.timestamp.strftime('%Y-%m-%d %H:%M:%S')}
                    </div>
                    <div class="content">{msg.content}</div>
                </div>
                """
            
            html_content += """
            </body>
            </html>
            """
            
            return Response(html_content, mimetype='text/html', headers={
                'Content-Disposition': f'attachment; filename=ai_ultra_export_{timestamp}.html'
            })
            
        else:
            return jsonify({'success': False, 'error': 'Invalid format'}), 400
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/share', methods=['POST'])
@limiter.limit("10 per minute")
def share_session():
    try:
        data = request.json
        session_id = data.get('sessionId')
        
        session = session_manager.get_session(session_id)
        if not session:
            return jsonify({'success': False, 'error': 'Session not found'}), 404
        
        share_id = str(uuid.uuid4())
        session.shared = True
        session.share_id = share_id
        
        session_manager.update_session(session)
        
        return jsonify({
            'success': True,
            'shareId': share_id,
            'shareUrl': f"/shared/{share_id}"
        })
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/stats', methods=['GET'])
@cache.cached(timeout=60)
def get_stats():
    try:
        total_sessions = len(session_manager.sessions)
        total_messages = sum(len(s.messages) for s in session_manager.sessions.values())
        total_tokens = sum(s.total_tokens for s in session_manager.sessions.values())
        
        key_stats = key_manager.get_stats()
        
        mode_distribution = defaultdict(int)
        for session in session_manager.sessions.values():
            mode_distribution[session.mode.value] += 1
        
        return jsonify({
            'success': True,
            'stats': {
                'sessions': {
                    'total': total_sessions,
                    'active': sum(1 for s in session_manager.sessions.values() 
                                if (datetime.now(timezone.utc) - s.updated_at).days < 7),
                    'shared': sum(1 for s in session_manager.sessions.values() if s.shared)
                },
                'messages': {
                    'total': total_messages,
                    'average_per_session': total_messages / total_sessions if total_sessions > 0 else 0
                },
                'tokens': {
                    'total': total_tokens,
                    'average_per_session': total_tokens / total_sessions if total_sessions > 0 else 0
                },
                'modes': dict(mode_distribution),
                'keys': key_stats,
                'performance': {
                    'uptime': '99.9%',
                    'average_response_time': key_stats.get('average_response_time', 0),
                    'success_rate': key_stats.get('success_rate', 0)
                }
            }
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now(timezone.utc).isoformat(),
        'version': '3.0.0',
        'features': [
            'streaming',
            'advanced_prompts',
            'extended_context',
            'multi_modal',
            'collaboration',
            'analytics',
            'export',
            'sharing'
        ]
    })

@app.errorhandler(404)
def not_found(error):
    return jsonify({'error': 'Not found'}), 404

@app.errorhandler(429)
def rate_limit_exceeded(error):
    return jsonify({'error': 'Rate limit exceeded. Please try again later.'}), 429

@app.errorhandler(500)
def internal_error(error):
    return jsonify({'error': 'Internal server error'}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    
    if os.environ.get('ENVIRONMENT') == 'production':
        from waitress import serve
        print(f"Starting Jack's AI Ultra Ultimate Edition on port {port} (Production Mode)")
        serve(app, host='0.0.0.0', port=port, threads=50, connection_limit=1000)
    else:
        print(f"Starting Jack's AI Ultra Ultimate Edition on port {port} (Development Mode)")

        app.run(host='0.0.0.0', port=port, debug=True)
